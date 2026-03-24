"""
PPTX to JSON Converter for PraisonAI PPT

Extracts slide content from a .pptx file and returns a dict conforming to the
praisonaippt JSON schema (the inverse of create_presentation).

All extraction is read-only — the source PPTX is never modified.
"""

import os
import re
import json
from pathlib import Path
from typing import Optional, Dict, Any, List, Tuple

# ── colour helpers ─────────────────────────────────────────────────────────────

def _rgb_to_hex(rgb) -> str:
    """Convert an RGBColor (or (r,g,b) tuple) to '#RRGGBB' string."""
    try:
        r, g, b = int(rgb.r), int(rgb.g), int(rgb.b)
    except AttributeError:
        r, g, b = int(rgb[0]), int(rgb[1]), int(rgb[2])
    return f"#{r:02X}{g:02X}{b:02X}"


def _color_distance(rgb1, rgb2) -> float:
    """Simple Euclidean distance in RGB space."""
    try:
        r1, g1, b1 = int(rgb1.r), int(rgb1.g), int(rgb1.b)
        r2, g2, b2 = int(rgb2.r), int(rgb2.g), int(rgb2.b)
    except AttributeError:
        return 999
    return ((r1 - r2) ** 2 + (g1 - g2) ** 2 + (b1 - b2) ** 2) ** 0.5


def _parse_hex_to_rgb_triple(hex_str: str) -> Tuple[int, int, int]:
    """Convert '#RRGGBB' or 'RRGGBB' to (r,g,b) ints."""
    h = hex_str.lstrip('#')
    return int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16)


# ── text-frame helpers ─────────────────────────────────────────────────────────

def _tf_text(tf) -> str:
    """Full text of a text frame, newline-separated by paragraph."""
    lines = []
    for para in tf.paragraphs:
        lines.append(''.join(run.text for run in para.runs))
    return '\n'.join(lines).strip()


def _all_runs(tf):
    """Yield every run in a text frame."""
    for para in tf.paragraphs:
        for run in para.runs:
            yield run


def _most_common(items):
    """Return the most-common item in a list (or None if empty)."""
    if not items:
        return None
    return max(set(items), key=items.count)


# ── reference detection ────────────────────────────────────────────────────────

# Covers:
#   English  → "John 3:16 (NIV)", "2 Corinthians 5:7", "Mark 10:30 (NKJV) - Part 1"
#   Tamil    → "மத்தேயு 9:20-22 (TAOVBSI)"
#   Hebrew   → "ரோமர் 12:2 (TAOVBSI)"
#   Numbered book → "1 Corinthians 11:30", "2 Kings 3:17"
#   Ranges   → "Matthew 8:5–10", "Exodus 12:40-41"
#   Multi-verse → "Exodus 7:7; Acts 7:23, 30"
REFERENCE_RE = re.compile(
    r'^\d?\s*'                                  # optional leading book number
    r'[\w\u0B80-\u0BFF\u0600-\u06FF]'          # first char: word or Tamil/Arabic
    r'[\w\u0B80-\u0BFF\u0600-\u06FF\s\'\-]{1,50}'  # book name
    r'\d+[:\.\-–]\d+'                           # chapter:verse or range
)


def _looks_like_reference(text: str) -> bool:
    """Return True if text looks like a Bible reference."""
    if not text or not text.strip():
        return False
    return bool(REFERENCE_RE.match(text.strip()))


# ── slide classification ───────────────────────────────────────────────────────

def _get_text_boxes(slide) -> List:
    """Return all shapes with text frames, sorted top-to-bottom."""
    boxes = [s for s in slide.shapes if s.has_text_frame]
    return sorted(boxes, key=lambda s: s.top)


def _classify_slide(slide, index: int) -> str:
    """
    Returns one of: 'title', 'section', 'verse', 'list'

    Rules:
      - slide 0 is always 'title'
      - single textbox, bold, >36pt, centered → 'section'
      - single textbox, short text (<60 chars), no Bible reference pattern → 'section'
      - text starts with bullet (•) or numbered prefix → 'list'
      - otherwise → 'verse'
    """
    if index == 0:
        return 'title'

    boxes = _get_text_boxes(slide)
    if not boxes:
        return 'verse'

    # Single textbox heuristic for section slide
    if len(boxes) == 1:
        tf = boxes[0].text_frame
        text = _tf_text(tf)
        if text:
            # Check first paragraph formatting
            paras = [p for p in tf.paragraphs if p.text.strip()]
            if paras:
                p = paras[0]
                runs = list(p.runs)
                if runs:
                    size_pt = runs[0].font.size
                    bold = runs[0].font.bold
                    size_val = size_pt.pt if size_pt else 0
                    if size_val >= 36 and bold:
                        return 'section'

            # Short, non-Bible-reference single textbox → section header
            stripped = text.strip()
            # Only short text (< 60 chars) with no verse pattern qualifies
            if (len(stripped) < 80 and '\n' not in stripped.strip()
                    and not re.search(r'\d+[:.\-–]\d+', stripped)):
                return 'section'

    # Bullet / numbered list detection
    full_text = '\n'.join(_tf_text(b.text_frame) for b in boxes)
    stripped = full_text.strip()
    if stripped.startswith('•') or stripped.startswith('–') or stripped.startswith('-'):
        return 'list'
    if re.match(r'^\d+\.\s', stripped):
        return 'list'

    return 'verse'


# ── slide_style extraction ─────────────────────────────────────────────────────

def _extract_slide_style(prs) -> dict:
    """
    Scan all slides to infer a slide_style dict:
      background_color, background_image (warning),
      text_color, reference_position, alignment,
      highlight_color, annotation_color, font_name
    """
    from pptx.oxml.ns import qn

    style = {}
    warnings = []

    # ── background ────────────────────────────────────────────────────────────
    has_bg_image = False
    bg_colors = []
    for slide in prs.slides:
        fill = slide.background.fill
        try:
            if fill.type is not None:
                # Solid fill
                try:
                    rgb = fill.fore_color.rgb
                    bg_colors.append((int(rgb.r), int(rgb.g), int(rgb.b)))
                except Exception:
                    pass
        except Exception:
            pass
        # Check for picture shape used as background
        for shape in slide.shapes:
            if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                # Check if it covers the full slide (background-sized)
                if (abs(shape.left) < 100000 and abs(shape.top) < 100000 and
                        abs(shape.width - prs.slide_width) < 500000):
                    has_bg_image = True
                    break

    if has_bg_image:
        warnings.append("background_image: file path not recoverable from PPTX binary")
        # We don't emit background_image (path unknown)
    elif bg_colors:
        r, g, b = _most_common(bg_colors)
        if not (r == 255 and g == 255 and b == 255):  # skip pure white
            style['background_color'] = f"#{r:02X}{g:02X}{b:02X}"

    # ── run-level analysis across all verse slides ─────────────────────────────
    body_colors = []
    run_colors = []
    font_names = []
    alignments = []
    ref_positions = []
    body_size_pts = []

    from pptx.util import Pt
    from pptx.enum.text import PP_ALIGN

    ALIGN_MAP = {PP_ALIGN.LEFT: 'left', PP_ALIGN.RIGHT: 'right',
                 PP_ALIGN.CENTER: 'center', PP_ALIGN.JUSTIFY: 'left'}

    for slide in prs.slides:
        boxes = _get_text_boxes(slide)
        if not boxes:
            continue
        for shape in boxes:
            tf = shape.text_frame
            for para in tf.paragraphs:
                if para.alignment and para.alignment in ALIGN_MAP:
                    alignments.append(ALIGN_MAP[para.alignment])
                for run in para.runs:
                    try:
                        color = run.font.color.rgb
                        run_colors.append((int(color.r), int(color.g), int(color.b)))
                    except Exception:
                        pass
                    if run.font.name:
                        font_names.append(run.font.name)
                    if run.font.size:
                        sz = run.font.size.pt
                        if 20 <= sz <= 40:
                            body_size_pts.append(sz)

        # reference position: look for small text box with reference pattern
        for shape in boxes:
            tf = shape.text_frame
            text = _tf_text(tf)
            if _looks_like_reference(text) or (text and len(text) < 60 and shape.height < prs.slide_height * 0.15):
                top_ratio = shape.top / prs.slide_height
                if top_ratio < 0.25:
                    ref_positions.append('top')
                else:
                    ref_positions.append('bottom')

    # Most common body color
    if run_colors:
        mc = _most_common(run_colors)
        r, g, b = mc
        if r > 200 and g > 200 and b > 200:
            style['text_color'] = 'white'
        elif r < 30 and g < 30 and b < 30:
            style['text_color'] = '#000000'
        else:
            style['text_color'] = f"#{r:02X}{g:02X}{b:02X}"
        body_rgb = mc
    else:
        body_rgb = (0, 0, 0)

    # Reference position
    if ref_positions:
        style['reference_position'] = _most_common(ref_positions)

    # Alignment
    if alignments:
        style['alignment'] = _most_common(alignments)

    # Font name
    if font_names:
        fn = _most_common(font_names)
        if fn and fn.lower() not in ('calibri', 'arial'):
            style['font_name'] = fn

    # Highlight color: runs whose color differs significantly from body_rgb
    highlight_candidates = []
    annotation_candidates = []
    body_r, body_g, body_b = body_rgb

    for slide in prs.slides:
        for shape in _get_text_boxes(slide):
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    try:
                        color = run.font.color.rgb
                        cr, cg, cb = int(color.r), int(color.g), int(color.b)
                        dist = ((cr - body_r)**2 + (cg - body_g)**2 + (cb - body_b)**2) ** 0.5
                        if dist > 40:
                            # Check if superscript (annotation)
                            rPr = run._r.find('{http://schemas.openxmlformats.org/drawingml/2006/main}rPr')
                            if rPr is not None and rPr.get('baseline') == '30000':
                                annotation_candidates.append(f"#{cr:02X}{cg:02X}{cb:02X}")
                            else:
                                highlight_candidates.append(f"#{cr:02X}{cg:02X}{cb:02X}")
                    except Exception:
                        pass

    if highlight_candidates:
        style['highlight_color'] = _most_common(highlight_candidates)
    if annotation_candidates:
        style['annotation_color'] = _most_common(annotation_candidates)

    return style, warnings


def _extract_slide_size(prs) -> Optional[str]:
    """Map presentation dimensions to a preset string or None for default."""
    from pptx.util import Inches
    PRESETS = {
        (13.33, 7.5): 'widescreen',
        (10.0,  7.5): 'standard',
        (12.8,  8.0): '16:10',
    }
    w_in = round(prs.slide_width.inches, 2)
    h_in = round(prs.slide_height.inches, 2)
    return PRESETS.get((w_in, h_in))


# ── per-slide extractors ───────────────────────────────────────────────────────

def _extract_title_slide(slide) -> Tuple[str, str]:
    """Return (title, subtitle) from the first slide."""
    boxes = _get_text_boxes(slide)
    if not boxes:
        return '', ''
    sizes = []
    for box in boxes:
        tf = box.text_frame
        for run in _all_runs(tf):
            sz = run.font.size
            sizes.append((sz.pt if sz else 0, _tf_text(tf), box))
    if not sizes:
        return _tf_text(boxes[0].text_frame), ''
    sizes.sort(key=lambda x: -x[0])
    title = sizes[0][1].strip()
    subtitle = sizes[1][1].strip() if len(sizes) > 1 else ''
    # Avoid returning the same text as both title and subtitle
    if subtitle == title:
        subtitle = ''
    return title, subtitle


def _extract_section_name(slide) -> str:
    """Return the section name from a section slide."""
    boxes = _get_text_boxes(slide)
    if not boxes:
        return ''
    return _tf_text(boxes[0].text_frame).strip()


def _reconstruct_text_with_verse_nums(tf) -> str:
    """
    Reconstruct verse text, converting superscript verse-number runs
    back into the 'NN text' prefix format used in source JSON.
    """
    result_lines = []
    for para in tf.paragraphs:
        line_parts = []
        for run in para.runs:
            try:
                rPr = run._r.find(
                    '{http://schemas.openxmlformats.org/drawingml/2006/main}rPr')
                is_superscript = (rPr is not None and rPr.get('baseline') == '30000')
            except Exception:
                is_superscript = False

            txt = run.text
            if not txt:
                continue
            if is_superscript:
                # Strip narrow space if present; this is a verse number
                num = txt.strip('\u2009').strip()
                if num.isdigit():
                    line_parts.append(num + ' ')
                # Else it's an annotation bubble (❶❷…) — skip
            else:
                line_parts.append(txt)
        line = ''.join(line_parts).rstrip()
        if line:
            result_lines.append(line)
    return '\n'.join(result_lines)


def _detect_list_type(text: str) -> Optional[str]:
    """Detect bullet or numbered list from text content."""
    stripped = text.strip()
    if stripped.startswith('•') or stripped.startswith('–'):
        return 'bullet'
    if re.match(r'^\d+\.\s', stripped):
        return 'numbered'
    return None


def _strip_list_prefix(text: str, list_type: str) -> str:
    """Remove bullet/number prefixes added by add_list_slide()."""
    lines = text.split('\n')
    cleaned = []
    for line in lines:
        line = line.strip()
        if list_type == 'bullet':
            line = re.sub(r'^[•–\-]\s*', '', line)
        elif list_type == 'numbered':
            line = re.sub(r'^\d+\.\s*', '', line)
        if line:
            cleaned.append(line)
    return '\n'.join(cleaned)


def _is_warm_color(r: int, g: int, b: int) -> bool:
    """Orange / gold / yellow family."""
    return r > 180 and b < 130 and r > g * 0.9


def _is_cool_color(r: int, g: int, b: int) -> bool:
    """Blue / cyan family."""
    return b > 120 and b > r and b > g * 0.8


def _detect_highlights_and_large_text(
        tf,
        body_color_rgb: tuple,
        highlight_color_hex: Optional[str]
) -> Tuple[list, Optional[dict]]:
    """
    Scan runs in a text frame. Return:
      - highlights: list of strings (primary color) or {"text","color"} dicts (secondary)
      - large_text: dict of {word: font_size_pt} for oversized runs, or None

    Detection strategy (in priority order):
      1. Large font (>140% of body_size) → large_text dict entry
      2. Explicitly colored run, different color from body:
           warm (orange/gold/yellow) → simple string (uses highlight_color)
           cool (blue)              → {"text": ..., "color": "#1E50C8"} (annotation)
           other non-reference      → simple string
           reference color #003366  → skipped (it's the ref line color)
      3. Bold run with SAME color as body → simple string (bold = emphasis in many PPTXs)
    """
    REFERENCE_RGB = (0x00, 0x33, 0x66)   # dark blue used for ref lines
    ANNOTATION_COLOR = "#1E50C8"          # our design secondary color
    NS = '{http://schemas.openxmlformats.org/drawingml/2006/main}'

    highlights_out = []
    large_text_out = {}
    seen_texts = set()

    # Determine body font size
    sizes = []
    for para in tf.paragraphs:
        for run in para.runs:
            if run.font.size:
                sizes.append(run.font.size.pt)
    body_size = _most_common(sizes) or 32

    br, bg_, bb = body_color_rgb

    # Determine if body text is predominantly bold
    # (if so, bold alone shouldn't trigger highlight)
    bold_flags = []
    for para in tf.paragraphs:
        for run in para.runs:
            if run.text.strip():
                bold_flags.append(bool(run.font.bold))
    all_bold = len(bold_flags) > 0 and all(bold_flags)

    for para in tf.paragraphs:
        for run in para.runs:
            txt = run.text
            if not txt or not txt.strip():
                continue

            # Skip superscript (verse number or annotation bubble)
            try:
                rPr = run._r.find(
                    '{http://schemas.openxmlformats.org/drawingml/2006/main}rPr')
                is_superscript = (rPr is not None and rPr.get('baseline') == '30000')
            except Exception:
                is_superscript = False
            if is_superscript:
                continue

            # Large text detection (significantly bigger than body)
            if run.font.size:
                sz = run.font.size.pt
                if sz > body_size * 1.4 and txt.strip() not in seen_texts:
                    large_text_out[txt.strip()] = min(int(sz), 80)
                    seen_texts.add(txt.strip())
                    continue

            word = txt.strip()
            if not word or word in seen_texts:
                continue

            # ── 2. <a:highlight> XML element (background fill highlight) ──────
            # Used in modern PowerPoint for text highlighter pen (yellow, green, etc.)
            try:
                rPr = run._r.find(NS + 'rPr')
                if rPr is not None:
                    hl_elem = rPr.find(NS + 'highlight')
                    if hl_elem is not None:
                        seen_texts.add(word)
                        highlights_out.append(word)
                        continue
            except Exception:
                pass

            # ── 3. Color-based highlight detection ───────────────────────────
            detected = False
            try:
                color = run.font.color.rgb
                cr, cg, cb = int(color.r), int(color.g), int(color.b)
                dist_from_body = ((cr-br)**2 + (cg-bg_)**2 + (cb-bb)**2) ** 0.5

                # Skip if same as body color
                if dist_from_body < 30:
                    # Fall through to bold check below
                    pass
                # Skip reference line color
                elif (abs(cr - REFERENCE_RGB[0]) < 20 and
                      abs(cg - REFERENCE_RGB[1]) < 20 and
                      abs(cb - REFERENCE_RGB[2]) < 20):
                    continue  # reference-colored text, not a highlight
                else:
                    # Color is distinct — map to our design palette
                    seen_texts.add(word)
                    if _is_cool_color(cr, cg, cb):
                        # Secondary (cool) color → annotation color object
                        highlights_out.append({"text": word, "color": ANNOTATION_COLOR})
                    else:
                        # Warm or other → primary highlight (simple string)
                        highlights_out.append(word)
                    detected = True

            except Exception:
                pass  # No explicit color set (theme color) → fall through to bold check

            # ── Bold-as-highlight fallback ─────────────────────────────────
            if not detected and not all_bold and bool(run.font.bold):
                seen_texts.add(word)
                highlights_out.append(word)

    # ── Isolated inner-run detection (same-color/bold) ─────────────────────
    # Runs in the middle of a 3+ run paragraph with same formatting = highlight artifact
    for para in tf.paragraphs:
        text_runs = [r for r in para.runs if r.text.strip()]
        if len(text_runs) < 3:
            continue
        # Check all runs have same bold+color
        def run_key(r):
            try: rgb = str(r.font.color.rgb)
            except: rgb = ''
            return (bool(r.font.bold), rgb)
        first_key = run_key(text_runs[0])
        if not all(run_key(r) == first_key for r in text_runs):
            continue  # mixed formatting → handled by color/bold detection above
        # Inner runs (not first, not last) are highlight candidates
        for inner_run in text_runs[1:-1]:
            word = inner_run.text.strip()
            if word and word not in seen_texts and len(word) > 2:
                seen_texts.add(word)
                highlights_out.append(word)

    return highlights_out, (large_text_out if large_text_out else None)



def _extract_verse_from_slide(
        slide, prs,
        body_color_rgb: tuple,
        highlight_color_hex: Optional[str]
) -> dict:
    """
    Extract a verse dict from a regular verse slide.
    Returns: {reference, text, highlights?, large_text?, list_type?}
    """
    from pptx.util import Inches

    boxes = _get_text_boxes(slide)
    if not boxes:
        return {'reference': '', 'text': ''}

    slide_height = prs.slide_height

    # Identify reference box vs body box
    # Reference box: height < 22% of slide AND positioned near top or bottom
    ref_box = None
    body_boxes = []
    for box in boxes:
        text = _tf_text(box.text_frame).strip()
        if not text:
            continue
        top_ratio = box.top / slide_height
        height_ratio = box.height / slide_height
        if height_ratio < 0.22 and (top_ratio < 0.25 or top_ratio > 0.72):
            if ref_box is None:
                ref_box = box
            else:
                body_boxes.append(box)
        else:
            body_boxes.append(box)

    # If only 1 box exists, it's the body (no reference)
    if not body_boxes and ref_box:
        body_boxes = [ref_box]
        ref_box = None

    reference = _tf_text(ref_box.text_frame).strip() if ref_box else ''
    body_tf = body_boxes[0].text_frame if body_boxes else None

    if body_tf is None:
        return {'reference': reference, 'text': ''}

    # Reconstruct text (preserving verse number superscripts)
    text = _reconstruct_text_with_verse_nums(body_tf)

    # Detect list type
    list_type = _detect_list_type(text)
    if list_type:
        text = _strip_list_prefix(text, list_type)

    # Detect highlights + large_text
    highlights, large_text = _detect_highlights_and_large_text(
        body_tf, body_color_rgb, highlight_color_hex)

    verse: Dict[str, Any] = {
        'reference': reference,
        'text': text,
    }
    if highlights:
        verse['highlights'] = highlights
    if list_type:
        verse['list_type'] = list_type
    if large_text:
        verse['large_text'] = large_text

    return verse


# ── main converter class ───────────────────────────────────────────────────────

class PPTXToJSONConverter:
    """
    Convert a .pptx file produced by praisonaippt back to its JSON schema.

    Usage:
        converter = PPTXToJSONConverter("presentation.pptx")
        data = converter.convert()
    """

    def __init__(self, pptx_path: str):
        pptx_path = str(pptx_path)
        if not pptx_path.lower().endswith(('.pptx', '.ppt')):
            raise ValueError("Input must be a .pptx or .ppt file")
        if not os.path.exists(pptx_path):
            raise FileNotFoundError(f"File not found: {pptx_path}")
        self._path = pptx_path

    def convert(self) -> dict:
        """
        Extract and return a praisonaippt-compatible dict from the PPTX.

        The returned dict may include a '_extraction_warnings' key listing
        features that could not be fully recovered (e.g. background_image path,
        annotation numbers in highlights).
        """
        from pptx import Presentation

        prs = Presentation(self._path)
        slides = list(prs.slides)

        extraction_warnings = []

        # ── slide meta ─────────────────────────────────────────────────────────
        slide_size = _extract_slide_size(prs)
        slide_style, style_warnings = _extract_slide_style(prs)
        extraction_warnings.extend(style_warnings)

        # Body colour for highlight detection
        body_color_rgb = (255, 255, 255)  # default white-on-dark
        tc = slide_style.get('text_color', '')
        if tc:
            if tc.lower() == 'white':
                body_color_rgb = (255, 255, 255)
            elif tc.startswith('#'):
                try:
                    body_color_rgb = _parse_hex_to_rgb_triple(tc)
                except Exception:
                    pass
        highlight_color_hex = slide_style.get('highlight_color')

        # ── extract title slide ────────────────────────────────────────────────
        title = ''
        subtitle = ''
        if slides:
            try:
                title, subtitle = _extract_title_slide(slides[0])
            except Exception:
                title = 'Untitled'
                subtitle = ''

        # ── iterate remaining slides ───────────────────────────────────────────
        sections: List[Dict[str, Any]] = []
        current_section: Optional[Dict[str, Any]] = None

        for idx, slide in enumerate(slides[1:], start=1):
            try:
                slide_type = _classify_slide(slide, idx)
            except Exception:
                slide_type = 'verse'

            if slide_type == 'section':
                # Push previous section
                if current_section is not None:
                    sections.append(current_section)
                try:
                    name = _extract_section_name(slide)
                except Exception:
                    name = ''
                current_section = {'section': name, 'verses': []}

            else:  # 'verse' or 'list'
                # Ensure there is a section container
                if current_section is None:
                    current_section = {'section': '', 'verses': []}

                try:
                    verse = _extract_verse_from_slide(
                        slide, prs, body_color_rgb, highlight_color_hex)
                except Exception as e:
                    verse = {
                        'reference': '',
                        'text': f'[extraction failed: {e}]',
                    }
                current_section['verses'].append(verse)

        # Flush last section
        if current_section is not None:
            sections.append(current_section)

        # ── assemble result ────────────────────────────────────────────────────
        result: Dict[str, Any] = {}

        if extraction_warnings:
            result['_extraction_warnings'] = extraction_warnings

        result['_source'] = 'extracted'
        result['presentation_title'] = title
        result['presentation_subtitle'] = subtitle

        if slide_size:
            result['slide_size'] = slide_size

        if slide_style:
            result['slide_style'] = slide_style

        result['sections'] = sections

        return result


# ── convenience function (mirrors convert_pptx_to_pdf) ─────────────────────────

def pptx_to_json(
    pptx_path: str,
    output_path: Optional[str] = None,
    pretty: bool = True,
) -> dict:
    """
    Extract a praisonaippt-compatible JSON dict from a PPTX file.

    This is the inverse of create_presentation(): given a PPTX file
    (ideally produced by this package), it reconstructs the JSON schema
    dict that could regenerate the presentation.

    Args:
        pptx_path: Path to the .pptx or .ppt file to read
        output_path: If provided, write the JSON to this file path
        pretty: If True (default), emit indented JSON (indent=2)

    Returns:
        dict conforming to the praisonaippt JSON schema

    Raises:
        FileNotFoundError: If pptx_path does not exist
        ValueError: If pptx_path is not a .pptx or .ppt file

    Examples:
        # In-memory
        from praisonaippt import pptx_to_json
        data = pptx_to_json("presentation.pptx")

        # Save to file
        pptx_to_json("presentation.pptx", output_path="output.json")

        # Feed directly back into create_presentation for a round-trip test
        from praisonaippt import create_presentation, pptx_to_json
        data = pptx_to_json("slides.pptx")
        create_presentation(data, output_file="roundtrip.pptx")

    Notes:
        - background_image file paths are not recoverable from the PPTX binary.
        - highlight annotation numbers (bubble chars ❶❷…) are not recoverable.
        - All other features are extracted losslessly or with best-effort.
    """
    converter = PPTXToJSONConverter(pptx_path)
    data = converter.convert()
    if output_path:
        Path(output_path).parent.mkdir(parents=True, exist_ok=True)
        with open(output_path, 'w', encoding='utf-8') as f:
            json.dump(data, f, indent=2 if pretty else None,
                      ensure_ascii=False)
        print(f"✓ JSON saved to: {output_path}")
    return data
