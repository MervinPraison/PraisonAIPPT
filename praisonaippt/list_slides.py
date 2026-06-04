"""Print a numbered outline of slides in a .pptx file."""

import re
from pathlib import Path

_SLIDE_TYPE_RE = re.compile(r"^slide_type:\s*(\S+)", re.MULTILINE)


def _slide_label(slide, texts, max_text_len=120):
    """Build one-line slide summary including avatar layout type from notes."""
    try:
        notes = slide.notes_slide.notes_text_frame.text or ""
        match = _SLIDE_TYPE_RE.search(notes)
        if match:
            kind = match.group(1)
            user_notes = _SLIDE_TYPE_RE.sub("", notes).strip()
            if user_notes:
                line = user_notes.replace("\n", " | ")
                if max_text_len:
                    line = line[:max_text_len]
                return f"[{kind}] {line}"
            return f"[{kind}]"
    except Exception:
        pass
    if texts:
        line = texts[0].replace("\n", " | ")
        if max_text_len:
            line = line[:max_text_len]
        return line
    return "(no text)"


def print_slide_outline(pptx_path, max_text_len=120):
    """
    Print slide count and a one-line summary per slide.

    Returns:
        int: 0 on success, 1 if the file is missing or unreadable.
    """
    path = Path(pptx_path)
    if not path.exists():
        print(f"MISSING: {path}")
        return 1

    from pptx import Presentation

    prs = Presentation(str(path))
    print(f"slides: {len(prs.slides)}")
    for i, slide in enumerate(prs.slides, 1):
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                texts.append(shape.text.strip())
        print(f"{i:2d}: {_slide_label(slide, texts, max_text_len)}")
    return 0
