"""PPTX to MP4 video export (compositor backend, Mac/Linux default)."""

from __future__ import annotations

import asyncio
import hashlib
import json
import logging
import shutil
import tempfile
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple

from pptx import Presentation

from .avatar_layouts import (
    AVATAR_SLIDE_TYPES,
    RegionBox,
    export_floating_pip_box,
    export_slide_regions,
    region_box_to_pixels,
)
from .ffmpeg_composer import (
    OverlaySpec,
    check_video_tools,
    concat_segments,
    ffprobe_duration,
    ffprobe_has_audio,
    is_image_path,
    is_video_path,
    pdf_to_png_pages,
    pick_video_encoder,
    render_slide_segment,
)
from .pdf_converter import PDFOptions, convert_pptx_to_pdf
from .exceptions import SchemaError
from .utils import resolve_asset_path

logger = logging.getLogger(__name__)

from .video_presets import VIDEO_PRESETS as _PRESETS


@dataclass
class VideoOptions:
    backend: str = "compositor"
    narration_mode: str = "fixed"
    output_path: Optional[str] = None
    width: int = 1920
    height: int = 1080
    fps: int = 30
    dpi: int = 192
    preset: str = "standard"
    slide_duration_sec: float = 5.0
    avatar_timeline: str = "auto"
    slide_range: Optional[Tuple[int, int]] = None
    keep_temp: bool = False
    avatar_fit: str = "cover"
    avatar_shape: str = "circle"
    avatar_crop_x_ratio: float = 0.5
    avatar_crop_y_ratio: float = 0.06
    avatar_zoom_ratio: float = 1.45
    media_fit_default: str = "contain"
    media_crop_y_ratio: float = 0.12
    media_zoom_ratio: float = 1.0
    loop_avatar_if_shorter: bool = True
    tts_provider: str = "edge"
    tts_voice: str = "en-GB-RyanNeural"
    captions_enabled: bool = True
    pdf_path: Optional[str] = None
    encoder: Optional[str] = None
    slide_cache: bool = True

    def validate(self) -> None:
        from .yaml_validate import validate_video_export

        try:
            validate_video_export({
                "backend": self.backend,
                "narration_mode": self.narration_mode,
                "preset": self.preset,
                "avatar_timeline": self.avatar_timeline,
                "avatar": {
                    "fit": self.avatar_fit,
                    "shape": self.avatar_shape,
                    "crop_x_ratio": self.avatar_crop_x_ratio,
                    "crop_y_ratio": self.avatar_crop_y_ratio,
                    "zoom_ratio": self.avatar_zoom_ratio,
                    "loop_if_shorter": self.loop_avatar_if_shorter,
                },
            })
        except SchemaError as exc:
            raise ValueError(str(exc)) from exc

    def __post_init__(self) -> None:
        if self.preset in _PRESETS:
            p = _PRESETS[self.preset]
            self.width = p["width"]
            self.height = p["height"]
            self.fps = p["fps"]
            self.dpi = p["dpi"]
        self.validate()

    @classmethod
    def from_dict(cls, raw: Optional[dict], deck: Optional[dict] = None) -> "VideoOptions":
        raw = dict(raw or {})
        deck = deck or {}
        opts = cls()
        if raw.get("backend"):
            opts.backend = str(raw["backend"])
        if raw.get("narration_mode"):
            opts.narration_mode = str(raw["narration_mode"])
        elif raw.get("audio_source"):
            mapped = narration_mode_from_audio_source(str(raw["audio_source"]))
            if mapped:
                opts.narration_mode = mapped
        if raw.get("output_path"):
            opts.output_path = str(raw["output_path"])
        if raw.get("preset"):
            opts.preset = str(raw["preset"])
            if opts.preset in _PRESETS:
                p = _PRESETS[opts.preset]
                opts.width, opts.height, opts.fps, opts.dpi = (
                    p["width"], p["height"], p["fps"], p["dpi"],
                )
        res = raw.get("resolution") or {}
        if isinstance(res, dict):
            if res.get("width"):
                opts.width = int(res["width"])
            if res.get("height"):
                opts.height = int(res["height"])
        if raw.get("fps") is not None:
            opts.fps = int(raw["fps"])
        if raw.get("dpi") is not None:
            opts.dpi = int(raw["dpi"])
        if raw.get("slide_duration_sec") is not None:
            opts.slide_duration_sec = float(raw["slide_duration_sec"])
        if raw.get("avatar_timeline"):
            opts.avatar_timeline = str(raw["avatar_timeline"])
        avatar = raw.get("avatar") or {}
        if isinstance(avatar, dict):
            if avatar.get("fit"):
                opts.avatar_fit = str(avatar["fit"])
            if avatar.get("shape"):
                opts.avatar_shape = str(avatar["shape"])
            if avatar.get("crop_x_ratio") is not None:
                opts.avatar_crop_x_ratio = float(avatar["crop_x_ratio"])
            if avatar.get("crop_y_ratio") is not None:
                opts.avatar_crop_y_ratio = float(avatar["crop_y_ratio"])
            if avatar.get("zoom_ratio") is not None:
                opts.avatar_zoom_ratio = float(avatar["zoom_ratio"])
            if avatar.get("loop_if_shorter") is not None:
                opts.loop_avatar_if_shorter = bool(avatar["loop_if_shorter"])
        media = raw.get("media") or {}
        if isinstance(media, dict):
            if media.get("fit"):
                opts.media_fit_default = str(media["fit"])
            if media.get("crop_y_ratio") is not None:
                opts.media_crop_y_ratio = float(media["crop_y_ratio"])
            if media.get("zoom_ratio") is not None:
                opts.media_zoom_ratio = max(1.0, float(media["zoom_ratio"]))
        tts = raw.get("tts") or {}
        if isinstance(tts, dict):
            if tts.get("provider"):
                opts.tts_provider = str(tts["provider"])
            if tts.get("voice"):
                opts.tts_voice = str(tts["voice"])
        caps = raw.get("captions") or {}
        if isinstance(caps, dict) and caps.get("enabled") is not None:
            opts.captions_enabled = bool(caps["enabled"])
        if raw.get("slide_cache") is not None:
            opts.slide_cache = bool(raw["slide_cache"])
        ts = deck.get("slide_timestamps")
        if ts:
            opts._slide_timestamps = list(ts)  # type: ignore[attr-defined]
        _merge_slide_style_pip(opts, deck.get("slide_style") or {})
        from .yaml_validate import validate_video_export

        validate_video_export(raw)
        opts.validate()
        return opts


def _merge_slide_style_pip(opts: VideoOptions, slide_style: dict) -> None:
    """Apply ``slide_style.layouts.pip`` crop/zoom/shape defaults to video options."""
    pip = (slide_style.get("layouts") or {}).get("pip") or {}
    if not isinstance(pip, dict):
        return
    if pip.get("crop_x_ratio") is not None:
        opts.avatar_crop_x_ratio = float(pip["crop_x_ratio"])
    if pip.get("crop_y_ratio") is not None:
        opts.avatar_crop_y_ratio = float(pip["crop_y_ratio"])
    if pip.get("zoom_ratio") is not None:
        opts.avatar_zoom_ratio = float(pip["zoom_ratio"])
    if pip.get("shape"):
        opts.avatar_shape = str(pip["shape"])


def resolve_video_backend(options: VideoOptions) -> str:
    """Resolve ``auto`` / reject unimplemented backends."""
    backend = options.backend
    if backend == "aspose_frames":
        raise NotImplementedError(
            "video backend 'aspose_frames' is not implemented; use compositor"
        )
    if backend == "auto":
        return "compositor"
    return backend


def required_video_tools(backend: str) -> List[str]:
    """Tools required for the resolved video backend."""
    keys = ["ffmpeg", "ffprobe"]
    if backend not in ("powerpoint",):
        keys.extend(["pdftoppm", "libreoffice"])
    return keys


def _pptx_cache_key(pptx_path: str, dpi: int) -> str:
    st = Path(pptx_path).stat()
    raw = f"{pptx_path}:{st.st_mtime_ns}:{st.st_size}:{dpi}"
    return hashlib.sha256(raw.encode()).hexdigest()[:16]


def _slide_cache_dir(pptx_path: str, dpi: int) -> Path:
    return Path.home() / ".praisonaippt" / "video_cache" / _pptx_cache_key(pptx_path, dpi)


@dataclass
class SlideVideoEntry:
    index: int
    slide_role: str
    slide_type: Optional[str]
    verse: Optional[dict]
    duration_sec: float = 5.0
    narration_mode: str = "fixed"
    audio_path: Optional[str] = None
    audio_start_sec: float = 0.0
    audio_primary: str = "none"
    avatar_video_path: Optional[str] = None
    media_path: Optional[str] = None
    media_fit: str = "contain"
    avatar_shape: Optional[str] = None
    skip_media_overlay: bool = False
    skip_avatar_overlay: bool = False
    avatar_crop_x_ratio: Optional[float] = None
    avatar_crop_y_ratio: Optional[float] = None
    avatar_zoom_ratio: Optional[float] = None
    media_crop_y_ratio: Optional[float] = None
    media_zoom_ratio: Optional[float] = None
    avatar_box_px: Optional[dict] = None
    media_box_px: Optional[dict] = None
    text_panel_px: Optional[dict] = None
    notes: str = ""
    caption_text: str = ""


def iter_slide_plan(data: dict, custom_title: Optional[str] = None):
    """Mirror create_presentation slide order for manifest building."""
    yield {"slide_role": "title", "verse": None, "slide_type": "title"}
    for section in data.get("sections", []):
        if section.get("section") and not custom_title:
            yield {
                "slide_role": "section",
                "verse": None,
                "slide_type": "section",
                "section": section.get("section"),
            }
        for verse in section.get("verses", []):
            if not isinstance(verse, dict):
                continue
            from .slide_renderers import resolve_renderer

            renderer = resolve_renderer(verse)
            if renderer.kind == "verse":
                from .layout_tokens import split_max_length_default
                from .utils import split_long_text

                style = data.get("slide_style") or {}
                max_len = int(
                    verse.get("split_max_length") or split_max_length_default(style)
                )
                parts = split_long_text(
                    verse.get("text", ""), max_length=max(max_len, 50)
                )
                for part_idx in range(len(parts)):
                    yield {
                        "slide_role": "content",
                        "verse": verse,
                        "slide_type": "verse",
                        "part_index": part_idx,
                    }
            else:
                yield {
                    "slide_role": "content",
                    "verse": verse,
                    "slide_type": verse.get("slide_type", renderer.kind),
                }


def ffprobe_has_audio_safe(path: str, source_file: Optional[str] = None) -> bool:
    try:
        resolved = resolve_asset_path(path, source_file=source_file)
        p = resolved if resolved else path
        if p and Path(p).is_file():
            return ffprobe_has_audio(p)
    except Exception:
        pass
    return False


# Optional video_export.audio_source alias (maps to narration_mode in VideoOptions.from_dict).
_AUDIO_SOURCE_TO_MODE = {
    "heygen_video": "avatar",
    "heygen": "avatar",
    "video": "avatar",
    "external": "audio_file",
    "separate": "audio_file",
    "mp3": "audio_file",
    "tts": "tts",
}


def narration_mode_from_audio_source(audio_source: str) -> Optional[str]:
    return _AUDIO_SOURCE_TO_MODE.get(str(audio_source).strip().lower())


def _resolve_narration_mode(
    verse: Optional[dict], deck_mode: str, source_file: Optional[str] = None,
) -> str:
    if verse and verse.get("narration_mode"):
        return str(verse["narration_mode"])
    if deck_mode != "auto":
        return deck_mode
    if not verse:
        return "fixed"
    av = verse.get("avatar_video_path")
    if av and ffprobe_has_audio_safe(av, source_file):
        return "avatar"
    if verse.get("audio_path"):
        return "audio_file"
    notes = (verse.get("notes") or "").strip()
    if notes:
        return "tts"
    return "fixed"


def _notes_text(verse: Optional[dict]) -> str:
    if not verse:
        return ""
    for key in ("notes", "text", "reference"):
        val = (verse.get(key) or "").strip()
        if val:
            return val
    return ""


def _box_px(
    box: Optional[RegionBox],
    slide_w_in: float,
    slide_h_in: float,
    out_w: int,
    out_h: int,
) -> Optional[dict]:
    if not box:
        return None
    return region_box_to_pixels(box, slide_w_in, slide_h_in, out_w, out_h)


def build_video_manifest(
    data: Optional[dict],
    prs: Presentation,
    options: VideoOptions,
    *,
    source_file: Optional[str] = None,
    custom_title: Optional[str] = None,
) -> List[SlideVideoEntry]:
    slide_style_base = dict((data or {}).get("slide_style") or {})
    video_export = dict((data or {}).get("video_export") or {})
    slide_w_in = prs.slide_width.inches
    slide_h_in = prs.slide_height.inches
    entries: List[SlideVideoEntry] = []

    if data:
        plan = list(iter_slide_plan(data, custom_title=custom_title))
    else:
        plan = [{"slide_role": "content", "verse": None, "slide_type": None}] * len(prs.slides)

    if len(plan) != len(prs.slides):
        if data:
            raise RuntimeError(
                f"YAML slide plan ({len(plan)} slides) != PPTX ({len(prs.slides)} slides). "
                "Rebuild the PPTX from the same YAML or pass matching custom_title."
            )
        logger.warning(
            "Slide plan count (%d) != PPTX slides (%d); using PPTX notes fallback",
            len(plan),
            len(prs.slides),
        )
        plan = []
        for i, slide in enumerate(prs.slides):
            notes = ""
            try:
                notes = slide.notes_slide.notes_text_frame.text or ""
            except Exception:
                pass
            plan.append({
                "slide_role": "content",
                "verse": {"notes": notes} if notes else None,
                "slide_type": None,
            })

    for idx, item in enumerate(plan):
        verse = item.get("verse")
        slide_type = item.get("slide_type")
        from .deck_slides import DECK_SLIDE_TYPES, resolve_deck_style

        style = slide_style_base
        if slide_type in DECK_SLIDE_TYPES:
            style = resolve_deck_style(slide_style_base, verse or {}, slide_type)
        style = dict(style)
        if source_file:
            style["_source_file"] = source_file
        mode = _resolve_narration_mode(verse, options.narration_mode, source_file)
        entry = SlideVideoEntry(
            index=idx,
            slide_role=item.get("slide_role", "content"),
            slide_type=slide_type,
            verse=verse,
            narration_mode=mode,
            notes=_notes_text(verse),
            caption_text=_notes_text(verse),
        )
        if verse:
            if verse.get("duration_sec") is not None:
                entry.duration_sec = float(verse["duration_sec"])
            if verse.get("audio_path"):
                entry.audio_path = str(verse["audio_path"])
            if verse.get("audio_start_sec") is not None:
                entry.audio_start_sec = float(verse["audio_start_sec"])
            if verse.get("avatar_video_path"):
                entry.avatar_video_path = str(verse["avatar_video_path"])
            if verse.get("media_path"):
                entry.media_path = str(verse["media_path"])
            if verse.get("media_fit"):
                entry.media_fit = str(verse["media_fit"])
            else:
                entry.media_fit = options.media_fit_default

        from .deck_slides import (
            DECK_SLIDE_TYPES,
            deck_avatar_shape,
            deck_skips_avatar_overlay,
            deck_skips_media_overlay,
            export_deck_slide_regions,
        )
        from .video_protocol import (
            apply_pixel_offset,
            region_from_placement,
            resolve_framing,
            resolve_slide_overlays,
        )

        framing_kind = (
            slide_type
            if slide_type in AVATAR_SLIDE_TYPES or slide_type in DECK_SLIDE_TYPES
            else "pip"
        )
        overlays = resolve_slide_overlays(
            verse=verse,
            slide_type=slide_type,
            style=style,
            video_export=video_export,
            framing_kind=framing_kind,
        )
        g_off = overlays.global_offset_px
        av_region = None
        media_region = None

        if slide_type in AVATAR_SLIDE_TYPES:
            regions = export_slide_regions(prs, slide_type, style)
            av_region = region_from_placement(
                regions.get("avatar"), overlays.avatar, slide_w_in, slide_h_in, style, framing_kind,
            )
            media_region = region_from_placement(
                regions.get("media"), overlays.media, slide_w_in, slide_h_in, style, framing_kind,
            )
            panel = regions.get("text_panel")
            entry.text_panel_px = _box_px(
                panel, slide_w_in, slide_h_in, options.width, options.height
            )
        elif slide_type in DECK_SLIDE_TYPES:
            regions = export_deck_slide_regions(prs, slide_type, style)
            av_region = region_from_placement(
                regions.get("avatar"), overlays.avatar, slide_w_in, slide_h_in, style, framing_kind,
            )
            entry.avatar_shape = deck_avatar_shape(
                slide_type, style, options.avatar_shape, box=av_region, verse=verse,
            )
            if deck_skips_avatar_overlay(slide_type):
                entry.skip_avatar_overlay = True
            if deck_skips_media_overlay(slide_type):
                entry.skip_media_overlay = True
            else:
                media_region = region_from_placement(
                    regions.get("media"), overlays.media, slide_w_in, slide_h_in, style, framing_kind,
                )
            panel = regions.get("text_panel") or regions.get("content")
            entry.text_panel_px = _box_px(
                panel, slide_w_in, slide_h_in, options.width, options.height
            )
        elif verse and verse.get("avatar_video_path"):
            pip = export_floating_pip_box(prs, style)
            av_region = region_from_placement(
                pip, overlays.avatar, slide_w_in, slide_h_in, style, "pip",
            )
            from .avatar_layouts import _pip_shape_kind

            pip_shape = _pip_shape_kind(style)
            if pip_shape in ("square", "rect", "rectangle"):
                entry.avatar_shape = "rect"
            else:
                entry.avatar_shape = "circle"

        av_off = (
            overlays.avatar.offset_px[0] + g_off[0],
            overlays.avatar.offset_px[1] + g_off[1],
        )
        md_off = (
            overlays.media.offset_px[0] + g_off[0],
            overlays.media.offset_px[1] + g_off[1],
        )
        entry.avatar_box_px = apply_pixel_offset(
            _box_px(av_region, slide_w_in, slide_h_in, options.width, options.height), av_off,
        )
        entry.media_box_px = apply_pixel_offset(
            _box_px(media_region, slide_w_in, slide_h_in, options.width, options.height), md_off,
        )

        if entry.avatar_video_path and entry.avatar_box_px:
            crop_x, crop_y, zoom, fit, shape = resolve_framing(
                overlays.avatar, style, framing_kind,
                default_crop=options.avatar_crop_y_ratio,
                default_zoom=options.avatar_zoom_ratio,
                default_fit=options.avatar_fit,
                default_shape=entry.avatar_shape or options.avatar_shape,
            )
            entry.avatar_crop_x_ratio = crop_x
            entry.avatar_crop_y_ratio = crop_y
            entry.avatar_zoom_ratio = zoom
            if overlays.avatar.shape:
                entry.avatar_shape = shape

        if entry.media_path and entry.media_box_px:
            _, mcrop, mzoom, mfit, _ = resolve_framing(
                overlays.media, style, framing_kind,
                default_crop=options.media_crop_y_ratio,
                default_zoom=options.media_zoom_ratio,
                default_fit=entry.media_fit,
                default_shape="rect",
            )
            entry.media_crop_y_ratio = mcrop
            entry.media_zoom_ratio = mzoom
            if overlays.media.fit:
                entry.media_fit = mfit

        entries.append(entry)
    return entries


def _timestamp_duration(ts: list, index: int) -> Optional[float]:
    if ts and index < len(ts) - 1:
        return float(ts[index + 1]) - float(ts[index])
    return None


def resolve_slide_durations(
    entries: List[SlideVideoEntry],
    options: VideoOptions,
    *,
    source_file: Optional[str] = None,
    temp_dir: Path,
) -> None:
    ts = getattr(options, "_slide_timestamps", None)
    for entry in entries:
        mode = entry.narration_mode
        verse = entry.verse or {}
        sf = source_file or verse.get("_source_file")
        explicit_duration = verse.get("duration_sec") is not None
        ts_dur = _timestamp_duration(ts, entry.index) if ts else None

        if entry.slide_role in ("title", "section"):
            entry.duration_sec = float(options.slide_duration_sec)
            entry.audio_primary = "none"
            entry.audio_start_sec = 0.0
            _apply_sync_mode(entry, verse, options, temp_dir, sf)
            continue

        if mode == "fixed":
            entry.duration_sec = float(
                verse.get("duration_sec") or options.slide_duration_sec
            )
            entry.audio_primary = "none"
        elif mode == "audio_file":
            ap = entry.audio_path or verse.get("audio_path")
            if ap:
                resolved = resolve_asset_path(ap, source_file=sf)
                path = resolved if resolved else ap
                entry.audio_path = path
                entry.audio_primary = "file"
                if explicit_duration:
                    entry.duration_sec = float(verse["duration_sec"])
                elif ts_dur is not None:
                    entry.duration_sec = ts_dur
                else:
                    entry.duration_sec = ffprobe_duration(path)
                if verse.get("audio_start_sec") is not None:
                    entry.audio_start_sec = float(verse["audio_start_sec"])
            elif ts_dur is not None:
                entry.duration_sec = ts_dur
                entry.audio_primary = "none"
            else:
                entry.duration_sec = float(
                    verse.get("duration_sec") or options.slide_duration_sec
                )
        elif mode == "avatar":
            av = entry.avatar_video_path or verse.get("avatar_video_path")
            if av:
                resolved = resolve_asset_path(av, source_file=sf)
                path = resolved if resolved else av
                entry.avatar_video_path = path
                entry.audio_primary = "avatar"
                if explicit_duration:
                    entry.duration_sec = float(verse["duration_sec"])
                elif ts_dur is not None:
                    entry.duration_sec = ts_dur
                else:
                    entry.duration_sec = ffprobe_duration(path)
                if verse.get("audio_start_sec") is not None:
                    entry.audio_start_sec = float(verse["audio_start_sec"])
            else:
                entry.duration_sec = float(
                    verse.get("duration_sec") or options.slide_duration_sec
                )
        elif mode == "tts":
            text = entry.notes or _notes_text(verse)
            if text:
                mp3 = temp_dir / f"slide_{entry.index:03d}.mp3"
                synthesise_tts(text, options.tts_voice, str(mp3), provider=options.tts_provider)
                entry.audio_path = str(mp3)
                entry.duration_sec = ffprobe_duration(str(mp3)) + 0.3
                entry.audio_primary = "tts"
            else:
                entry.duration_sec = options.slide_duration_sec
        else:
            entry.duration_sec = options.slide_duration_sec

        _apply_sync_mode(entry, verse, options, temp_dir, sf)


def _apply_sync_mode(
    entry: SlideVideoEntry,
    verse: dict,
    options: VideoOptions,
    temp_dir: Path,
    source_file: Optional[str],
) -> None:
    """Adjust duration when verse defines sync_mode and multiple sources exist."""
    sync = verse.get("sync_mode")
    if not sync or sync not in ("avatar_lead", "notes_lead", "longest"):
        return
    if verse.get("duration_sec") is not None:
        return

    durations: List[float] = [entry.duration_sec]
    av = entry.avatar_video_path or verse.get("avatar_video_path")
    if av:
        try:
            resolved = resolve_asset_path(av, source_file=source_file) or av
            if Path(resolved).is_file():
                durations.append(ffprobe_duration(resolved))
        except Exception:
            pass

    notes = (entry.notes or _notes_text(verse)).strip()
    if notes and entry.narration_mode != "tts":
        try:
            mp3 = temp_dir / f"sync_{entry.index:03d}.mp3"
            synthesise_tts(notes, options.tts_voice, str(mp3), provider=options.tts_provider)
            durations.append(ffprobe_duration(str(mp3)) + 0.3)
        except Exception:
            pass

    if sync == "avatar_lead" and len(durations) > 1 and av:
        entry.duration_sec = durations[1]
    elif sync == "notes_lead" and len(durations) > 1:
        entry.duration_sec = durations[-1]
    elif sync == "longest":
        entry.duration_sec = max(durations)


def synthesise_tts(text: str, voice: str, output: str, *, provider: str = "edge") -> None:
    if provider == "azure":
        try:
            import azure.cognitiveservices.speech as speechsdk  # type: ignore
        except ImportError as e:
            raise RuntimeError(
                "Azure Speech not installed. Install with: pip install praisonaippt[video-tts-azure]"
            ) from e
        speech_config = speechsdk.SpeechConfig(subscription="", region="")
        synthesizer = speechsdk.SpeechSynthesizer(
            speech_config=speech_config, audio_config=speechsdk.audio.AudioOutputConfig(filename=output),
        )
        result = synthesizer.speak_text_async(text).get()
        if result.reason != speechsdk.ResultReason.SynthesizingAudioCompleted:
            raise RuntimeError(f"Azure TTS failed: {result.reason}")
        return
    if provider != "edge":
        raise RuntimeError(
            f"TTS provider {provider!r} not supported. Use edge or azure."
        )
    try:
        import edge_tts  # type: ignore
    except ImportError as e:
        raise RuntimeError(
            "edge-tts not installed. Install with: pip install praisonaippt[video-tts]"
        ) from e

    async def _run() -> None:
        communicate = edge_tts.Communicate(text, voice)
        await communicate.save(output)

    asyncio.run(_run())


def write_srt(entries: List[SlideVideoEntry], path: str) -> None:
    lines: List[str] = []
    t = 0.0
    idx = 1
    for entry in entries:
        if not entry.caption_text.strip():
            t += entry.duration_sec
            continue
        start = _srt_time(t)
        end = _srt_time(t + entry.duration_sec)
        lines.append(str(idx))
        idx += 1
        lines.append(f"{start} --> {end}")
        lines.append(entry.caption_text.strip())
        lines.append("")
        t += entry.duration_sec
    Path(path).write_text("\n".join(lines), encoding="utf-8")


def _srt_time(seconds: float) -> str:
    ms = int(round(seconds * 1000))
    h, rem = divmod(ms, 3600000)
    m, rem = divmod(rem, 60000)
    s, ms = divmod(rem, 1000)
    return f"{h:02d}:{m:02d}:{s:02d},{ms:03d}"


def rasterise_slides(
    pptx_path: str,
    temp_dir: Path,
    options: VideoOptions,
    *,
    pdf_path: Optional[str] = None,
) -> List[str]:
    prs = Presentation(pptx_path)
    slide_count = len(prs.slides)
    sr = options.slide_range
    page_first, page_last = 1, slide_count
    if sr:
        page_first, page_last = sr[0], sr[1]

    cache_dir = _slide_cache_dir(pptx_path, options.dpi) if options.slide_cache else None
    if cache_dir and cache_dir.is_dir():
        cached = sorted(cache_dir.glob("slide-*.png")) or sorted(cache_dir.glob("slide*.png"))
        if len(cached) == slide_count:
            return [str(p) for p in cached]

    pdf = pdf_path or str(temp_dir / "deck.pdf")
    if not pdf_path or not Path(pdf).is_file():
        convert_pptx_to_pdf(pptx_path, pdf, backend="auto", options=PDFOptions())

    png_dir = temp_dir / "slides"
    if sr:
        partial = pdf_to_png_pages(
            pdf, png_dir, dpi=options.dpi, first_page=page_first, last_page=page_last,
        )
        if len(partial) != page_last - page_first + 1:
            raise RuntimeError(
                f"Expected {page_last - page_first + 1} PNG pages, got {len(partial)}"
            )
        if cache_dir and options.slide_cache:
            cache_dir.mkdir(parents=True, exist_ok=True)
            for offset, page in enumerate(partial, start=page_first):
                dest = cache_dir / f"slide-{offset:03d}.png"
                shutil.copyfile(page, dest)
            cached = sorted(cache_dir.glob("slide-*.png"))
            if len(cached) == slide_count:
                return [str(p) for p in cached]
        full_pages = [""] * slide_count
        for i, page in enumerate(partial, start=page_first):
            full_pages[i - 1] = page
        return full_pages

    pages = pdf_to_png_pages(pdf, png_dir, dpi=options.dpi)
    if len(pages) != slide_count:
        raise RuntimeError(
            f"PDF page count ({len(pages)}) != slide count ({slide_count})"
        )

    if cache_dir and options.slide_cache:
        cache_dir.mkdir(parents=True, exist_ok=True)
        for src, i in zip(pages, range(1, slide_count + 1)):
            dest = cache_dir / f"slide-{i:03d}.png"
            if not dest.is_file():
                shutil.copyfile(src, dest)
        pages = [str(cache_dir / f"slide-{i:03d}.png") for i in range(1, slide_count + 1)]

    return pages


def _resolve_path(path: str, source_file: Optional[str]) -> Optional[str]:
    resolved = resolve_asset_path(path, source_file=source_file)
    p = resolved if resolved else path
    return p if p and Path(p).is_file() else None


def _overlays_for_entry(
    entry: SlideVideoEntry,
    options: VideoOptions,
    *,
    source_file: Optional[str] = None,
) -> List[OverlaySpec]:
    overlays: List[OverlaySpec] = []
    sf = source_file
    logger = logging.getLogger(__name__)

    if entry.media_box_px and entry.media_path and not entry.skip_media_overlay:
        p = _resolve_path(entry.media_path, sf)
        if p:
            box = entry.media_box_px
            m_crop = (
                entry.media_crop_y_ratio
                if entry.media_crop_y_ratio is not None
                else options.media_crop_y_ratio
            )
            m_zoom = (
                entry.media_zoom_ratio
                if entry.media_zoom_ratio is not None
                else options.media_zoom_ratio
            )
            overlays.append(
                OverlaySpec(
                    path=p,
                    x=box["x"],
                    y=box["y"],
                    width=box["width"],
                    height=box["height"],
                    is_video=is_video_path(p),
                    fit=entry.media_fit,
                    crop_y_ratio=m_crop,
                    zoom_ratio=max(1.0, m_zoom),
                )
            )
        else:
            logger.warning("Media overlay skipped; file not found: %s", entry.media_path)

    if entry.avatar_box_px and entry.avatar_video_path and not entry.skip_avatar_overlay:
        p = _resolve_path(entry.avatar_video_path, sf)
        if p:
            box = entry.avatar_box_px
            shape = entry.avatar_shape or options.avatar_shape
            crop_x = (
                entry.avatar_crop_x_ratio
                if entry.avatar_crop_x_ratio is not None
                else options.avatar_crop_x_ratio
            )
            crop_y = (
                entry.avatar_crop_y_ratio
                if entry.avatar_crop_y_ratio is not None
                else options.avatar_crop_y_ratio
            )
            zoom = (
                entry.avatar_zoom_ratio
                if entry.avatar_zoom_ratio is not None
                else options.avatar_zoom_ratio
            )
            overlays.append(
                OverlaySpec(
                    path=p,
                    x=box["x"],
                    y=box["y"],
                    width=box["width"],
                    height=box["height"],
                    is_video=True,
                    fit=options.avatar_fit,
                    shape=shape,
                    crop_x_ratio=crop_x,
                    crop_y_ratio=crop_y,
                    zoom_ratio=zoom,
                )
            )
        else:
            logger.warning("Avatar overlay skipped; file not found: %s", entry.avatar_video_path)
    return overlays


def _resolve_avatar_timeline(options: VideoOptions, entries: List[SlideVideoEntry]) -> str:
    """``auto`` → continuous when one shared avatar file spans content slides."""
    mode = options.avatar_timeline
    if mode != "auto":
        return mode
    paths = {
        e.avatar_video_path
        for e in entries
        if e.avatar_video_path and e.slide_role == "content"
    }
    if len(paths) == 1:
        return "continuous"
    return "per_slide"


def _apply_avatar_overlay_timing(
    overlays: List[OverlaySpec],
    entry: SlideVideoEntry,
    options: VideoOptions,
    avatar_offset: float,
    *,
    timeline: str,
) -> float:
    """Set continuous timeline offset and loop flags on avatar video overlays."""
    if not entry.avatar_video_path:
        return avatar_offset
    av_path = entry.avatar_video_path
    for ov in overlays:
        if not ov.is_video or ov.path != av_path:
            continue
        if timeline == "continuous":
            seek = (entry.verse or {}).get("audio_start_sec")
            if seek is not None and entry.slide_role not in ("title", "section"):
                ov.video_start_sec = float(seek)
            else:
                ov.video_start_sec = avatar_offset
        if options.loop_avatar_if_shorter:
            try:
                if ffprobe_duration(ov.path) < entry.duration_sec:
                    ov.loop_video = True
            except Exception:
                pass
    if timeline == "continuous":
        return avatar_offset + entry.duration_sec
    return avatar_offset


def compose_video(
    entries: List[SlideVideoEntry],
    png_paths: List[str],
    output: str,
    options: VideoOptions,
    temp_dir: Path,
    *,
    source_file: Optional[str] = None,
) -> str:
    encoder = options.encoder or pick_video_encoder()
    segments: List[str] = []
    pairs = list(zip(entries, png_paths))
    sr = options.slide_range
    if sr:
        lo, hi = sr
        pairs = [(e, p) for e, p in pairs if lo <= e.index + 1 <= hi]
    avatar_offset = 0.0
    timeline = _resolve_avatar_timeline(options, [e for e, _ in pairs])
    for entry, png in pairs:
        seg = str(temp_dir / f"seg_{entry.index:03d}.mp4")
        overlays = _overlays_for_entry(entry, options, source_file=source_file)
        avatar_offset = _apply_avatar_overlay_timing(
            overlays, entry, options, avatar_offset, timeline=timeline,
        )
        audio = None
        avatar_start = 0.0
        if entry.audio_primary in ("file", "tts") and entry.audio_path:
            audio = entry.audio_path
        elif entry.audio_primary == "avatar" and entry.avatar_video_path:
            audio = entry.avatar_video_path
            seek = (entry.verse or {}).get("audio_start_sec")
            if seek is not None:
                avatar_start = float(seek)
            elif timeline == "continuous":
                for ov in overlays:
                    if ov.is_video and ov.path == entry.avatar_video_path:
                        avatar_start = ov.video_start_sec
                        break
        render_slide_segment(
            png,
            entry.duration_sec,
            seg,
            fps=options.fps,
            width=options.width,
            height=options.height,
            encoder=encoder,
            overlays=overlays,
            audio_path=audio if entry.audio_primary != "avatar" else None,
            audio_start_sec=entry.audio_start_sec if entry.audio_primary in ("file", "tts") else 0.0,
        )
        if entry.audio_primary == "avatar" and entry.avatar_video_path:
            seg_av = str(temp_dir / f"seg_{entry.index:03d}_av.mp4")
            _mux_avatar_audio(
                seg, entry.avatar_video_path, seg_av, entry.duration_sec,
                start_sec=avatar_start,
            )
            seg = seg_av
        segments.append(seg)

    concat_segments(segments, output)
    return output


def _mux_avatar_audio(
    video_seg: str, avatar_path: str, output: str, duration: float,
    *, start_sec: float = 0.0,
) -> None:
    from .ffmpeg_composer import _run

    cmd = [
        "ffmpeg", "-y", "-hide_banner", "-loglevel", "error",
        "-i", video_seg,
    ]
    if start_sec > 0:
        cmd += ["-ss", f"{start_sec:.3f}"]
    cmd += [
        "-i", avatar_path,
        "-map", "0:v:0", "-map", "1:a:0",
        "-c:v", "copy", "-c:a", "aac", "-shortest",
        "-t", f"{duration:.3f}",
        output,
    ]
    proc = _run(cmd, timeout=int(duration) + 120)
    if proc.returncode != 0:
        logger.warning("Avatar audio mux failed; segment will be silent")
        shutil.copyfile(video_seg, output)


def convert_pptx_to_video(
    pptx_path: str,
    output_path: Optional[str] = None,
    *,
    data: Optional[dict] = None,
    options: Optional[VideoOptions] = None,
    pdf_path: Optional[str] = None,
    custom_title: Optional[str] = None,
) -> str:
    """Convert PPTX to MP4 using compositor backend."""
    opts = options or VideoOptions()
    if pdf_path:
        opts.pdf_path = pdf_path

    vex = (data or {}).get("video_export") or {}
    if vex and not options:
        opts = VideoOptions.from_dict(vex, data)

    resolved_backend = resolve_video_backend(opts)
    tools = check_video_tools()
    for key in required_video_tools(resolved_backend):
        if not tools[key].found:
            raise RuntimeError(
                f"Missing dependency {key}. Run: praisonaippt convert-video --check"
            )

    if resolved_backend == "powerpoint":
        from .workers.ppt_com import create_video_via_powerpoint
        out = output_path or opts.output_path or str(Path(pptx_path).with_suffix(".mp4"))
        return create_video_via_powerpoint(pptx_path, out)
    out = output_path or opts.output_path or str(Path(pptx_path).with_suffix(".mp4"))

    if not data:
        logger.warning(
            "No source YAML supplied; avatar/media PiP overlays require deck YAML with paths"
        )

    source_file = (data or {}).get("_source_file")
    prs = Presentation(pptx_path)

    temp_dir = Path(tempfile.mkdtemp(prefix="praison_video_"))
    try:
        entries = build_video_manifest(
            data, prs, opts, source_file=source_file, custom_title=custom_title,
        )
        resolve_slide_durations(entries, opts, source_file=source_file, temp_dir=temp_dir)
        pngs = rasterise_slides(pptx_path, temp_dir, opts, pdf_path=opts.pdf_path or pdf_path)
        compose_video(entries, pngs, out, opts, temp_dir, source_file=source_file)
        srt_entries = entries
        if opts.slide_range:
            lo, hi = opts.slide_range
            srt_entries = [e for e in entries if lo <= e.index + 1 <= hi]
        if opts.captions_enabled and any(e.caption_text for e in srt_entries):
            write_srt(srt_entries, str(Path(out).with_suffix(".srt")))
        if opts.keep_temp:
            manifest_path = temp_dir / "manifest.json"
            manifest_path.write_text(
                json.dumps([e.__dict__ for e in entries], indent=2),
                encoding="utf-8",
            )
            print(f"Debug artefacts kept in: {temp_dir}")
        else:
            shutil.rmtree(temp_dir, ignore_errors=True)
        return out
    except Exception:
        if not opts.keep_temp:
            shutil.rmtree(temp_dir, ignore_errors=True)
        raise


def convert_deck_to_video(
    data: dict,
    pptx_path: str,
    *,
    video_options: Optional[VideoOptions] = None,
    pdf_path: Optional[str] = None,
    custom_title: Optional[str] = None,
) -> str:
    vex = data.get("video_export") or {}
    opts = video_options or VideoOptions.from_dict(vex, data)
    out = opts.output_path or str(Path(pptx_path).with_suffix(".mp4"))
    return convert_pptx_to_video(
        pptx_path, out, data=data, options=opts, pdf_path=pdf_path,
        custom_title=custom_title,
    )
