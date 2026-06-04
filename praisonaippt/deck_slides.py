"""HeyGen-style designed deck slide layouts with colour presets."""

from __future__ import annotations

from typing import Any, Dict, List, Optional

from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

from .avatar_layouts import (
    RegionBox,
    _box_lengths,
    _draw_filled_rect,
    _hex_rgb,
    _pip_box_at,
    _place_avatar_in_box,
    _place_media_in_box,
)
from .layout_tokens import layout_in, typography_pt

DECK_SLIDE_TYPES = (
    "deck_title_split",
    "deck_exec_summary",
    "deck_split_performance",
    "deck_region_grid",
    "deck_product_columns",
    "deck_channel_analysis",
    "deck_customer_segments",
    "deck_thank_you",
    "deck_agenda",
    "deck_intro_split",
    "deck_opportunity_cards",
    "deck_forecast_split",
)

# Full-bleed rectangular avatar embeds (not circular PiP).
DECK_RECT_AVATAR_TYPES = frozenset({
    "deck_title_split",
    "deck_thank_you",
    "deck_split_performance",
    "deck_channel_analysis",
})

# Slide-level media is baked into PPTX for all deck layouts — skip FFmpeg duplicate overlay.
DECK_BAKED_MEDIA_TYPES = frozenset(DECK_SLIDE_TYPES)

# Full-panel avatar still is embedded in the PPTX — skip a second FFmpeg avatar layer.
# ``deck_thank_you`` needs a live overlay so the avatar keeps playing (still-only looks frozen).
DECK_BAKED_AVATAR_TYPES = frozenset({
    "deck_title_split",
})

# Named colour schemes (verse `color_scheme` or deck `slide_style.color_scheme`).
DECK_COLOR_PRESETS: Dict[str, Dict[str, str]] = {
    "sales_blue": {
        "background_color": "#1E40AF",
        "text_color": "#FFFFFF",
        "title_color": "#FFFFFF",
        "subtitle_color": "#93C5FD",
        "accent_color": "#2563EB",
        "badge_color": "#2563EB",
        "panel_color": "#FFFFFF",
        "panel_text_color": "#111827",
    },
    "exec_grey": {
        "background_color": "#E8EEF4",
        "text_color": "#1F2937",
        "title_color": "#111827",
        "subtitle_color": "#4B5563",
        "accent_color": "#2563EB",
        "badge_color": "#BFDBFE",
        "badge_text_color": "#1E40AF",
    },
    "split_blue": {
        "background_color": "#FFFFFF",
        "text_color": "#111827",
        "title_color": "#FFFFFF",
        "left_panel_color": "#4338CA",
        "accent_color": "#6366F1",
        "badge_color": "#E0E7FF",
        "badge_text_color": "#3730A3",
    },
    "region_navy": {
        "background_color": "#0F2744",
        "text_color": "#FFFFFF",
        "title_color": "#FFFFFF",
        "subtitle_color": "#93C5FD",
        "accent_color": "#60A5FA",
        "metric_color": "#FFFFFF",
    },
    "product_lavender": {
        "background_color": "#EDE7F6",
        "text_color": "#1F2937",
        "title_color": "#111827",
        "subtitle_color": "#4B5563",
        "accent_color": "#7C3AED",
        "metric_color": "#6D28D9",
        "badge_color": "#DDD6FE",
    },
    "channel_violet": {
        "background_color": "#FFFFFF",
        "text_color": "#111827",
        "title_color": "#FFFFFF",
        "subtitle_color": "#C4B5FD",
        "left_panel_color": "#5B52E8",
        "badge_color": "#DBEAFE",
        "badge_text_color": "#111827",
        "metric_color": "#111827",
    },
    "segments_sky": {
        "background_color": "#C6D4FF",
        "text_color": "#111827",
        "title_color": "#111827",
        "subtitle_color": "#374151",
        "metric_color": "#111827",
        "accent_color": "#2563EB",
    },
    "thank_you_blue": {
        "background_color": "#2563EB",
        "text_color": "#FFFFFF",
        "title_color": "#FFFFFF",
        "subtitle_color": "#BFDBFE",
        "panel_color": "#FFFFFF",
        "panel_text_color": "#1E40AF",
        "left_panel_color": "#2563EB",
    },
    "agenda_periwinkle": {
        "background_color": "#B0C4DE",
        "text_color": "#111827",
        "title_color": "#111827",
        "badge_color": "#FFFFFF",
        "badge_text_color": "#374151",
    },
    "intro_grey": {
        "background_color": "#F3F4F6",
        "text_color": "#111827",
        "title_color": "#111827",
        "panel_color": "#FFFFFF",
    },
    "opportunity_grey": {
        "background_color": "#F3F4F6",
        "text_color": "#111827",
        "title_color": "#111827",
        "badge_color": "#C6D8FF",
        "badge_text_color": "#374151",
    },
    "forecast_grey": {
        "background_color": "#F8F9FB",
        "text_color": "#111827",
        "title_color": "#111827",
        "badge_color": "#C6D8FF",
        "badge_text_color": "#374151",
    },
}


def _deep_merge(base: dict, override: dict) -> dict:
    out = dict(base)
    for key, val in override.items():
        if key in out and isinstance(out[key], dict) and isinstance(val, dict):
            out[key] = _deep_merge(out[key], val)
        else:
            out[key] = val
    return out


def resolve_deck_style(deck_style: dict, verse: dict, kind: str) -> dict:
    """Merge deck style, verse slide_style, and optional colour preset."""
    style = _deep_merge(dict(deck_style or {}), dict(verse.get("slide_style") or {}))
    preset = verse.get("color_scheme") or style.get("color_scheme")
    if not preset:
        preset = layout_in(style, kind, "color_scheme", None)
    if preset and preset in DECK_COLOR_PRESETS:
        style = _deep_merge(style, DECK_COLOR_PRESETS[preset])
    return style


def _layout_circle_pip(
    style: dict, kind: str, verse: Optional[dict] = None,
) -> bool:
    """True when verse or ``layouts.<kind>`` requests a circular PiP (not full-bleed rect)."""
    if verse and verse.get("avatar_shape"):
        return str(verse["avatar_shape"]).lower() in ("circle", "round", "rounded")
    raw = layout_in(style, kind, "avatar_shape", None)
    if raw is None and kind != "pip":
        raw = layout_in(style, "pip", "avatar_shape", None)
    return str(raw or "").lower() in ("circle", "round", "rounded")


def _deck_pip_anchor(style: dict, kind: str, default: str = "bottom_right") -> str:
    return str(layout_in(style, kind, "pip_position", default) or default)


def deck_avatar_shape(
    kind: str,
    style: dict,
    default: str = "circle",
    *,
    box: Optional[RegionBox] = None,
    verse: Optional[dict] = None,
) -> str:
    """Resolve FFmpeg avatar mask from YAML ``avatar_shape`` or layout-aware ``auto``."""
    from .avatar_layouts import resolve_avatar_shape, shape_for_video_overlay

    if _layout_circle_pip(style, kind, verse):
        return shape_for_video_overlay("circle")
    if kind in DECK_RECT_AVATAR_TYPES:
        return shape_for_video_overlay("h_rect")
    shape = resolve_avatar_shape(
        style or {}, layout_kind=kind, box=box, verse=verse,
    )
    if str(shape).lower() == "auto":
        shape = default
    return shape_for_video_overlay(shape)


def deck_skips_media_overlay(kind: str) -> bool:
    """Deck slides bake images in PPTX; compositor should not duplicate media overlays."""
    return kind in DECK_BAKED_MEDIA_TYPES


def deck_skips_avatar_overlay(kind: str) -> bool:
    """Full-bleed deck slides already embed the avatar still — avoid double avatar in video."""
    return kind in DECK_BAKED_AVATAR_TYPES


def _accent_rgb(style: dict, key: str = "accent_color", fallback: str = "#2563EB") -> RGBColor:
    return _hex_rgb(str(style.get(key) or style.get("highlight_color") or fallback))


def _badge_rgb(style: dict) -> RGBColor:
    return _hex_rgb(str(style.get("badge_color") or style.get("accent_color") or "#2563EB"))


def _metric_rgb(style: dict) -> RGBColor:
    return _hex_rgb(str(style.get("metric_color") or style.get("accent_color") or "#2563EB"))


def _panel_rgb(style: dict) -> RGBColor:
    return _hex_rgb(str(style.get("panel_color") or "#FFFFFF"))


def _resolve_items(raw: Any, *, label_key: str = "label") -> List[dict]:
    """Normalise verse list fields (strings or dicts) for deck protocols."""
    if not isinstance(raw, list):
        return []
    out: List[dict] = []
    for i, entry in enumerate(raw):
        if isinstance(entry, dict):
            item = dict(entry)
        else:
            item = {label_key: str(entry)}
        if not item.get("badge"):
            item["badge"] = f"{i + 1:02d}"
        if not item.get(label_key):
            item[label_key] = item.get("heading") or (
                item.get("text") if "heading" not in item and label_key != "text" else ""
            )
        out.append(item)
    return out


def _resolve_rows(raw: Any) -> List[dict]:
    """Normalise row lists for performance / channel analysis layouts."""
    if not isinstance(raw, list):
        return []
    out: List[dict] = []
    for entry in raw:
        if isinstance(entry, dict):
            out.append(dict(entry))
        else:
            out.append({"text": str(entry)})
    return out


def _media_path_from_verse(verse: dict) -> Optional[str]:
    return verse.get("media_path") or verse.get("image_path")


def _top_bottom_split(sh: float, style: dict, kind: str) -> tuple[float, float]:
    ratio = float(layout_in(style, kind, "top_height_ratio", 0.45))
    top_h = sh * ratio
    return top_h, sh - top_h


def _draw_index_badge(slide, left_in: float, top_in: float, size_in: float, label: str, style: dict) -> None:
    """White (or preset) numbered badge for agenda-style rows."""
    left, top, w, h = Inches(left_in), Inches(top_in), Inches(size_in), Inches(size_in * 0.72)
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
    try:
        shape.adjustments[0] = 0.2
    except (IndexError, AttributeError):
        pass
    shape.fill.solid()
    shape.fill.fore_color.rgb = _badge_rgb(style)
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = False
    tf.margin_left = tf.margin_right = Pt(2)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = str(label)
    p.alignment = PP_ALIGN.CENTER
    badge_text = style.get("badge_text_color") or style.get("text_color") or "#374151"
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = _hex_rgb(str(badge_text))


def _draw_percent_badge(
    slide, left_in: float, top_in: float, width_in: float, height_in: float, percent: str, style: dict
) -> float:
    left, top = Inches(left_in), Inches(top_in)
    w, h = Inches(width_in), Inches(height_in)
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
    try:
        shape.adjustments[0] = 0.18
    except (IndexError, AttributeError):
        pass
    shape.fill.solid()
    shape.fill.fore_color.rgb = _badge_rgb(style)
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = False
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = str(percent)
    p.alignment = PP_ALIGN.CENTER
    p.font.bold = True
    p.font.size = Pt(22)
    p.font.color.rgb = _hex_rgb(str(style.get("badge_text_color") or style.get("metric_color") or "#111827"))
    return width_in


def _add_dual_tone_title(
    slide, left_in: float, top_in: float, width_in: float, line1: str, line2: str, style: dict, theme: dict
) -> float:
    from .core import _write_body_paragraph

    h_in = _dual_tone_title_height(line1, line2, width_in, style)
    tb = slide.shapes.add_textbox(Inches(left_in), Inches(top_in), Inches(width_in), Inches(h_in))
    tf = tb.text_frame
    tf.word_wrap = True
    head_pt = int(typography_pt(style, "title_size_pt", 44))
    sub_pt = int(typography_pt(style, "subtitle_size_pt", 28) * 1.05)
    t1 = dict(theme)
    t1["body"] = theme.get("title") or theme.get("body")
    t2 = dict(theme)
    t2["body"] = _hex_rgb(str(style.get("subtitle_color") or "#BFDBFE"))
    p = tf.paragraphs[0]
    _write_body_paragraph(p, line1, head_pt, t1, style=style, alignment=PP_ALIGN.LEFT)
    if line2:
        p2 = tf.add_paragraph()
        _write_body_paragraph(p2, line2, sub_pt, t2, style=style, alignment=PP_ALIGN.LEFT)
    return h_in


def _add_contact_bar(slide, left_in: float, top_in: float, width_in: float, contact: str, style: dict) -> None:
    h_in = 0.55
    box = RegionBox(left_in, top_in, width_in, h_in)
    _draw_filled_rect(slide, box, _panel_rgb(style))
    tb = slide.shapes.add_textbox(Inches(left_in), Inches(top_in), Inches(width_in), Inches(h_in))
    tf = tb.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = contact
    p.alignment = PP_ALIGN.LEFT
    p.font.size = Pt(14)
    p.font.color.rgb = _hex_rgb(str(style.get("panel_text_color") or "#1E40AF"))


def _grid_area_clear_of_pip(
    sw: float,
    sh: float,
    margin: float,
    grid_top: float,
    pip: RegionBox,
    pip_anchor: str,
) -> RegionBox:
    """2×2 grid bounds that stay clear of the circular PiP (no overlap with bottom-left cells)."""
    gap = 0.18
    left, top = margin, grid_top
    right, bottom = sw - margin, sh - margin
    anchor = (pip_anchor or "bottom_left").lower().replace("-", "_")
    if pip:
        if anchor in ("bottom_left", "bl"):
            left = max(left, pip.left_in + pip.width_in + gap)
            bottom = min(bottom, pip.top_in - gap)
        elif anchor in ("bottom_right", "br"):
            right = min(right, pip.left_in - gap)
            bottom = min(bottom, pip.top_in - gap)
        elif anchor in ("top_left", "tl"):
            left = max(left, pip.left_in + pip.width_in + gap)
            top = max(top, pip.top_in + pip.height_in + gap)
        elif anchor in ("top_right", "tr"):
            right = min(right, pip.left_in - gap)
            top = max(top, pip.top_in + pip.height_in + gap)
    return RegionBox(left, top, max(0.5, right - left), max(0.5, bottom - top))


def _badge_width_in(label: str, size_in: float) -> float:
    """Wide enough for Q1-style labels on one line (square size_in is too narrow)."""
    n = len(str(label).strip()) or 1
    return max(float(size_in), 0.22 + n * 0.13)


def _draw_badge(slide, left_in: float, top_in: float, size_in: float, label: str, style: dict) -> float:
    """Rounded badge; returns drawn width in inches for row layout."""
    text = str(label)
    width_in = _badge_width_in(text, size_in)
    height_in = float(size_in) * 0.72
    left, top = Inches(left_in), Inches(top_in)
    w, h = Inches(width_in), Inches(height_in)
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
    try:
        shape.adjustments[0] = 0.22
    except (IndexError, AttributeError):
        pass
    shape.fill.solid()
    shape.fill.fore_color.rgb = _badge_rgb(style)
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = False
    tf.margin_left = tf.margin_right = Pt(2)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.CENTER
    badge_text = style.get("badge_text_color") or style.get("text_color") or "#1E40AF"
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = _hex_rgb(str(badge_text))
    return width_in


def _deck_font_name(style: dict, theme: dict) -> str:
    return str(theme.get("font_name") or style.get("font_name") or "Calibri")


def _shrink_font_pt(text: str, width_in: float, max_pt: int, min_pt: int = 16) -> int:
    from .core import _estimate_text_lines

    pt = int(max_pt)
    while pt > min_pt and _estimate_text_lines(text or "", width_in, pt) > 3:
        pt -= 2
    return max(min_pt, pt)


def _set_paragraph_font(p, pt: int, rgb: RGBColor, style: dict, theme: dict, *, bold: bool = False) -> None:
    p.font.size = Pt(pt)
    p.font.color.rgb = rgb
    p.font.bold = bold
    fn = _deck_font_name(style, theme)
    if fn:
        p.font.name = fn


def _title_block_height(title: str, subtitle: str, width_in: float, style: dict) -> float:
    """Estimated title + subtitle block height in inches (word-wrapped)."""
    from .core import _estimate_text_lines

    head_pt = _shrink_font_pt(title, width_in, int(typography_pt(style, "title_size_pt", 44)), min_pt=22)
    sub_pt = int(typography_pt(style, "subtitle_size_pt", 28))
    if subtitle:
        sub_pt = _shrink_font_pt(subtitle, width_in, sub_pt, min_pt=14)
    title_lines = _estimate_text_lines(title, width_in, head_pt)
    block_h = 0.12 + title_lines * (head_pt / 72.0) * 1.45
    if subtitle:
        sub_lines = _estimate_text_lines(subtitle, width_in, sub_pt)
        block_h += 0.1 + sub_lines * (sub_pt / 72.0) * 1.45
    return min(4.0, max(0.55, block_h))


def _dual_tone_title_height(line1: str, line2: str, width_in: float, style: dict) -> float:
    h = _title_block_height(line1, "", width_in, style)
    if (line2 or "").strip():
        sub_pt = int(typography_pt(style, "subtitle_size_pt", 28) * 1.05)
        from .core import _estimate_text_lines

        lines = _estimate_text_lines(line2, width_in, sub_pt)
        h += 0.1 + lines * (sub_pt / 72.0) * 1.45
    return min(4.0, max(0.55, h))


def _title_content_gap(style: dict, kind: str) -> float:
    return float(layout_in(style, kind, "title_content_gap_in", 0.28))


def _content_top_after_title(
    margin: float, title_h: float, kind: str, style: dict, *, min_top_key: str, min_default: float
) -> float:
    """Place body content directly below the title block (not a fixed slide floor)."""
    return margin + title_h + _title_content_gap(style, kind)


def _centre_in_band(band_top: float, band_bottom: float, block_h: float) -> float:
    """Vertically centre a block within a vertical band."""
    gap = band_bottom - band_top
    if block_h >= gap or gap <= 0:
        return band_top
    return band_top + (gap - block_h) * 0.5


def _header_block_height(header: str, width_in: float, style: dict, *, scale: float = 0.72) -> float:
    from .core import _estimate_text_lines

    if not (header or "").strip():
        return 0.0
    head_pt = _shrink_font_pt(
        header, width_in, int(typography_pt(style, "title_size_pt", 36) * scale), min_pt=18,
    )
    lines = _estimate_text_lines(header, width_in, head_pt)
    return 0.12 + lines * (head_pt / 72.0) * 1.4


def _add_title_block(
    slide,
    left_in: float,
    top_in: float,
    width_in: float,
    title: str,
    subtitle: str,
    style: dict,
    theme: dict,
    *,
    max_height_in: Optional[float] = None,
) -> float:
    from .core import _write_body_paragraph

    head_pt = _shrink_font_pt(title, width_in, int(typography_pt(style, "title_size_pt", 44)), min_pt=22)
    sub_pt = int(typography_pt(style, "subtitle_size_pt", 28))
    if subtitle:
        sub_pt = _shrink_font_pt(subtitle, width_in, sub_pt, min_pt=14)
    h_in = _title_block_height(title, subtitle, width_in, style)
    if max_height_in is not None:
        h_in = min(h_in, max(0.8, float(max_height_in)))
    tb = slide.shapes.add_textbox(Inches(left_in), Inches(top_in), Inches(width_in), Inches(h_in))
    tf = tb.text_frame
    tf.word_wrap = True
    title_theme = dict(theme)
    title_theme["body"] = theme.get("title") or theme.get("body")
    title_theme["font_name"] = _deck_font_name(style, theme)
    sub_theme = dict(theme)
    sub_theme["body"] = theme.get("subtitle") or theme.get("reference") or theme.get("body")
    sub_theme["font_name"] = title_theme["font_name"]
    p = tf.paragraphs[0]
    _write_body_paragraph(p, title, head_pt, title_theme, style=style, alignment=PP_ALIGN.LEFT)
    if subtitle:
        p2 = tf.add_paragraph()
        _write_body_paragraph(p2, subtitle, sub_pt, sub_theme, style=style, alignment=PP_ALIGN.LEFT)
    return h_in


def export_deck_slide_regions(prs, kind: str, style: dict) -> Dict[str, Optional[RegionBox]]:
    sw, sh = prs.slide_width.inches, prs.slide_height.inches
    regions: Dict[str, Optional[RegionBox]] = {
        "media": None,
        "avatar": None,
        "text_panel": None,
        "content": None,
    }
    if kind not in DECK_SLIDE_TYPES:
        return regions
    margin = float(layout_in(style, kind, "margin_in", 0.65))

    if kind == "deck_title_split":
        ratio = float(layout_in(style, kind, "avatar_width_ratio", 0.5))
        split = sw * ratio
        regions["text_panel"] = RegionBox(margin, margin, split - margin * 1.5, sh - margin * 2)
        regions["avatar"] = RegionBox(split, 0.0, sw - split, sh)
    elif kind == "deck_exec_summary":
        regions["text_panel"] = RegionBox(margin, margin, sw * 0.55, 1.4)
        regions["avatar"] = _pip_box_at(0, 0, sw, sh, style, kind, "top_right")
        col_top = float(layout_in(style, kind, "columns_top_in", 1.0))
        col_h = sh - col_top - margin
        regions["content"] = RegionBox(margin, col_top, sw - 2 * margin, col_h)
    elif kind == "deck_split_performance":
        ratio = float(layout_in(style, kind, "left_width_ratio", 0.45))
        lw = sw * ratio
        regions["text_panel"] = RegionBox(margin, margin, lw - margin * 2, 1.6)
        if _layout_circle_pip(style, kind):
            regions["avatar"] = _pip_box_at(
                0, 0, sw, sh, style, kind, _deck_pip_anchor(style, kind),
            )
        else:
            av_h = float(layout_in(style, kind, "avatar_height_in", 3.0))
            regions["avatar"] = RegionBox(0.35, sh - av_h - 0.35, lw - 0.5, av_h)
        regions["content"] = RegionBox(lw, 0.0, sw - lw, sh)
    elif kind == "deck_region_grid":
        regions["text_panel"] = RegionBox(margin, margin, sw * 0.55, 1.2)
        pip_anchor = str(layout_in(style, kind, "pip_position", "bottom_left"))
        regions["avatar"] = _pip_box_at(0, 0, sw, sh, style, kind, pip_anchor)
        grid_top = float(layout_in(style, kind, "grid_top_in", 1.55))
        regions["content"] = _grid_area_clear_of_pip(
            sw, sh, margin, grid_top, regions["avatar"], pip_anchor
        )
    elif kind == "deck_product_columns":
        regions["text_panel"] = RegionBox(margin, margin, sw * 0.55, 1.35)
        regions["avatar"] = _pip_box_at(0, 0, sw, sh, style, kind, "top_right")
        col_top = float(layout_in(style, kind, "columns_top_in", 1.85))
        regions["content"] = RegionBox(margin, col_top, sw - 2 * margin, sh - col_top - margin)
    elif kind == "deck_channel_analysis":
        ratio = float(layout_in(style, kind, "left_width_ratio", 0.45))
        lw = sw * ratio
        regions["text_panel"] = RegionBox(margin, margin, lw - margin * 2, 1.6)
        if _layout_circle_pip(style, kind):
            regions["avatar"] = _pip_box_at(
                0, 0, sw, sh, style, kind, _deck_pip_anchor(style, kind),
            )
        else:
            av_h = float(layout_in(style, kind, "avatar_height_in", 3.0))
            regions["avatar"] = RegionBox(0.35, sh - av_h - 0.35, lw - 0.5, av_h)
        regions["content"] = RegionBox(lw, 0.0, sw - lw, sh)
    elif kind == "deck_customer_segments":
        regions["text_panel"] = RegionBox(margin, margin, sw * 0.55, 1.2)
        regions["avatar"] = _pip_box_at(0, 0, sw, sh, style, kind, "top_right")
        col_top = float(layout_in(style, kind, "columns_top_in", 1.55))
        regions["content"] = RegionBox(margin, col_top, sw - 2 * margin, sh - col_top - margin)
    elif kind == "deck_thank_you":
        if _layout_circle_pip(style, kind):
            regions["text_panel"] = RegionBox(margin, margin, sw * 0.72, sh - margin * 2)
            regions["avatar"] = _pip_box_at(
                0, 0, sw, sh, style, kind, _deck_pip_anchor(style, kind),
            )
        else:
            ratio = float(layout_in(style, kind, "avatar_width_ratio", 0.5))
            split = sw * ratio
            regions["text_panel"] = RegionBox(margin, margin, split - margin * 1.5, sh - margin * 2)
            regions["avatar"] = RegionBox(split, 0.0, sw - split, sh)
    elif kind == "deck_agenda":
        list_top = float(layout_in(style, kind, "list_top_in", 1.55))
        regions["text_panel"] = RegionBox(margin, margin, sw * 0.4, 1.0)
        regions["content"] = RegionBox(margin, list_top, sw - 2 * margin, sh - list_top - margin)
    elif kind in ("deck_intro_split", "deck_forecast_split"):
        top_h, bot_h = _top_bottom_split(sh, style, kind)
        regions["text_panel"] = RegionBox(margin, margin, sw - 2 * margin, top_h - margin)
        regions["media"] = RegionBox(0.0, top_h, sw, bot_h)
    elif kind == "deck_opportunity_cards":
        regions["text_panel"] = RegionBox(margin, margin, sw * 0.45, 1.0)
        col_top = float(layout_in(style, kind, "columns_top_in", 1.25))
        regions["content"] = RegionBox(margin, col_top, sw - 2 * margin, sh - col_top - margin)
    return regions


def _render_title_split(slide, prs, verse: dict, style: dict, theme: dict, *, source_file: Optional[str]) -> None:
    sw, sh = prs.slide_width.inches, prs.slide_height.inches
    ratio = float(layout_in(style, "deck_title_split", "avatar_width_ratio", 0.5))
    split = sw * ratio
    left_bg = str(style.get("left_panel_color") or style.get("background_color") or "#1E40AF")
    _draw_filled_rect(slide, RegionBox(0, 0, split, sh), _hex_rgb(left_bg))
    title = str(verse.get("text") or verse.get("headline") or "")
    subtitle = str(verse.get("reference") or verse.get("subheader") or "")
    margin = float(layout_in(style, "deck_title_split", "margin_in", 0.75))
    _add_title_block(slide, margin, sh * 0.28, split - margin * 1.6, title, subtitle, style, theme)
    avatar_box = RegionBox(split, 0.0, sw - split, sh)
    panel_bg = str(style.get("avatar_panel_color") or "#111827")
    _place_avatar_in_box(
        slide, avatar_box, verse.get("avatar_video_path"),
        poster_path=verse.get("avatar_poster_path"), source_file=source_file, style=style,
        layout_kind="deck_title_split", verse=verse, panel_fill_rgb=_hex_rgb(panel_bg),
    )


def _render_exec_summary(slide, prs, verse: dict, style: dict, theme: dict, *, source_file: Optional[str]) -> None:
    sw, sh = prs.slide_width.inches, prs.slide_height.inches
    margin = float(layout_in(style, "deck_exec_summary", "margin_in", 0.65))
    title = str(verse.get("text") or verse.get("headline") or "")
    subtitle = str(verse.get("reference") or verse.get("subheader") or "")
    pip = _pip_box_at(0, 0, sw, sh, style, "deck_exec_summary", "top_right")
    _place_avatar_in_box(
        slide, pip, verse.get("avatar_video_path"),
        poster_path=verse.get("avatar_poster_path"), source_file=source_file, style=style,
        layout_kind="deck_exec_summary", verse=verse, draw_frame=False,
    )
    title_w = max(2.0, pip.left_in - margin - 0.25)
    title_h = _add_title_block(slide, margin, margin, title_w, title, subtitle, style, theme)
    items = _resolve_items(verse.get("items") or [])
    col_count = max(len(items), 1)
    gap = float(layout_in(style, "deck_exec_summary", "column_gap_in", 0.35))
    col_top = _content_top_after_title(
        margin, title_h, "deck_exec_summary", style, min_top_key="columns_top_in", min_default=1.0,
    )
    band_bottom = sh - margin - 0.5
    block_h = float(layout_in(style, "deck_exec_summary", "column_block_height_in", 1.35))
    col_y = _centre_in_band(col_top, band_bottom, block_h)
    usable = sw - 2 * margin - gap * (col_count - 1)
    col_w = usable / col_count
    row_w = col_count * col_w + gap * (col_count - 1)
    col_start = max(margin, (sw - row_w) / 2.0)
    body_pt = int(typography_pt(style, "body_size_pt", 22) * 0.82)
    from .core import _write_body_paragraph

    for i, item in enumerate(items[:3]):
        x = col_start + i * (col_w + gap)
        badge = str(item.get("badge") or f"{i + 1:02d}")
        badge_w = _draw_badge(slide, x, col_y, 0.42, badge, style)
        text_x = x + badge_w + 0.12
        text_w = max(0.5, col_w - badge_w - 0.12)
        tb = slide.shapes.add_textbox(
            Inches(text_x), Inches(col_y + 0.48), Inches(text_w), Inches(block_h - 0.5),
        )
        tf = tb.text_frame
        tf.word_wrap = True
        head = str(item.get("heading") or item.get("label") or "")
        body = str(item.get("text") or "")
        if head:
            p = tf.paragraphs[0]
            _write_body_paragraph(p, head, body_pt + 2, theme, style=style, alignment=PP_ALIGN.LEFT)
            p.font.bold = True
        if body:
            para = tf.add_paragraph() if head else tf.paragraphs[0]
            _write_body_paragraph(para, body, body_pt, theme, style=style, alignment=PP_ALIGN.LEFT)


def _render_split_performance(slide, prs, verse: dict, style: dict, theme: dict, *, source_file: Optional[str]) -> None:
    sw, sh = prs.slide_width.inches, prs.slide_height.inches
    ratio = float(layout_in(style, "deck_split_performance", "left_width_ratio", 0.45))
    lw = sw * ratio
    left_bg = str(
        layout_in(style, "deck_split_performance", "left_bg_color", None)
        or style.get("left_panel_color")
        or "#4338CA"
    )
    _draw_filled_rect(slide, RegionBox(0, 0, lw, sh), _hex_rgb(left_bg))
    _draw_filled_rect(slide, RegionBox(lw, 0, sw - lw, sh), _panel_rgb(style))
    left_theme = dict(theme)
    left_theme["body"] = RGBColor(0xFF, 0xFF, 0xFF)
    left_theme["title"] = RGBColor(0xFF, 0xFF, 0xFF)
    left_theme["font_name"] = _deck_font_name(style, theme)
    margin = float(layout_in(style, "deck_split_performance", "margin_in", 0.65))
    title = str(verse.get("text") or verse.get("headline") or "")
    subtitle = str(verse.get("reference") or verse.get("subheader") or "")
    av_h = float(layout_in(style, "deck_split_performance", "avatar_height_in", 3.0))
    title_w = lw - margin * 2
    av_top = sh - av_h - 0.35
    max_title_h = max(0.8, av_top - margin - 0.2)
    title_h = _add_title_block(
        slide, margin, margin, title_w, title, subtitle, style, left_theme, max_height_in=max_title_h,
    )
    if _layout_circle_pip(style, "deck_split_performance", verse):
        pip = _pip_box_at(
            0, 0, sw, sh, style, "deck_split_performance",
            _deck_pip_anchor(style, "deck_split_performance"),
        )
        _place_avatar_in_box(
            slide, pip, verse.get("avatar_video_path"),
            poster_path=verse.get("avatar_poster_path"), source_file=source_file, style=style,
            layout_kind="deck_split_performance", verse=verse, draw_frame=False,
        )
    else:
        avatar_box = RegionBox(0.35, sh - av_h - 0.35, lw - 0.5, av_h)
        _place_avatar_in_box(
            slide, avatar_box, verse.get("avatar_video_path"),
            poster_path=verse.get("avatar_poster_path"), source_file=source_file, style=style,
            layout_kind="deck_split_performance", verse=verse, panel_fill_rgb=_hex_rgb(left_bg),
        )
    header = str(verse.get("header") or "")
    rows = _resolve_rows(verse.get("rows") or [])
    rx = lw + margin
    rw = sw - lw - margin * 2
    body_rgb = theme.get("title") or theme.get("body")
    header_h = 0.0
    if header:
        header_h = _header_block_height(header, rw, style, scale=0.82)
        head_pt = _shrink_font_pt(header, rw, int(typography_pt(style, "title_size_pt", 36) * 0.82), min_pt=20)
        ht = slide.shapes.add_textbox(Inches(rx), Inches(margin), Inches(rw), Inches(header_h))
        ht.text_frame.word_wrap = True
        hp = ht.text_frame.paragraphs[0]
        hp.text = header
        _set_paragraph_font(hp, head_pt, body_rgb, style, theme, bold=True)
    panel_top = margin + header_h + (_title_content_gap(style, "deck_split_performance") if header else 0.12)
    row_h = float(layout_in(style, "deck_split_performance", "row_height_in", 0.72))
    row_gap = float(layout_in(style, "deck_split_performance", "row_gap_in", 0.28))
    n_rows = max(len(rows), 1)
    block_h = n_rows * row_h + max(0, n_rows - 1) * row_gap
    row_top = _centre_in_band(panel_top, sh - margin, block_h)
    from .core import _write_body_paragraph

    row_theme = dict(theme)
    row_theme["font_name"] = _deck_font_name(style, theme)
    body_pt = int(typography_pt(style, "body_size_pt", 22) * 0.78)

    for i, row in enumerate(rows):
        y = row_top + i * (row_h + row_gap)
        badge = str(row.get("badge") or row.get("label") or "")
        badge_w = 0.0
        if badge:
            badge_w = _draw_badge(slide, rx, y, 0.38, badge[:8], style)
        num = str(row.get("number") or row.get("metric") or "")
        desc = str(row.get("text") or "")
        nx = rx + (badge_w + 0.12 if badge_w else 0.0)
        num_w = 0.0
        if num:
            num_w = 1.2
            nt = slide.shapes.add_textbox(Inches(nx), Inches(y), Inches(num_w), Inches(0.55))
            np = nt.text_frame.paragraphs[0]
            np.text = num
            _set_paragraph_font(np, min(26, body_pt + 4), body_rgb, style, theme, bold=True)
        if desc:
            desc_x = nx + (num_w + 0.12 if num_w else 0.0)
            desc_w = max(1.0, rw - (desc_x - rx))
            dt = slide.shapes.add_textbox(Inches(desc_x), Inches(y), Inches(desc_w), Inches(row_h - 0.08))
            dt.text_frame.word_wrap = True
            dp = dt.text_frame.paragraphs[0]
            desc_pt = _shrink_font_pt(desc, desc_w, body_pt, min_pt=12)
            _write_body_paragraph(dp, desc, desc_pt, row_theme, style=style, alignment=PP_ALIGN.LEFT)


def _render_region_grid(slide, prs, verse: dict, style: dict, theme: dict, *, source_file: Optional[str]) -> None:
    sw, sh = prs.slide_width.inches, prs.slide_height.inches
    margin = float(layout_in(style, "deck_region_grid", "margin_in", 0.65))
    title = str(verse.get("text") or verse.get("headline") or "")
    title_w = sw * 0.5
    title_h = _add_title_block(slide, margin, margin, title_w, title, "", style, theme)
    if verse.get("map_path"):
        map_top = _content_top_after_title(
            margin, title_h, "deck_region_grid", style, min_top_key="grid_top_in", min_default=1.55,
        )
        map_h = 1.1
        _place_media_in_box(
            slide,
            RegionBox(margin, map_top, sw * 0.45, map_h),
            verse.get("map_path"),
            fit="contain",
            source_file=source_file,
        )
    pip_anchor = str(layout_in(style, "deck_region_grid", "pip_position", "bottom_left"))
    pip = _pip_box_at(0, 0, sw, sh, style, "deck_region_grid", pip_anchor)
    _place_avatar_in_box(
        slide, pip, verse.get("avatar_video_path"),
        poster_path=verse.get("avatar_poster_path"), source_file=source_file, style=style,
        layout_kind="deck_region_grid", verse=verse, draw_frame=False,
    )
    cells = verse.get("cells") or verse.get("columns") or []
    grid_top = _content_top_after_title(
        margin, title_h, "deck_region_grid", style, min_top_key="grid_top_in", min_default=1.55,
    )
    grid = _grid_area_clear_of_pip(sw, sh, margin, grid_top, pip, pip_anchor)
    col_gap = 0.3
    row_gap = 0.15
    cell_w = grid.width_in / 2 - col_gap / 2
    cell_h = (grid.height_in - row_gap) / 2
    from .core import _write_body_paragraph

    for i, cell in enumerate(cells[:4]):
        if not isinstance(cell, dict):
            cell = {"text": str(cell)}
        col, row = i % 2, i // 2
        x = grid.left_in + col * (cell_w + col_gap)
        y = grid.top_in + row * (cell_h + row_gap)
        num = str(cell.get("number") or cell.get("metric") or "")
        label = str(cell.get("label") or cell.get("heading") or "")
        desc = str(cell.get("text") or "")
        tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(cell_w), Inches(cell_h))
        tf = tb.text_frame
        tf.word_wrap = True
        if num:
            p = tf.paragraphs[0]
            p.text = num
            p.font.bold = True
            p.font.size = Pt(34)
            p.font.color.rgb = _metric_rgb(style)
        if label:
            para = tf.add_paragraph() if num else tf.paragraphs[0]
            _write_body_paragraph(para, label, 20, theme, style=style, alignment=PP_ALIGN.LEFT)
            para.font.bold = True
        if desc:
            para = tf.add_paragraph()
            _write_body_paragraph(para, desc, 14, theme, style=style, alignment=PP_ALIGN.LEFT)


def _render_product_columns(slide, prs, verse: dict, style: dict, theme: dict, *, source_file: Optional[str]) -> None:
    sw, sh = prs.slide_width.inches, prs.slide_height.inches
    margin = float(layout_in(style, "deck_product_columns", "margin_in", 0.65))
    title = str(verse.get("text") or verse.get("headline") or "")
    subtitle = str(verse.get("reference") or verse.get("subheader") or "")
    pip = _pip_box_at(0, 0, sw, sh, style, "deck_product_columns", "top_right")
    _place_avatar_in_box(
        slide, pip, verse.get("avatar_video_path"),
        poster_path=verse.get("avatar_poster_path"), source_file=source_file, style=style,
        layout_kind="deck_product_columns", verse=verse, draw_frame=False,
    )
    title_w = max(2.0, pip.left_in - margin - 0.25)
    title_h = _add_title_block(slide, margin, margin, title_w, title, subtitle, style, theme)
    cols = verse.get("columns") or []
    col_top = _content_top_after_title(
        margin, title_h, "deck_product_columns", style, min_top_key="columns_top_in", min_default=1.0,
    )
    gap = float(layout_in(style, "deck_product_columns", "column_gap_in", 0.25))
    count = max(len(cols), 1)
    col_w = (sw - 2 * margin - gap * (count - 1)) / count
    band_bottom = sh - margin - 0.5
    block_h = float(layout_in(style, "deck_product_columns", "column_block_height_in", 1.45))
    col_y = _centre_in_band(col_top, band_bottom, block_h)
    from .core import _write_body_paragraph

    for i, col in enumerate(cols[:4]):
        if not isinstance(col, dict):
            col = {"text": str(col)}
        x = margin + i * (col_w + gap)
        num = str(col.get("number") or col.get("metric") or "")
        label = str(col.get("label") or col.get("heading") or "")
        desc = str(col.get("text") or "")
        tb = slide.shapes.add_textbox(Inches(x), Inches(col_y), Inches(col_w), Inches(block_h))
        tf = tb.text_frame
        tf.word_wrap = True
        if num:
            p = tf.paragraphs[0]
            p.text = num
            p.font.bold = True
            p.font.size = Pt(30)
            p.font.color.rgb = _metric_rgb(style)
        if label:
            para = tf.add_paragraph() if num else tf.paragraphs[0]
            _write_body_paragraph(para, label, 18, theme, style=style, alignment=PP_ALIGN.LEFT)
            para.font.bold = True
        if desc:
            para = tf.add_paragraph()
            _write_body_paragraph(para, desc, 13, theme, style=style, alignment=PP_ALIGN.LEFT)


def _render_channel_analysis(slide, prs, verse: dict, style: dict, theme: dict, *, source_file: Optional[str]) -> None:
    sw, sh = prs.slide_width.inches, prs.slide_height.inches
    ratio = float(layout_in(style, "deck_channel_analysis", "left_width_ratio", 0.45))
    lw = sw * ratio
    left_bg = str(style.get("left_panel_color") or "#5B52E8")
    _draw_filled_rect(slide, RegionBox(0, 0, lw, sh), _hex_rgb(left_bg))
    _draw_filled_rect(slide, RegionBox(lw, 0, sw - lw, sh), _panel_rgb(style))
    margin = float(layout_in(style, "deck_channel_analysis", "margin_in", 0.65))
    line1 = str(verse.get("text") or verse.get("headline") or "")
    line2 = str(verse.get("reference") or verse.get("subheader") or "")
    title_w = lw - margin * 2
    _add_dual_tone_title(slide, margin, margin, title_w, line1, line2, style, theme)
    if _layout_circle_pip(style, "deck_channel_analysis", verse):
        pip = _pip_box_at(
            0, 0, sw, sh, style, "deck_channel_analysis",
            _deck_pip_anchor(style, "deck_channel_analysis"),
        )
        _place_avatar_in_box(
            slide, pip, verse.get("avatar_video_path"),
            poster_path=verse.get("avatar_poster_path"), source_file=source_file, style=style,
            layout_kind="deck_channel_analysis", verse=verse, draw_frame=False,
        )
    else:
        av_h = float(layout_in(style, "deck_channel_analysis", "avatar_height_in", 3.0))
        avatar_box = RegionBox(0.35, sh - av_h - 0.35, lw - 0.5, av_h)
        _place_avatar_in_box(
            slide, avatar_box, verse.get("avatar_video_path"),
            poster_path=verse.get("avatar_poster_path"), source_file=source_file, style=style,
            layout_kind="deck_channel_analysis", verse=verse, panel_fill_rgb=_hex_rgb(left_bg),
        )
    header = str(verse.get("header") or "")
    rows = _resolve_rows(verse.get("rows") or [])
    rx = lw + margin
    rw = sw - lw - margin * 2
    header_h = 0.0
    if header:
        header_h = _header_block_height(header, rw, style, scale=0.75)
        ht = slide.shapes.add_textbox(Inches(rx), Inches(margin), Inches(rw), Inches(header_h))
        ht.text_frame.word_wrap = True
        hp = ht.text_frame.paragraphs[0]
        hp.text = header
        hp.font.bold = True
        hp.font.size = Pt(int(typography_pt(style, "title_size_pt", 36) * 0.75))
        hp.font.color.rgb = theme.get("title") or theme.get("body")
    panel_top = margin + header_h + (_title_content_gap(style, "deck_channel_analysis") if header else 0.12)
    row_h = float(layout_in(style, "deck_channel_analysis", "row_height_in", 0.88))
    row_gap = float(layout_in(style, "deck_channel_analysis", "row_gap_in", 0.32))
    badge_h = float(layout_in(style, "deck_channel_analysis", "badge_height_in", 0.55))
    badge_max_w = float(layout_in(style, "deck_channel_analysis", "badge_width_in", 0.95))
    n_rows = max(len(rows), 1)
    block_h = n_rows * row_h + max(0, n_rows - 1) * row_gap
    row_top = _centre_in_band(panel_top, sh - margin, block_h)
    from .core import _write_body_paragraph

    for i, row in enumerate(rows):
        y = row_top + i * (row_h + row_gap)
        num = str(row.get("number") or row.get("metric") or "")
        label = str(row.get("label") or row.get("heading") or "")
        desc = str(row.get("text") or "")
        badge_w = 0.0
        if num:
            fit_w = max(0.55, min(badge_max_w, _badge_width_in(num, badge_h) + 0.1))
            badge_w = _draw_percent_badge(slide, rx, y, fit_w, badge_h, num, style)
        tx = rx + badge_w + 0.15
        text_w = max(1.0, rw - (tx - rx))
        if label or desc:
            dt = slide.shapes.add_textbox(Inches(tx), Inches(y), Inches(text_w), Inches(row_h - 0.08))
            tf = dt.text_frame
            tf.word_wrap = True
            if label:
                lp = tf.paragraphs[0]
                _write_body_paragraph(lp, label, 15, theme, style=style, alignment=PP_ALIGN.LEFT)
                lp.font.bold = True
            if desc:
                para = tf.add_paragraph() if label else tf.paragraphs[0]
                _write_body_paragraph(para, desc, 13, theme, style=style, alignment=PP_ALIGN.LEFT)


def _render_customer_segments(slide, prs, verse: dict, style: dict, theme: dict, *, source_file: Optional[str]) -> None:
    sw, sh = prs.slide_width.inches, prs.slide_height.inches
    margin = float(layout_in(style, "deck_customer_segments", "margin_in", 0.65))
    title = str(verse.get("text") or verse.get("headline") or "")
    pip = _pip_box_at(0, 0, sw, sh, style, "deck_customer_segments", "top_right")
    _place_avatar_in_box(
        slide, pip, verse.get("avatar_video_path"),
        poster_path=verse.get("avatar_poster_path"), source_file=source_file, style=style,
        layout_kind="deck_customer_segments", verse=verse, draw_frame=False,
    )
    title_w = max(2.0, pip.left_in - margin - 0.25)
    title_h = _add_title_block(slide, margin, margin, title_w, title, "", style, theme)
    cols = verse.get("columns") or []
    min_col_top = float(layout_in(style, "deck_customer_segments", "columns_top_in", 1.55))
    col_top = _content_top_after_title(
        margin, title_h, "deck_customer_segments", style, min_top_key="columns_top_in", min_default=min_col_top,
    )
    gap = float(layout_in(style, "deck_customer_segments", "column_gap_in", 0.35))
    count = max(len(cols), 1)
    col_w = (sw - 2 * margin - gap * (count - 1)) / count
    from .core import _write_body_paragraph

    for i, col in enumerate(cols[:3]):
        if not isinstance(col, dict):
            col = {"text": str(col)}
        x = margin + i * (col_w + gap)
        num = str(col.get("number") or col.get("metric") or "")
        label = str(col.get("label") or col.get("heading") or "")
        desc = str(col.get("text") or "")
        tb = slide.shapes.add_textbox(Inches(x), Inches(col_top), Inches(col_w), Inches(sh - col_top - margin))
        tf = tb.text_frame
        tf.word_wrap = True
        if num:
            p = tf.paragraphs[0]
            p.text = num
            p.font.bold = True
            p.font.size = Pt(36)
            p.font.color.rgb = _metric_rgb(style)
        if label:
            para = tf.add_paragraph() if num else tf.paragraphs[0]
            _write_body_paragraph(para, label, 20, theme, style=style, alignment=PP_ALIGN.LEFT)
            para.font.bold = True
        if desc:
            para = tf.add_paragraph()
            _write_body_paragraph(para, desc, 14, theme, style=style, alignment=PP_ALIGN.LEFT)


def _render_thank_you(slide, prs, verse: dict, style: dict, theme: dict, *, source_file: Optional[str]) -> None:
    sw, sh = prs.slide_width.inches, prs.slide_height.inches
    left_bg = str(style.get("left_panel_color") or style.get("background_color") or "#2563EB")
    margin = float(layout_in(style, "deck_thank_you", "margin_in", 0.75))
    line1 = str(verse.get("text") or verse.get("headline") or "THANK")
    line2 = str(verse.get("reference") or verse.get("subheader") or "YOU")
    contact = str(verse.get("contact") or verse.get("email") or "")

    if _layout_circle_pip(style, "deck_thank_you", verse):
        _draw_filled_rect(slide, RegionBox(0, 0, sw, sh), _hex_rgb(left_bg))
        title_top = sh * 0.22
        title_w = sw * 0.72 - margin * 1.6
        title_h = _add_dual_tone_title(slide, margin, title_top, title_w, line1, line2, style, theme)
        if contact:
            _add_contact_bar(slide, margin, title_top + title_h + 0.35, title_w, contact, style)
        pip = _pip_box_at(
            0, 0, sw, sh, style, "deck_thank_you", _deck_pip_anchor(style, "deck_thank_you"),
        )
        _place_avatar_in_box(
            slide, pip, verse.get("avatar_video_path"),
            poster_path=verse.get("avatar_poster_path"), source_file=source_file, style=style,
            layout_kind="deck_thank_you", verse=verse, draw_frame=False,
        )
        return

    ratio = float(layout_in(style, "deck_thank_you", "avatar_width_ratio", 0.5))
    split = sw * ratio
    _draw_filled_rect(slide, RegionBox(0, 0, split, sh), _hex_rgb(left_bg))
    _draw_filled_rect(slide, RegionBox(split, 0, sw - split, sh), _hex_rgb(left_bg))
    title_top = sh * 0.22
    title_w = split - margin * 1.6
    title_h = _add_dual_tone_title(slide, margin, title_top, title_w, line1, line2, style, theme)
    if contact:
        _add_contact_bar(slide, margin, title_top + title_h + 0.35, split - margin * 2, contact, style)
    avatar_box = RegionBox(split, 0.0, sw - split, sh)
    _place_avatar_in_box(
        slide, avatar_box, verse.get("avatar_video_path"),
        poster_path=verse.get("avatar_poster_path"), source_file=source_file, style=style,
        layout_kind="deck_thank_you", verse=verse, panel_fill_rgb=_hex_rgb(left_bg),
    )


def _render_agenda(slide, prs, verse: dict, style: dict, theme: dict, *, source_file: Optional[str]) -> None:
    sw, sh = prs.slide_width.inches, prs.slide_height.inches
    margin = float(layout_in(style, "deck_agenda", "margin_in", 0.65))
    title = str(verse.get("text") or verse.get("headline") or "Agenda")
    title_w = sw * 0.4
    title_h = _add_title_block(slide, margin, margin, title_w, title, "", style, theme)
    items = _resolve_items(verse.get("items") or verse.get("agenda") or [])
    col_count = int(layout_in(style, "deck_agenda", "agenda_columns", 2))
    agenda_list_top = float(layout_in(style, "deck_agenda", "list_top_in", 1.55))
    list_top = _content_top_after_title(
        margin, title_h, "deck_agenda", style, min_top_key="list_top_in", min_default=agenda_list_top,
    )
    min_row_h = float(layout_in(style, "deck_agenda", "row_height_in", 0.52))
    rows = max(1, (len(items) + col_count - 1) // col_count)
    available = sh - list_top - margin
    row_h = max(min_row_h, available / rows)
    badge_w = float(layout_in(style, "deck_agenda", "badge_width_in", 0.42))
    badge_h = badge_w * 0.72
    gap = float(layout_in(style, "deck_agenda", "column_gap_in", 0.55))
    usable_w = sw - 2 * margin - gap * (col_count - 1)
    col_w = usable_w / col_count
    from .core import _write_body_paragraph

    for i, item in enumerate(items):
        col, row = i % col_count, i // col_count
        x = margin + col * (col_w + gap)
        y = list_top + row * row_h
        badge_y = y + max(0.0, (row_h - badge_h) * 0.12)
        badge = str(item.get("badge") or f"{i + 1:02d}")
        label = str(item.get("label") or item.get("heading") or item.get("text") or "")
        _draw_index_badge(slide, x, badge_y, badge_w, badge, style)
        tb = slide.shapes.add_textbox(
            Inches(x + badge_w + 0.22), Inches(y + 0.06), Inches(col_w - badge_w - 0.28), Inches(row_h - 0.08)
        )
        p = tb.text_frame.paragraphs[0]
        _write_body_paragraph(p, label, 16, theme, style=style, alignment=PP_ALIGN.LEFT)
        p.font.bold = True


def _render_intro_split(slide, prs, verse: dict, style: dict, theme: dict, *, source_file: Optional[str]) -> None:
    sw, sh = prs.slide_width.inches, prs.slide_height.inches
    margin = float(layout_in(style, "deck_intro_split", "margin_in", 0.65))
    top_h, bot_h = _top_bottom_split(sh, style, "deck_intro_split")
    panel_bg = str(style.get("panel_color") or style.get("background_color") or "#F3F4F6")
    _draw_filled_rect(slide, RegionBox(0, 0, sw, top_h), _hex_rgb(panel_bg))
    title = str(verse.get("text") or verse.get("headline") or "")
    body = str(
        verse.get("reference") or verse.get("subheader") or verse.get("body") or verse.get("description") or ""
    )
    title_w = float(layout_in(style, "deck_intro_split", "title_width_ratio", 0.38)) * sw
    _add_title_block(slide, margin, margin + 0.15, title_w - margin, title, "", style, theme)
    if body:
        from .core import _write_body_paragraph

        body_x = title_w + 0.15
        body_w = sw - body_x - margin
        tb = slide.shapes.add_textbox(Inches(body_x), Inches(margin + 0.2), Inches(body_w), Inches(top_h - margin * 1.5))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        body_pt = int(typography_pt(style, "body_size_pt", 22) * 0.75)
        _write_body_paragraph(p, body, body_pt, theme, style=style, alignment=PP_ALIGN.LEFT)
    media = _media_path_from_verse(verse)
    fit = str(verse.get("media_fit") or verse.get("image_fit") or layout_in(style, "deck_intro_split", "media_fit", "cover"))
    _place_media_in_box(
        slide, RegionBox(0, top_h, sw, bot_h), media, fit=fit, source_file=source_file,
    )


def _render_opportunity_cards(slide, prs, verse: dict, style: dict, theme: dict, *, source_file: Optional[str]) -> None:
    sw, sh = prs.slide_width.inches, prs.slide_height.inches
    margin = float(layout_in(style, "deck_opportunity_cards", "margin_in", 0.65))
    title = str(verse.get("text") or verse.get("headline") or "")
    title_w = sw * 0.45
    title_h = _add_title_block(slide, margin, margin, title_w, title, "", style, theme)
    cards = _resolve_items(verse.get("columns") or verse.get("items") or [])
    col_top = _content_top_after_title(
        margin, title_h, "deck_opportunity_cards", style, min_top_key="columns_top_in", min_default=1.25,
    )
    gap = float(layout_in(style, "deck_opportunity_cards", "column_gap_in", 0.3))
    img_h = float(layout_in(style, "deck_opportunity_cards", "image_height_in", 1.35))
    count = max(len(cards), 1)
    col_w = (sw - 2 * margin - gap * (count - 1)) / min(count, 3)
    from .core import _write_body_paragraph

    for i, card in enumerate(cards[:3]):
        x = margin + i * (col_w + gap)
        badge = str(card.get("badge") or f"{i + 1:02d}")
        _draw_badge(slide, x, col_top, 0.38, badge, style)
        img_path = card.get("image_path") or card.get("media_path")
        img_top = col_top + 0.48
        if img_path:
            _place_media_in_box(
                slide, RegionBox(x, img_top, col_w, img_h), img_path,
                fit=str(card.get("image_fit") or card.get("media_fit") or "cover"),
                source_file=source_file,
            )
        text_top = img_top + img_h + 0.12
        heading = str(card.get("heading") or card.get("label") or "")
        desc = str(card.get("text") or card.get("description") or "")
        tb = slide.shapes.add_textbox(Inches(x), Inches(text_top), Inches(col_w), Inches(sh - text_top - margin))
        tf = tb.text_frame
        tf.word_wrap = True
        if heading:
            p = tf.paragraphs[0]
            _write_body_paragraph(p, heading, 18, theme, style=style, alignment=PP_ALIGN.LEFT)
            p.font.bold = True
        if desc:
            para = tf.add_paragraph() if heading else tf.paragraphs[0]
            _write_body_paragraph(para, desc, 13, theme, style=style, alignment=PP_ALIGN.LEFT)


def _render_forecast_split(slide, prs, verse: dict, style: dict, theme: dict, *, source_file: Optional[str]) -> None:
    sw, sh = prs.slide_width.inches, prs.slide_height.inches
    margin = float(layout_in(style, "deck_forecast_split", "margin_in", 0.65))
    top_h, bot_h = _top_bottom_split(sh, style, "deck_forecast_split")
    panel_bg = str(style.get("panel_color") or style.get("background_color") or "#F8F9FB")
    _draw_filled_rect(slide, RegionBox(0, 0, sw, top_h), _hex_rgb(panel_bg))
    title = str(verse.get("text") or verse.get("headline") or "")
    title_w = sw * 0.35
    title_h = _add_title_block(slide, margin, margin, title_w, title, "", style, theme)
    items = _resolve_items(verse.get("items") or [])
    col_top = _content_top_after_title(
        margin, title_h, "deck_forecast_split", style, min_top_key="items_top_in", min_default=1.15,
    )
    gap = float(layout_in(style, "deck_forecast_split", "column_gap_in", 0.35))
    count = max(len(items), 1)
    col_w = (sw - 2 * margin - gap * (count - 1)) / min(count, 3)
    from .core import _write_body_paragraph

    for i, item in enumerate(items[:3]):
        x = margin + i * (col_w + gap)
        badge = str(item.get("badge") or f"{i + 1:02d}")
        _draw_badge(slide, x, col_top, 0.38, badge, style)
        body = str(item.get("text") or item.get("label") or item.get("heading") or "")
        tb = slide.shapes.add_textbox(
            Inches(x), Inches(col_top + 0.48), Inches(col_w), Inches(top_h - col_top - 0.55)
        )
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        body_pt = int(typography_pt(style, "body_size_pt", 22) * 0.72)
        _write_body_paragraph(p, body, body_pt, theme, style=style, alignment=PP_ALIGN.LEFT)
    media = _media_path_from_verse(verse)
    fit = str(verse.get("media_fit") or layout_in(style, "deck_forecast_split", "media_fit", "cover"))
    _place_media_in_box(
        slide, RegionBox(0, top_h, sw, bot_h), media, fit=fit, source_file=source_file,
    )


def render_deck_slide(
    prs, kind: str, verse: dict, deck_style: Optional[dict] = None, *, source_file: Optional[str] = None
):
    from .core import _apply_slide_background, _resolve_theme

    if kind not in DECK_SLIDE_TYPES:
        raise ValueError(f"Unknown deck slide kind: {kind}")
    style = resolve_deck_style(deck_style or {}, verse, kind)
    if source_file:
        style["_source_file"] = source_file
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _apply_slide_background(slide, style, prs)
    theme = _resolve_theme(style)
    if kind == "deck_title_split":
        _render_title_split(slide, prs, verse, style, theme, source_file=source_file)
    elif kind == "deck_exec_summary":
        _render_exec_summary(slide, prs, verse, style, theme, source_file=source_file)
    elif kind == "deck_split_performance":
        _render_split_performance(slide, prs, verse, style, theme, source_file=source_file)
    elif kind == "deck_region_grid":
        _render_region_grid(slide, prs, verse, style, theme, source_file=source_file)
    elif kind == "deck_product_columns":
        _render_product_columns(slide, prs, verse, style, theme, source_file=source_file)
    elif kind == "deck_channel_analysis":
        _render_channel_analysis(slide, prs, verse, style, theme, source_file=source_file)
    elif kind == "deck_customer_segments":
        _render_customer_segments(slide, prs, verse, style, theme, source_file=source_file)
    elif kind == "deck_thank_you":
        _render_thank_you(slide, prs, verse, style, theme, source_file=source_file)
    elif kind == "deck_agenda":
        _render_agenda(slide, prs, verse, style, theme, source_file=source_file)
    elif kind == "deck_intro_split":
        _render_intro_split(slide, prs, verse, style, theme, source_file=source_file)
    elif kind == "deck_opportunity_cards":
        _render_opportunity_cards(slide, prs, verse, style, theme, source_file=source_file)
    elif kind == "deck_forecast_split":
        _render_forecast_split(slide, prs, verse, style, theme, source_file=source_file)
    return slide
