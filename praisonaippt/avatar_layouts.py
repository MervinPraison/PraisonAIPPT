"""Avatar and media layout slides (speaking-head video + media regions)."""

from __future__ import annotations

import base64
import io
import os
import subprocess
import tempfile
from dataclasses import dataclass
from typing import Dict, Optional, Tuple

from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

from .layout_tokens import layout_in, typography_pt
from .text_panel_anchors import HERO_PANEL_ANCHORS, TEXT_PANEL_ANCHORS
from .utils import resolve_asset_path

AVATAR_SLIDE_TYPES = (
    "avatar_only",
    "media_only",
    "avatar_media_1",
    "avatar_media_2",
    "avatar_media_3",
    "avatar_name_card",
    "avatar_headline",
    "avatar_headline_full",
    "avatar_quote",
    "avatar_border",
    "media_border",
    "avatar_media_border_1",
    "avatar_media_border_2",
    "avatar_media_border_3",
    "avatar_intro",
    "avatar_outro",
)

_BORDER_KINDS = frozenset({
    "avatar_border",
    "media_border",
    "avatar_media_border_1",
    "avatar_media_border_2",
    "avatar_media_border_3",
})

_SPLIT_KINDS = frozenset({
    "avatar_media_1",
    "avatar_media_2",
    "avatar_media_border_1",
    "avatar_media_border_2",
})

_PIP_KINDS = frozenset({
    "avatar_media_3",
    "avatar_media_border_3",
    "avatar_quote",
})

# Live PiP is composited in video export; baking a still here duplicates FFmpeg overlay.
_AVATAR_PIP_VIDEO_OVERLAY_ONLY = frozenset({"avatar_quote", "avatar_media_3"})

_VIDEO_EXTS = {".mp4": "video/mp4", ".mov": "video/quicktime", ".m4v": "video/mp4", ".webm": "video/webm"}
_IMAGE_EXTS = {".png", ".jpg", ".jpeg", ".gif", ".webp", ".bmp", ".tif", ".tiff"}

_GREY_POSTER_PNG = base64.b64decode(
    "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAQAAAC1HAwCAAAAC0lEQVR42mNk+A8AAQUBAScY42YAAAAASUVORK5CYII="
)

_AVATAR_GREY = RGBColor(0xB8, 0xB8, 0xB8)
_PIP_BACKDROP_GREY = RGBColor(0xC8, 0xC8, 0xC8)
_MEDIA_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
_PANEL_NAVY = RGBColor(0x1E, 0x3A, 0x5F)

_INTRO_COLOURS = {
    "cream": RGBColor(0xF5, 0xF0, 0xE6),
    "tan": RGBColor(0xC4, 0xA8, 0x82),
    "navy": RGBColor(0x1E, 0x3A, 0x5F),
    "sage": RGBColor(0x8F, 0xA8, 0x88),
}


@dataclass(frozen=True)
class RegionBox:
    left_in: float
    top_in: float
    width_in: float
    height_in: float
    rounded: bool = False
    corner_radius_in: float = 0.12


def _hex_rgb(hex_color: str) -> RGBColor:
    h = (hex_color or "#1E3A5F").lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def _border_inset(style: dict, kind: str) -> float:
    return float(layout_in(style, kind, "border_inset_in", layout_in(style, "avatar_border", "border_inset_in", 0.25)))


def _content_area(prs, style: dict, kind: str) -> Tuple[float, float, float, float]:
    """Return (left_in, top_in, width_in, height_in) for inner content."""
    sw, sh = prs.slide_width.inches, prs.slide_height.inches
    if kind in _BORDER_KINDS:
        inset = _border_inset(style, kind)
        return inset, inset, sw - 2 * inset, sh - 2 * inset
    return 0.0, 0.0, sw, sh


def _split_boxes(
    cx: float, cy: float, cw: float, ch: float, media_ratio: float, gap_in: float
) -> Tuple[RegionBox, RegionBox]:
    gap = max(gap_in, 0.0)
    usable = cw - gap
    media_w = usable * media_ratio
    avatar_w = usable - media_w
    media = RegionBox(cx, cy, media_w, ch)
    avatar = RegionBox(cx + media_w + gap, cy, avatar_w, ch)
    return media, avatar


def _pip_box(cx: float, cy: float, cw: float, ch: float, style: dict, kind: str) -> RegionBox:
    anchor = str(
        layout_in(style, kind, "pip_position", layout_in(style, "pip", "position", "bottom_right"))
    ).lower()
    return _pip_box_at(cx, cy, cw, ch, style, kind, anchor)


def _pip_box_at(
    cx: float, cy: float, cw: float, ch: float, style: dict, kind: str, anchor: str
) -> RegionBox:
    ratio_raw = layout_in(style, kind, "pip_width_ratio", None)
    margin_raw = layout_in(style, kind, "pip_margin_in", None)
    ratio = float(ratio_raw if ratio_raw is not None else layout_in(style, "pip", "width_ratio", 0.20))
    margin = float(margin_raw if margin_raw is not None else layout_in(style, "pip", "margin_in", 0.38))
    size = cw * ratio
    shape = _pip_shape_kind(style, kind)
    rounded = shape in ("circle", "round", "rounded")
    radius = size / 2 if shape == "circle" else float(layout_in(style, "pip", "corner_radius_in", 0.12))
    anchor = (anchor or "bottom_right").lower().replace("-", "_")
    if anchor in ("top_right", "tr"):
        left, top = cx + cw - size - margin, cy + margin
    elif anchor in ("bottom_left", "bl"):
        left, top = cx + margin, cy + ch - size - margin
    elif anchor in ("top_left", "tl"):
        left, top = cx + margin, cy + margin
    else:
        left, top = cx + cw - size - margin, cy + ch - size - margin
    return RegionBox(left, top, size, size, rounded=rounded, corner_radius_in=radius)


def export_floating_pip_box(prs, style: dict) -> RegionBox:
    """Bottom-right PiP box used for avatar overlays on any slide type."""
    cx, cy, cw, ch = 0.0, 0.0, prs.slide_width.inches, prs.slide_height.inches
    return _pip_box(cx, cy, cw, ch, style, "pip")


AVATAR_SHAPE_VALUES = frozenset({
    "auto", "circle", "round", "rounded", "square", "rect", "rectangle",
    "h_rect", "horizontal", "wide", "v_rect", "vertical", "tall",
})

_PIP_LAYOUT_KINDS = frozenset({
    "deck_exec_summary", "deck_region_grid", "deck_product_columns",
    "deck_customer_segments",
})

_STRIP_LAYOUT_KINDS = frozenset({"deck_split_performance", "deck_channel_analysis"})


def _shape_from_box_aspect(box: RegionBox) -> str:
    ar = box.width_in / max(box.height_in, 0.01)
    if ar >= 1.2:
        return "h_rect"
    if ar <= 0.82:
        return "v_rect"
    return "square"


def default_avatar_shape_for_layout(kind: str, box: Optional[RegionBox]) -> str:
    """Layout-aware default when ``avatar_shape: auto``."""
    try:
        from .deck_slides import DECK_RECT_AVATAR_TYPES

        if kind in DECK_RECT_AVATAR_TYPES:
            return _shape_from_box_aspect(box) if box else "h_rect"
    except ImportError:
        pass
    if kind in _STRIP_LAYOUT_KINDS:
        return "h_rect"
    if kind in _PIP_LAYOUT_KINDS:
        return "circle"
    if box:
        return _shape_from_box_aspect(box)
    return "circle"


def resolve_avatar_shape(
    style: dict,
    *,
    layout_kind: str = "pip",
    box: Optional[RegionBox] = None,
    verse: Optional[dict] = None,
) -> str:
    """verse.avatar_shape > layouts.<kind>.avatar_shape > layouts.pip > layout default."""
    if verse and verse.get("avatar_shape"):
        raw = str(verse["avatar_shape"]).lower().strip()
    else:
        raw = layout_in(style, layout_kind, "avatar_shape", None)
        if raw is None and layout_kind != "pip":
            raw = layout_in(style, "pip", "avatar_shape", None)
        if raw is None:
            raw = layout_in(style, layout_kind, "shape", None) or layout_in(
                style, "pip", "shape", None,
            )
        raw = str(raw or "auto").lower().strip() if raw is not None else "auto"
    if raw == "auto":
        return default_avatar_shape_for_layout(layout_kind, box)
    return raw


def shape_uses_circle_mask(shape: str) -> bool:
    return str(shape).lower() in ("circle", "round", "rounded")


def shape_for_video_overlay(shape: str) -> str:
    """FFmpeg overlay mask: circle vs rectangular."""
    s = str(shape).lower()
    if shape_uses_circle_mask(s):
        return "circle"
    return "rect"


def _pip_shape_kind(style: dict, kind: str = "pip") -> str:
    """Legacy PiP shape from ``shape`` / ``pip_shape`` keys."""
    raw = layout_in(style, kind, "pip_shape", None) or layout_in(style, kind, "shape", None)
    if raw is None and kind != "pip":
        raw = layout_in(style, "pip", "shape", "circle")
    return str(raw or "circle").lower()


def _content_beside_pip(
    cx: float, cy: float, cw: float, ch: float, pip: RegionBox, margin: float
) -> RegionBox:
    """Text area that leaves the bottom-right PiP corner clear."""
    text_w = max(cw - pip.width_in - 2 * margin, 1.0)
    text_h = max(pip.top_in - cy - margin, 1.0)
    return RegionBox(cx + margin, cy + margin, text_w, text_h)


def _text_panel_box(prs, style: dict, kind: str, position: str) -> RegionBox:
    cx, cy, cw, ch = _content_area(prs, style, kind)
    margin = float(layout_in(style, kind, "panel_margin_in", 0.35))
    pw = cw * float(layout_in(style, kind, "panel_width_ratio", 0.42))
    ph = float(layout_in(style, kind, "panel_height_in", 1.2))
    if position == "top":
        return RegionBox(cx + margin, cy + margin, pw, ph)
    return RegionBox(cx + margin, cy + ch - ph - margin, pw, ph)


_TEXT_ANCHORS = HERO_PANEL_ANCHORS
_HERO_LAYOUTS = frozenset({"stacked", "full_bleed"})
_TEXT_STYLES = frozenset({"navy_panel", "overlay", "semi_panel"})


def _hero_layout_mode(style: dict, verse: Optional[dict], kind: str) -> str:
    tp = (verse or {}).get("text_panel") if isinstance((verse or {}).get("text_panel"), dict) else {}
    raw = tp.get("hero_layout") or layout_in(style, kind, "hero_layout", "stacked")
    mode = str(raw).lower().strip()
    return mode if mode in _HERO_LAYOUTS else "stacked"


def _text_style_mode(style: dict, verse: Optional[dict], kind: str) -> str:
    tp = (verse or {}).get("text_panel") if isinstance((verse or {}).get("text_panel"), dict) else {}
    raw = tp.get("style") or layout_in(style, kind, "text_style", "navy_panel")
    mode = str(raw).lower().strip()
    return mode if mode in _TEXT_STYLES else "navy_panel"


def _verse_text_panel_cfg(style: dict, verse: Optional[dict], kind: str) -> dict:
    tp = (verse or {}).get("text_panel") if isinstance((verse or {}).get("text_panel"), dict) else {}
    anchor = str(tp.get("anchor") or layout_in(style, kind, "text_anchor", "top_left")).lower().strip()
    if anchor == "auto":
        resolved = (verse or {}).get("_hero_panel_anchor")
        anchor = str(resolved).lower().strip() if resolved else "top_left"
    if anchor not in _TEXT_ANCHORS:
        anchor = "top_left"
    return {
        "anchor": anchor,
        "width_ratio": float(
            tp.get("width_ratio") or layout_in(style, kind, "panel_width_ratio", 0.42)
        ),
        "height_in": float(tp.get("height_in") or layout_in(style, kind, "panel_height_in", 0.9)),
        "margin_in": float(tp.get("margin_in") or layout_in(style, kind, "panel_margin_in", 0.35)),
        "max_width_ratio": tp.get("max_width_ratio"),
    }


def _estimate_panel_height(headline: str, subheader: str, width_in: float) -> float:
    """Minimum panel height (inches) so headline + subheader stay inside the navy box."""
    w = max(2.4, float(width_in))
    chars = max(12, int(w * 5.6))
    head = (headline or "").strip()
    sub = (subheader or "").strip()
    head_lines = max(1, (len(head) + chars - 1) // chars)
    sub_lines = max(1, (len(sub) + chars - 1) // chars) if sub else 0
    return 0.26 + head_lines * 0.40 + sub_lines * 0.30 + 0.24


def _text_panel_box_anchored(
    prs,
    style: dict,
    kind: str,
    cfg: dict,
    *,
    pip: Optional[RegionBox] = None,
    verse: Optional[dict] = None,
) -> RegionBox:
    cx, cy, cw, ch = _content_area(prs, style, kind)
    margin = float(cfg["margin_in"])
    pw = cw * float(cfg["width_ratio"])
    if cfg.get("max_width_ratio") is not None:
        pw = min(pw, cw * float(cfg["max_width_ratio"]))
    ph = float(cfg["height_in"])
    if verse is not None:
        ph = max(
            ph,
            _estimate_panel_height(
                str(verse.get("headline") or ""),
                str(verse.get("subheader") or ""),
                pw,
            ),
        )
    anchor = cfg["anchor"]
    pip_gap = float(layout_in(style, kind, "text_pip_gap_in", 0.12))
    pip_left = pip.left_in - pip_gap if pip else cx + cw

    if anchor == "top_left":
        left, top = cx + margin, cy + margin
    elif anchor == "top":
        left, top = cx + (cw - pw) / 2, cy + margin
    elif anchor == "top_right":
        left, top = cx + cw - pw - margin, cy + margin
    elif anchor == "bottom_left":
        left = cx + margin
        top = max(cy + margin, cy + ch - ph - margin)
        if pip and left + pw > pip_left and top + ph > pip.top_in - pip_gap:
            top = min(top, pip.top_in - ph - pip_gap)
    elif anchor == "bottom_right":
        left = cx + cw - pw - margin
        top = max(cy + margin, cy + ch - ph - margin)
        if pip and left < pip_left + pip.width_in + pip_gap:
            left = max(cx + margin, pip_left - pw - pip_gap)
    else:
        left, top = cx + margin, cy + ch - ph - margin

    box = RegionBox(
        left, top, max(1.0, pw), max(0.55, ph),
        rounded=True, corner_radius_in=0.08,
    )
    if pip:
        gap = pip_gap
        while (
            box.left_in + box.width_in > pip.left_in - gap
            and box.top_in + box.height_in > pip.top_in - gap
            and box.left_in < pip.left_in + pip.width_in + gap
            and box.top_in < pip.top_in + pip.height_in + gap
        ):
            moved = False
            new_top = pip.top_in - gap - box.height_in
            if new_top >= cy + margin:
                box = RegionBox(box.left_in, new_top, box.width_in, box.height_in)
                moved = True
            elif box.left_in + box.width_in > pip.left_in - gap:
                new_left = pip.left_in - gap - box.width_in
                if new_left >= cx + margin:
                    box = RegionBox(new_left, box.top_in, box.width_in, box.height_in)
                    moved = True
            if not moved:
                break
    if pip and anchor.startswith("top"):
        max_bottom = pip.top_in - pip_gap
        if box.top_in + box.height_in > max_bottom:
            box = RegionBox(box.left_in, box.top_in, box.width_in, max(0.55, max_bottom - box.top_in))
    if pip and anchor.startswith("bottom"):
        max_bottom = cy + ch - margin
        if box.top_in + box.height_in > max_bottom:
            box = RegionBox(box.left_in, max(cy + margin, max_bottom - box.height_in), box.width_in, box.height_in)
    return box


def _add_hero_headline(
    slide,
    box: RegionBox,
    headline: str,
    subheader: str,
    style: dict,
    theme: dict,
    *,
    text_style: str,
) -> None:
    if text_style == "overlay":
        _add_headline_content(slide, box, headline, subheader, style, theme)
    else:
        _add_text_panel(slide, box, headline, subheader, style, theme)


def _media_region_below_panel(
    cx: float,
    cy: float,
    cw: float,
    ch: float,
    panel: RegionBox,
    style: dict,
    kind: str,
) -> RegionBox:
    """Hero screenshot band under the top headline panel (``avatar_media_3``)."""
    gap = float(layout_in(style, kind, "panel_media_gap_in", 0.06))
    bottom_margin = float(layout_in(style, kind, "pip_margin_in", 0.32))
    media_top = panel.top_in + panel.height_in + gap
    media_h = (cy + ch) - media_top - bottom_margin
    return RegionBox(cx, media_top, cw, max(1.0, media_h))


def _name_card_pill_boxes(
    cx: float, cy: float, cw: float, ch: float, style: dict, kind: str, *, has_title: bool
) -> Tuple[RegionBox, Optional[RegionBox]]:
    """Bottom-left name and title pills (separate navy panels)."""
    margin = float(layout_in(style, kind, "panel_margin_in", 0.35))
    pw = cw * float(layout_in(style, kind, "panel_width_ratio", 0.42))
    name_h = float(layout_in(style, kind, "name_pill_height_in", 0.72))
    title_h = float(layout_in(style, kind, "title_pill_height_in", 0.55))
    gap = float(layout_in(style, kind, "pill_gap_in", 0.12))
    radius = float(layout_in(style, kind, "pill_radius_in", 0.08))
    bottom = cy + ch - margin
    if has_title:
        title_box = RegionBox(
            cx + margin, bottom - title_h, pw * 0.92, title_h, rounded=True, corner_radius_in=radius
        )
        name_box = RegionBox(
            cx + margin,
            title_box.top_in - gap - name_h,
            pw,
            name_h,
            rounded=True,
            corner_radius_in=radius,
        )
        return name_box, title_box
    return RegionBox(
        cx + margin, bottom - name_h, pw, name_h, rounded=True, corner_radius_in=radius
    ), None


def _slide_regions(
    prs, kind: str, style: dict, verse: Optional[dict] = None,
) -> Dict[str, Optional[RegionBox]]:
    cx, cy, cw, ch = _content_area(prs, style, kind)
    full = RegionBox(cx, cy, cw, ch)
    regions: Dict[str, Optional[RegionBox]] = {
        "media": None,
        "avatar": None,
        "text_panel": None,
        "name_pill": None,
        "title_pill": None,
    }

    if kind == "avatar_only":
        regions["avatar"] = full
    elif kind == "media_only":
        regions["media"] = full
    elif kind in _SPLIT_KINDS:
        ratio = float(layout_in(style, kind, "media_width_ratio"))
        gap_raw = layout_in(style, kind, "gap_in")
        if gap_raw is None:
            gap_raw = layout_in(style, kind, "inner_gap_in")
        gap = float(gap_raw or 0)
        rounded = kind.startswith("avatar_media_border")
        radius_raw = layout_in(style, kind, "inner_radius_in") if rounded else 0
        radius = float(radius_raw or 0.12) if rounded else 0.0
        media, avatar = _split_boxes(cx, cy, cw, ch, ratio, gap)
        regions["media"] = RegionBox(
            media.left_in, media.top_in, media.width_in, media.height_in, rounded, radius
        )
        regions["avatar"] = RegionBox(
            avatar.left_in, avatar.top_in, avatar.width_in, avatar.height_in, rounded, radius
        )
    elif kind in _PIP_KINDS:
        if kind == "avatar_media_3":
            pip = _pip_box(cx, cy, cw, ch, style, kind)
            regions["avatar"] = pip
            if _hero_layout_mode(style, verse, kind) == "full_bleed":
                regions["media"] = full
                cfg = _verse_text_panel_cfg(style, verse, kind)
                regions["text_panel"] = _text_panel_box_anchored(
                    prs, style, kind, cfg, pip=pip, verse=verse,
                )
            else:
                panel = _text_panel_box(prs, style, kind, "top")
                regions["text_panel"] = panel
                regions["media"] = _media_region_below_panel(cx, cy, cw, ch, panel, style, kind)
        else:
            regions["media"] = full if kind != "avatar_quote" else None
            regions["avatar"] = _pip_box(cx, cy, cw, ch, style, kind)
    elif kind == "avatar_name_card":
        regions["avatar"] = full
        name_box, title_box = _name_card_pill_boxes(
            cx, cy, cw, ch, style, kind, has_title=True
        )
        regions["name_pill"] = name_box
        regions["title_pill"] = title_box
    elif kind == "avatar_headline":
        pip = _pip_box(cx, cy, cw, ch, style, kind)
        regions["avatar"] = pip
        margin = float(layout_in(style, kind, "panel_margin_in", 0.6))
        regions["text_panel"] = _content_beside_pip(cx, cy, cw, ch, pip, margin)
    elif kind == "avatar_headline_full":
        regions["avatar"] = full
        regions["text_panel"] = _text_panel_box(prs, style, kind, "top")
    elif kind == "avatar_intro":
        pass
    elif kind == "avatar_outro":
        regions["avatar"] = full
    elif kind == "avatar_border":
        regions["avatar"] = full
    elif kind == "media_border":
        regions["media"] = full
    return regions


def export_slide_regions(
    prs, kind: str, style: dict, verse: Optional[dict] = None,
) -> Dict[str, Optional[RegionBox]]:
    """Public wrapper for layout region geometry (inches on slide)."""
    return _slide_regions(prs, kind, style, verse=verse)


def region_box_to_pixels(
    box: RegionBox,
    slide_width_in: float,
    slide_height_in: float,
    out_width: int,
    out_height: int,
) -> Dict[str, int]:
    """Map inch region box to output pixel coordinates."""
    pad_x, pad_y, cw, ch = letterbox_content_rect(
        slide_width_in, slide_height_in, out_width, out_height,
    )
    sx = cw / slide_width_in
    sy = ch / slide_height_in
    return {
        "x": pad_x + int(round(box.left_in * sx)),
        "y": pad_y + int(round(box.top_in * sy)),
        "width": max(1, int(round(box.width_in * sx))),
        "height": max(1, int(round(box.height_in * sy))),
    }


def letterbox_content_rect(
    slide_width_in: float,
    slide_height_in: float,
    out_width: int,
    out_height: int,
) -> Tuple[int, int, int, int]:
    """Return pad_x, pad_y, content_w, content_h for letterboxed slide in output frame."""
    slide_aspect = slide_width_in / slide_height_in
    out_aspect = out_width / out_height
    if slide_aspect > out_aspect:
        content_w = out_width
        content_h = max(1, int(round(out_width / slide_aspect)))
        pad_x, pad_y = 0, (out_height - content_h) // 2
    else:
        content_h = out_height
        content_w = max(1, int(round(out_height * slide_aspect)))
        pad_x, pad_y = (out_width - content_w) // 2, 0
    return pad_x, pad_y, content_w, content_h


def _poster_bytes(poster_path: Optional[str], source_file: Optional[str]) -> io.BytesIO:
    if poster_path:
        resolved = resolve_asset_path(poster_path, source_file=source_file)
        path = resolved if resolved else poster_path
        if path and os.path.isfile(path):
            with open(path, "rb") as fh:
                return io.BytesIO(fh.read())
    return io.BytesIO(_GREY_POSTER_PNG)


def _is_video_path(path: str) -> bool:
    return os.path.splitext(path)[1].lower() in _VIDEO_EXTS


def _video_mime(path: str) -> str:
    return _VIDEO_EXTS.get(os.path.splitext(path)[1].lower(), "video/unknown")


def _box_lengths(box: RegionBox):
    return Inches(box.left_in), Inches(box.top_in), Inches(box.width_in), Inches(box.height_in)


def _draw_filled_rect(slide, box: RegionBox, rgb: RGBColor, rounded: bool = False) -> None:
    left, top, width, height = _box_lengths(box)
    if rounded:
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        try:
            short_in = min(box.width_in, box.height_in)
            adj = min(0.5, max(0.02, box.corner_radius_in / short_in)) if short_in > 0 else 0.08
            shape.adjustments[0] = adj
        except (IndexError, AttributeError):
            pass
    else:
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb
    shape.line.fill.background()


def _place_empty_region(slide, box: RegionBox, zone: str) -> None:
    colour = _MEDIA_WHITE if zone == "media" else _AVATAR_GREY
    _draw_filled_rect(slide, box, colour, rounded=box.rounded)


def _jpeg_pip_preview_enabled(style: Optional[dict], verse: dict) -> bool:
    if verse.get("jpeg_show_pip_preview") is False:
        return False
    if verse.get("jpeg_show_pip_preview") is True:
        return True
    return bool((style or {}).get("_jpeg_show_pip_preview"))


def _place_pip_preview_placeholder(slide, box: RegionBox, style: Optional[dict]) -> None:
    """Grey circle for JPEG layout QA; live HeyGen PiP is composited in MP4 export."""
    from pptx.enum.shapes import MSO_SHAPE

    left, top, width, height = _box_lengths(box)
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = _PIP_BACKDROP_GREY
    shape.line.fill.background()


def _place_overlay_only_pip(
    slide, box: RegionBox, style: Optional[dict], verse: dict,
) -> None:
    """PiP drawn in FFmpeg only — optional grey placeholder for JPEG QA (no white ring)."""
    if _jpeg_pip_preview_enabled(style, verse):
        _place_pip_preview_placeholder(slide, box, style)


def _fit_movie_in_box(slide, movie_path: str, poster: io.BytesIO, box: RegionBox, mime: str) -> None:
    left, top, width, height = _box_lengths(box)
    slide.shapes.add_movie(
        movie_path,
        left,
        top,
        width,
        height,
        poster_frame_image=poster,
        mime_type=mime,
    )


def _pip_still_pixels(box: RegionBox, dpi: int = 150) -> tuple[int, int]:
    px = max(96, int(round(box.width_in * dpi)))
    return px, max(96, int(round(box.height_in * dpi)))


# Extra vertical crop for circular PiP — lowers empty space above the head (lower crop_y).
_CIRCLE_HEADSPACE_TRIM = 0.012


def _framing_shape(style: dict, layout_kind: str) -> str:
    kind_shape = layout_in(style, layout_kind, "shape", None)
    if kind_shape is not None:
        return str(kind_shape).lower()
    return _pip_shape_kind(style)


def avatar_framing(style: dict, layout_kind: str = "pip") -> Tuple[float, float, float]:
    """Per-layout avatar crop/zoom; falls back to ``layouts.pip`` defaults."""
    crop_x_kind = layout_in(style, layout_kind, "crop_x_ratio", None)
    crop_y_kind = layout_in(style, layout_kind, "crop_y_ratio", None)
    zoom_kind = layout_in(style, layout_kind, "zoom_ratio", None)
    crop_x = float(
        crop_x_kind if crop_x_kind is not None else layout_in(style, "pip", "crop_x_ratio", 0.5)
    )
    crop_y = float(
        crop_y_kind if crop_y_kind is not None else layout_in(style, "pip", "crop_y_ratio", 0.02)
    )
    zoom = float(
        zoom_kind if zoom_kind is not None else layout_in(style, "pip", "zoom_ratio", 1.47)
    )
    shape = _framing_shape(style, layout_kind)
    if shape in ("circle", "round", "rounded"):
        crop_y = max(0.0, crop_y - _CIRCLE_HEADSPACE_TRIM)
    return crop_x, crop_y, zoom


def _pip_video_seek_sec(verse: Optional[dict], *, default: float = 0.5) -> float:
    """Match video overlay timing: sample near ``audio_start_sec`` when set."""
    if verse and verse.get("audio_start_sec") is not None:
        return max(0.0, float(verse["audio_start_sec"]) + 0.35)
    return default


def _save_pip_still_png(
    video_path: str,
    box: RegionBox,
    style: dict,
    *,
    source_file: Optional[str] = None,
    poster: Optional[io.BytesIO] = None,
    rect_mask: bool = False,
    layout_kind: str = "pip",
    seek_sec: Optional[float] = None,
    verse: Optional[dict] = None,
) -> Optional[str]:
    """Extract a face-centred still PNG; circle mask only for PiP (not full-bleed panels)."""
    w, h = _pip_still_pixels(box)
    crop_x, crop_y, zoom = avatar_framing(style, layout_kind)
    zoom = max(1.0, float(zoom))
    shape = _pip_shape_kind(style)
    is_circle = not rect_mask and shape in ("circle", "round", "rounded")
    tmp = tempfile.NamedTemporaryFile(suffix=".png", delete=False)
    tmp.close()
    out_path = tmp.name
    try:
        from .ffmpeg_composer import _circle_alpha_filter, _cover_scale_filter

        vf = _cover_scale_filter(
            w, h, crop_x_ratio=crop_x, crop_y_ratio=crop_y, zoom_ratio=zoom,
        )
        if is_circle:
            vf = f"{vf},{_circle_alpha_filter()}"
        ss = _pip_video_seek_sec(verse) if seek_sec is None else max(0.0, float(seek_sec))
        cmd = [
            "ffmpeg", "-y", "-hide_banner", "-loglevel", "error",
            "-ss", f"{ss:.3f}", "-i", video_path,
            "-vframes", "1", "-vf", vf, out_path,
        ]
        proc = subprocess.run(cmd, capture_output=True, text=True, timeout=60)
        if proc.returncode == 0 and os.path.isfile(out_path):
            return out_path
    except Exception:
        pass
    if poster and poster.getbuffer().nbytes > 64:
        try:
            from PIL import Image, ImageDraw

            poster.seek(0)
            img = Image.open(poster).convert("RGBA")
            img = img.resize((w, h), Image.Resampling.LANCZOS)
            if is_circle:
                mask = Image.new("L", (w, h), 0)
                ImageDraw.Draw(mask).ellipse((0, 0, w - 1, h - 1), fill=255)
                img.putalpha(mask)
            img.save(out_path, "PNG")
            return out_path
        except Exception:
            pass
    try:
        os.unlink(out_path)
    except OSError:
        pass
    return None


def _place_avatar_still_in_box(
    slide,
    video_path: str,
    poster: io.BytesIO,
    box: RegionBox,
    style: dict,
    *,
    source_file: Optional[str] = None,
    rect_mask: bool = False,
    layout_kind: str = "pip",
    verse: Optional[dict] = None,
) -> bool:
    still = _save_pip_still_png(
        video_path, box, style, source_file=source_file, poster=poster,
        rect_mask=rect_mask, layout_kind=layout_kind, verse=verse,
    )
    if still:
        _place_picture_in_box(slide, still, box, fit="cover" if rect_mask else "fill")
        try:
            os.unlink(still)
        except OSError:
            pass
        return True
    return False


def _place_picture_in_box(slide, path: str, box: RegionBox, fit: str) -> None:
    from .core import _fit_picture_in_box

    left, top, width, height = _box_lengths(box)
    fit = (fit or "contain").lower()
    if fit == "fill":
        slide.shapes.add_picture(path, left, top, width=width, height=height)
        return
    pic = slide.shapes.add_picture(path, left, top, width=width)
    _fit_picture_in_box(pic, left, top, width, height, fit if fit in ("contain", "cover") else "contain")


def _place_media_in_box(
    slide,
    box: RegionBox,
    media_path: Optional[str],
    *,
    fit: str = "contain",
    poster_path: Optional[str] = None,
    source_file: Optional[str] = None,
) -> None:
    if not media_path:
        _place_empty_region(slide, box, "media")
        return
    resolved = resolve_asset_path(media_path, source_file=source_file)
    path = resolved if resolved else media_path
    if not path or not os.path.isfile(path):
        print(f"Warning: Media not found: {media_path}")
        _place_empty_region(slide, box, "media")
        return
    if _is_video_path(path):
        poster = _poster_bytes(poster_path, source_file)
        _fit_movie_in_box(slide, path, poster, box, _video_mime(path))
    else:
        _place_picture_in_box(slide, path, box, fit)


def _pip_backdrop_box(box: RegionBox, pad_in: float = 0.03) -> RegionBox:
    return RegionBox(
        box.left_in - pad_in,
        box.top_in - pad_in,
        box.width_in + 2 * pad_in,
        box.height_in + 2 * pad_in,
        rounded=box.rounded,
        corner_radius_in=box.corner_radius_in,
    )


def _draw_pip_backdrop(slide, box: RegionBox, style: dict) -> None:
    """Grey fill behind circle PiP avatar."""
    backdrop = _pip_backdrop_box(box)
    left, top, width, height = _box_lengths(backdrop)
    shape_kind = str(layout_in(style, "pip", "shape", "circle")).lower()
    if shape_kind == "circle":
        shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, width, height)
    else:
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        try:
            shape.adjustments[0] = 0.18
        except (IndexError, AttributeError):
            pass
    shape.fill.solid()
    shape.fill.fore_color.rgb = _PIP_BACKDROP_GREY
    shape.line.fill.background()


def _draw_pip_frame(slide, box: RegionBox, style: dict) -> None:
    """Ring around the avatar PiP (circle, rounded rect, or square)."""
    left, top, width, height = _box_lengths(box)
    shape_kind = _pip_shape_kind(style)
    if shape_kind in ("circle", "round", "rounded"):
        shape = slide.shapes.add_shape(
            MSO_SHAPE.OVAL if shape_kind == "circle" else MSO_SHAPE.ROUNDED_RECTANGLE,
            left, top, width, height,
        )
        if shape_kind != "circle":
            try:
                shape.adjustments[0] = 0.18
            except (IndexError, AttributeError):
                pass
    else:
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.background()
    border = str(layout_in(style, "pip", "border_color", "#FFFFFF"))
    shape.line.color.rgb = _hex_rgb(border)
    shape.line.width = Pt(float(layout_in(style, "pip", "border_width_pt", 2.5)))


def _place_avatar_in_box(
    slide,
    box: RegionBox,
    avatar_path: Optional[str],
    *,
    poster_path: Optional[str] = None,
    source_file: Optional[str] = None,
    style: Optional[dict] = None,
    draw_frame: Optional[bool] = None,
    layout_kind: str = "pip",
    verse: Optional[dict] = None,
    panel_fill_rgb: Optional[RGBColor] = None,
) -> None:
    if not avatar_path:
        _place_empty_region(slide, box, "avatar")
        return
    resolved = resolve_asset_path(avatar_path, source_file=source_file)
    path = resolved if resolved else avatar_path
    if not path or not os.path.isfile(path):
        print(f"Warning: Avatar video not found: {avatar_path}")
        _place_empty_region(slide, box, "avatar")
        return
    poster = _poster_bytes(poster_path, source_file)
    pip_style = style or {}
    shape_kind = resolve_avatar_shape(
        pip_style, layout_kind=layout_kind, box=box, verse=verse,
    )
    is_circle = shape_uses_circle_mask(shape_kind)
    is_square = shape_kind in ("square", "rect", "rectangle", "h_rect", "v_rect", "horizontal", "vertical", "wide", "tall")

    if draw_frame is None:
        draw_frame = is_circle

    if panel_fill_rgb is not None and not is_circle:
        _draw_filled_rect(slide, box, panel_fill_rgb)

    if is_circle:
        if not _place_avatar_still_in_box(
            slide, path, poster, box, pip_style, source_file=source_file,
            layout_kind=layout_kind, verse=verse,
        ):
            _fit_movie_in_box(slide, path, poster, box, _video_mime(path))
        if draw_frame:
            _draw_pip_frame(slide, box, pip_style)
        return

    if not _place_avatar_still_in_box(
        slide, path, poster, box, pip_style, source_file=source_file, rect_mask=True,
        layout_kind=layout_kind, verse=verse,
    ):
        _fit_movie_in_box(slide, path, poster, box, _video_mime(path))
    if draw_frame and is_square:
        _draw_pip_frame(slide, box, pip_style)


def place_floating_avatar_pip(
    slide,
    verse: dict,
    style: dict,
    *,
    prs,
    source_file: Optional[str] = None,
) -> None:
    """Avatar PiP on standard slides — drawn last so it floats above the slide."""
    path = verse.get("avatar_video_path")
    if not path:
        return
    style = dict(style or {})
    if source_file:
        style["_source_file"] = source_file
    box = export_floating_pip_box(prs, style)
    _place_overlay_only_pip(slide, box, style, verse)


def _draw_border_frame(slide, prs, style: dict, kind: str) -> None:
    inset = _border_inset(style, kind)
    colour = _hex_rgb(str(layout_in(style, kind, "border_color", "#1E3A5F")))
    width_pt = float(layout_in(style, kind, "border_width_pt", 8))
    sw, sh = prs.slide_width.inches, prs.slide_height.inches
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(inset),
        Inches(inset),
        Inches(sw - 2 * inset),
        Inches(sh - 2 * inset),
    )
    shape.fill.background()
    shape.line.color.rgb = colour
    shape.line.width = Pt(width_pt)


def _panel_colour(style: dict) -> RGBColor:
    raw = (style or {}).get("avatar_panel_color")
    return _hex_rgb(raw) if raw else _PANEL_NAVY


def _format_subheader_lines(subheader: str) -> list[str]:
    parts = [s.strip() for s in subheader.replace(" · ", "·").split("·") if s.strip()]
    if len(parts) <= 1:
        return [subheader.strip()] if subheader.strip() else []
    return [f"•  {part}" for part in parts]


def _headline_font_sizes(style: dict, headline: str, subheader: str) -> tuple[int, int]:
    sub_len = len(subheader or "")
    head_pt = int(typography_pt(style, "title_size_pt", 44) * 0.78)
    if len(headline) > 48:
        head_pt = max(head_pt - 6, 28)
    if sub_len > 140:
        sub_pt = 18
    elif sub_len > 90:
        sub_pt = 20
    elif sub_len > 55:
        sub_pt = 22
    else:
        sub_pt = int(typography_pt(style, "body_size_pt", 32) * 0.72)
    return head_pt, sub_pt


def _add_headline_content(
    slide, box: RegionBox, headline: str, subheader: str, style: dict, theme: dict
) -> None:
    from .core import _write_body_paragraph

    left, top, width, height = _box_lengths(box)
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    tf.margin_left = tf.margin_right = Pt(10)
    tf.margin_top = tf.margin_bottom = Pt(6)
    tf.vertical_anchor = MSO_ANCHOR.TOP

    head_pt, sub_pt = _headline_font_sizes(style, headline, subheader)
    head_theme = dict(theme)
    head_theme["body"] = theme.get("title") or theme.get("body")
    sub_theme = dict(theme)
    sub_theme["body"] = theme.get("subtitle") or theme.get("reference") or theme.get("body")

    p = tf.paragraphs[0]
    _write_body_paragraph(p, headline, head_pt, head_theme, style=style, alignment=PP_ALIGN.LEFT)
    p.space_after = Pt(12)

    for line in _format_subheader_lines(subheader):
        para = tf.add_paragraph()
        _write_body_paragraph(para, line, sub_pt, sub_theme, style=style, alignment=PP_ALIGN.LEFT)
        para.space_after = Pt(8)


def _add_single_pill_text(
    slide, box: RegionBox, text: str, font_pt: int, style: dict, theme: dict
) -> None:
    from .core import _write_body_paragraph

    _draw_filled_rect(slide, box, _panel_colour(style), rounded=box.rounded)
    left, top, width, height = _box_lengths(box)
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Pt(14)
    tf.margin_top = tf.margin_bottom = Pt(6)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    panel_theme = dict(theme)
    panel_theme["body"] = RGBColor(0xFF, 0xFF, 0xFF)
    _write_body_paragraph(
        tf.paragraphs[0], text, font_pt, panel_theme, style=style, alignment=PP_ALIGN.LEFT
    )


def _add_name_card_pills(
    slide,
    prs,
    verse: dict,
    style: dict,
    theme: dict,
    *,
    cx: float,
    cy: float,
    cw: float,
    ch: float,
) -> None:
    headline = str(verse.get("headline", ""))
    subheader = str(verse.get("subheader") or "").strip()
    name_box, title_box = _name_card_pill_boxes(
        cx, cy, cw, ch, style, "avatar_name_card", has_title=bool(subheader)
    )
    name_pt = int(typography_pt(style, "title_size_pt", 44) * 0.68)
    title_pt = int(typography_pt(style, "subtitle_size_pt", 28) * 0.82)
    _add_single_pill_text(slide, name_box, headline, name_pt, style, theme)
    if subheader and title_box:
        _add_single_pill_text(slide, title_box, subheader, title_pt, style, theme)


def _draw_rotated_block(
    slide, left_in: float, top_in: float, w_in: float, h_in: float, rgb: RGBColor, rotation_deg: float
) -> None:
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(left_in), Inches(top_in), Inches(w_in), Inches(h_in)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb
    shape.line.fill.background()
    shape.rotation = rotation_deg


def _render_avatar_intro(slide, prs, style: dict) -> None:
    """Decorative diagonal colour blocks (HeyGen-style intro)."""
    sw, sh = prs.slide_width.inches, prs.slide_height.inches
    _draw_rotated_block(slide, -0.4, -0.2, sw + 0.8, sh * 0.55, _INTRO_COLOURS["cream"], -8)
    _draw_rotated_block(slide, sw * 0.35, -0.5, sw * 0.85, sh * 0.75, _INTRO_COLOURS["navy"], 22)
    _draw_rotated_block(slide, -0.3, sh * 0.35, sw * 0.7, sh * 0.55, _INTRO_COLOURS["tan"], -18)
    _draw_rotated_block(slide, sw * 0.05, sh * 0.62, sw * 0.45, sh * 0.38, _INTRO_COLOURS["sage"], 12)


def _draw_centre_diamond(slide, prs, style: dict) -> None:
    sw, sh = prs.slide_width.inches, prs.slide_height.inches
    size = float(layout_in(style, "avatar_outro", "diamond_size_in", 1.85))
    left = (sw - size) / 2
    top = (sh - size) / 2
    shape = slide.shapes.add_shape(MSO_SHAPE.DIAMOND, Inches(left), Inches(top), Inches(size), Inches(size))
    shape.fill.solid()
    shape.fill.fore_color.rgb = _MEDIA_WHITE
    shape.line.fill.background()


def _add_text_panel(slide, box: RegionBox, headline: str, subheader: str, style: dict, theme: dict) -> None:
    from pptx.enum.shapes import MSO_SHAPE

    from .core import _write_body_paragraph

    left, top, width, height = _box_lengths(box)
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height,
    )
    if box.rounded:
        try:
            short_in = min(box.width_in, box.height_in)
            adj = min(0.5, max(0.02, box.corner_radius_in / short_in)) if short_in > 0 else 0.08
            shape.adjustments[0] = adj
        except (IndexError, AttributeError):
            pass
    shape.fill.solid()
    shape.fill.fore_color.rgb = _panel_colour(style)
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Pt(14)
    tf.margin_top = Pt(12)
    tf.margin_bottom = Pt(10)
    tf.vertical_anchor = MSO_ANCHOR.TOP
    head_pt = int(typography_pt(style, "title_size_pt", 44) * 0.72)
    sub_pt = int(typography_pt(style, "subtitle_size_pt", 28) * 0.78)
    if box.width_in < 3.8:
        head_pt = max(22, head_pt - 2)
        sub_pt = max(16, sub_pt - 2)
    panel_theme = dict(theme)
    panel_theme["body"] = RGBColor(0xFF, 0xFF, 0xFF)
    p = tf.paragraphs[0]
    p.space_after = Pt(6)
    _write_body_paragraph(p, headline, head_pt, panel_theme, style=style, alignment=PP_ALIGN.LEFT)
    if subheader:
        p2 = tf.add_paragraph()
        p2.space_before = Pt(2)
        _write_body_paragraph(p2, subheader, sub_pt, panel_theme, style=style, alignment=PP_ALIGN.LEFT)


def _render_avatar_quote(slide, prs, verse: dict, style: dict, regions: dict, theme: dict) -> None:
    from .core import _resolve_alignment, _write_body_paragraph

    bg = str(layout_in(style, "avatar_quote", "quote_bg_color", "#1E3A5F"))
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = _hex_rgb(bg)

    pip = regions.get("avatar")
    cx, cy, cw, ch = _content_area(prs, style, "avatar_quote")
    margin = float(layout_in(style, "avatar_quote", "margin_in", 0.6))
    pip_w = pip.width_in if pip else 0.0
    pip_top = pip.top_in if pip else ch
    text_w_in = max(cw - pip_w - 2 * margin, 4.0)
    left_in = cx + margin

    quote_pt = int(verse.get("font_size") or typography_pt(style, "quote_size_pt", 36))
    top_in = float(layout_in(style, "avatar_quote", "top_in", 1.6))
    available_h = max(pip_top - top_in - 0.35, 2.0)
    align = _resolve_alignment(verse.get("alignment", "center"))

    quote_h = available_h * 0.62
    tb = slide.shapes.add_textbox(Inches(left_in), Inches(top_in), Inches(text_w_in), Inches(quote_h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Pt(12)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    quote_theme = dict(theme)
    quote_theme["body"] = RGBColor(0xFF, 0xFF, 0xFF)
    quote_theme["reference"] = RGBColor(0xDD, 0xDD, 0xDD)
    _write_body_paragraph(tf.paragraphs[0], verse["text"], quote_pt, quote_theme, style=style, alignment=align)

    ref_str = (verse.get("reference") or "").strip()
    if ref_str:
        ref_pt = int(typography_pt(style, "reference_size_pt", 24))
        ref_top = top_in + quote_h + 0.15
        ref_h = max(min(available_h - quote_h - 0.2, 1.1), 0.55)
        rt = slide.shapes.add_textbox(Inches(left_in), Inches(ref_top), Inches(text_w_in), Inches(ref_h))
        tf_ref = rt.text_frame
        tf_ref.word_wrap = True
        tf_ref.vertical_anchor = MSO_ANCHOR.TOP
        rp = tf_ref.paragraphs[0]
        rp.text = ref_str
        rp.alignment = align
        rp.font.size = Pt(ref_pt)
        rp.font.italic = True
        rp.font.color.rgb = quote_theme["reference"]
        if theme.get("font_name"):
            rp.font.name = theme["font_name"]


def render_avatar_slide(prs, kind: str, verse: dict, style=None, *, source_file: Optional[str] = None):
    """Build one avatar layout slide."""
    from .core import _apply_slide_background, _resolve_theme

    style = dict(style or {})
    if source_file:
        style["_source_file"] = source_file
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    if kind not in ("avatar_quote", "avatar_intro"):
        _apply_slide_background(slide, style, prs)
    theme = _resolve_theme(style)
    regions = _slide_regions(prs, kind, style, verse=verse)

    media_box = regions.get("media")
    avatar_box = regions.get("avatar")
    media_fit = verse.get("media_fit", "contain")

    if kind == "avatar_intro":
        _render_avatar_intro(slide, prs, style)
        _stamp_slide_type_note(slide, kind, verse)
        return slide

    if kind == "avatar_quote":
        _render_avatar_quote(slide, prs, verse, style, regions, theme)
        if avatar_box:
            _place_overlay_only_pip(slide, avatar_box, style, verse)
        _stamp_slide_type_note(slide, kind, verse)
        return slide

    if kind == "avatar_headline":
        panel = regions.get("text_panel")
        if panel:
            _add_headline_content(
                slide,
                panel,
                str(verse.get("headline", "")),
                str(verse.get("subheader") or ""),
                style,
                theme,
            )
        if avatar_box:
            _place_avatar_in_box(
                slide,
                avatar_box,
                verse.get("avatar_video_path"),
                poster_path=verse.get("avatar_poster_path"),
                source_file=source_file,
                style=style,
            )
        _stamp_slide_type_note(slide, kind, verse)
        return slide

    if kind == "avatar_headline_full":
        if avatar_box:
            _place_avatar_in_box(
                slide,
                avatar_box,
                verse.get("avatar_video_path"),
                poster_path=verse.get("avatar_poster_path"),
                source_file=source_file,
                style=style,
                draw_frame=False,
            )
        panel = regions.get("text_panel")
        if panel:
            _add_text_panel(
                slide,
                panel,
                str(verse.get("headline", "")),
                str(verse.get("subheader") or ""),
                style,
                theme,
            )
        _stamp_slide_type_note(slide, kind, verse)
        return slide

    if media_box and kind in _SPLIT_KINDS and media_box.rounded:
        _place_empty_region(slide, media_box, "media")
    if avatar_box and kind in _SPLIT_KINDS and avatar_box.rounded:
        _place_empty_region(slide, avatar_box, "avatar")

    if media_box:
        _place_media_in_box(
            slide,
            media_box,
            verse.get("media_path"),
            fit=media_fit,
            poster_path=verse.get("media_poster_path"),
            source_file=source_file,
        )
    if avatar_box:
        if kind in _AVATAR_PIP_VIDEO_OVERLAY_ONLY:
            _place_overlay_only_pip(slide, avatar_box, style, verse)
        else:
            _place_avatar_in_box(
                slide,
                avatar_box,
                verse.get("avatar_video_path"),
                poster_path=verse.get("avatar_poster_path"),
                source_file=source_file,
                style=style,
            )

    if kind == "avatar_name_card":
        cx, cy, cw, ch = _content_area(prs, style, kind)
        _add_name_card_pills(
            slide, prs, verse, style, theme, cx=cx, cy=cy, cw=cw, ch=ch
        )
    else:
        panel = regions.get("text_panel")
        if panel:
            tstyle = _text_style_mode(style, verse, kind) if kind == "avatar_media_3" else "navy_panel"
            _add_hero_headline(
                slide,
                panel,
                str(verse.get("headline", "")),
                str(verse.get("subheader") or ""),
                style,
                theme,
                text_style=tstyle,
            )

    if kind == "avatar_outro":
        _draw_centre_diamond(slide, prs, style)

    if kind in _BORDER_KINDS:
        _draw_border_frame(slide, prs, style, kind)
    _stamp_slide_type_note(slide, kind, verse)
    return slide


def _stamp_slide_type_note(slide, kind: str, verse: dict) -> None:
    """Store layout kind in speaker notes for list-slides (hidden from normal notes)."""
    try:
        notes = (verse.get("notes") or "").strip()
        meta = f"slide_type: {kind}"
        frame = slide.notes_slide.notes_text_frame
        if notes:
            if meta not in notes:
                frame.text = f"{notes}\n{meta}"
        else:
            frame.text = meta
    except Exception:
        pass
