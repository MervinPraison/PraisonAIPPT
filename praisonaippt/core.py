"""
Core presentation creation logic for Bible verses PowerPoint generator.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from .utils import split_long_text, sanitize_filename, resolve_asset_path
from .pdf_converter import PDFOptions, convert_pptx_to_pdf
from .layout_tokens import (
    layout_in,
    typography_pt,
    content_box,
    content_width_inches,
    split_max_length_default,
    title_custom_threshold,
)


def _apply_slide_background(slide, style: dict, prs=None):
    """
    Apply background to a slide from a slide_style dict.

    Supported keys:
        background_color (str): Hex color e.g. '#1A1A2E'
        background_image (str): Absolute or relative path to an image file
    """
    if not style:
        return

    img_path = style.get('background_image')
    bg_color = style.get('background_color')
    source_file = style.get('_source_file')

    if img_path:
        import os
        resolved = resolve_asset_path(img_path, source_file=source_file)
        img_path = resolved if resolved else img_path
        if os.path.exists(img_path):
            # Use prs dimensions or fall back to standard 16:9 (13.33 x 7.5 in)
            if prs is not None:
                w, h = prs.slide_width, prs.slide_height
            else:
                w, h = Inches(13.33), Inches(7.5)
            pic = slide.shapes.add_picture(img_path, 0, 0, w, h)
            # Move picture to back
            sp_tree = slide.shapes._spTree
            sp_tree.remove(pic._element)
            sp_tree.insert(2, pic._element)
        elif bg_color:
            img_path = None
    if not img_path and bg_color:
        fill = slide.background.fill
        fill.solid()
        hex_c = bg_color.lstrip('#')
        r, g, b = int(hex_c[0:2], 16), int(hex_c[2:4], 16), int(hex_c[4:6], 16)
        fill.fore_color.rgb = RGBColor(r, g, b)


def _apply_speaker_notes(slide, notes):
    """Set presenter notes on a slide when ``notes`` is non-empty."""
    text = (notes or "").strip()
    if not text:
        return
    try:
        slide.notes_slide.notes_text_frame.text = text
    except AttributeError:
        pass


def _write_body_paragraph(p, text, font_size, theme, style=None, alignment=None,
                          highlights=None, large_text=None):
    """Fill a paragraph with plain or highlighted body text."""
    if alignment is not None and hasattr(alignment, "value"):
        align = alignment
    else:
        align = _resolve_alignment(alignment) if alignment else PP_ALIGN.LEFT
    p.alignment = align
    fn = theme.get("font_name")
    ann_pt = int(typography_pt(style or {}, "annotation_size_pt", 46))
    if (highlights and len(highlights) > 0) or (large_text and len(large_text) > 0):
        _apply_highlights(
            p, text or "", highlights, large_text,
            body_rgb=theme["body"],
            highlight_rgb=theme["highlight"],
            annotation_rgb=theme["annotation"],
            font_name=fn,
            base_font_size=int(font_size),
            annotation_size_pt=ann_pt,
        )
    else:
        run = p.add_run()
        run.text = text or ""
        run.font.size = Pt(font_size)
        run.font.color.rgb = theme["body"]
        if fn:
            run.font.name = fn


def _two_column_layout(prs, style, kind="two_column"):
    """Return geometry for side-by-side columns."""
    left, width, width_in, _ = content_box(prs, style, kind)
    gap_in = float(layout_in(style, kind, "column_gap_in", 0.4))
    col_w_in = (width_in - gap_in) / 2.0
    left_x = left.inches
    right_x = left.inches + col_w_in + gap_in
    top_in = float(layout_in(style, kind, "top_in", 0.9))
    bottom_in = float(layout_in(style, kind, "bottom_reserve_in", 0.5))
    height_in = prs.slide_height.inches - top_in - bottom_in
    return left_x, right_x, col_w_in, top_in, height_in, gap_in


def add_title_only_slide(prs, title, subtitle=None, style=None, font_size=None):
    """Title-only slide (PowerPoint Title Only / Google TITLE_ONLY)."""
    style = style or {}
    theme = _resolve_theme(style)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _apply_slide_background(slide, style, prs)
    title_pt = int(font_size or typography_pt(style, "title_size_pt", 44))
    left, width, width_in, _ = content_box(prs, style, "title_only")
    title_lines = _estimate_text_lines(title or "", width_in, title_pt)
    title_h_in = min(0.55 + title_lines * 0.52, 3.2)
    sub = (subtitle or "").strip()
    sub_pt = int(typography_pt(style, "subtitle_size_pt", 28))
    sub_h_in = 0.0
    if sub:
        sub_h_in = min(0.45 + _estimate_text_lines(sub, width_in, sub_pt) * 0.38, 2.5)
    gap_in = float(layout_in(style, "title", "subtitle_gap_in", 0.25)) if sub else 0.0
    block_h_in = title_h_in + (gap_in + sub_h_in if sub else 0)
    top_in = (prs.slide_height.inches - block_h_in) / 2.0
    tb = slide.shapes.add_textbox(left, Inches(top_in), width, Inches(title_h_in))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title or ""
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(title_pt)
    p.font.bold = True
    p.font.color.rgb = theme["title"]
    if theme["font_name"]:
        p.font.name = theme["font_name"]
    if sub:
        tb2 = slide.shapes.add_textbox(
            left, Inches(top_in + title_h_in + gap_in), width, Inches(sub_h_in),
        )
        tf2 = tb2.text_frame
        tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        p2.text = sub
        p2.alignment = PP_ALIGN.CENTER
        p2.font.size = Pt(sub_pt)
        p2.font.color.rgb = theme["subtitle"]
        if theme["font_name"]:
            p2.font.name = theme["font_name"]
    return slide


def add_two_column_slide(prs, left_text, right_text, style=None, font_size=28,
                         alignment="left", left_highlights=None, right_highlights=None):
    """Two Content layout — side-by-side body columns."""
    style = style or {}
    theme = _resolve_theme(style)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _apply_slide_background(slide, style, prs)
    lx, rx, col_w, top_in, height_in, _ = _two_column_layout(prs, style, "two_column")
    for x_in, text, hl in (
        (lx, left_text, left_highlights),
        (rx, right_text, right_highlights),
    ):
        tb = slide.shapes.add_textbox(Inches(x_in), Inches(top_in), Inches(col_w), Inches(height_in))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.TOP
        p = tf.paragraphs[0]
        _write_body_paragraph(p, text, font_size, theme, style=style, alignment=alignment, highlights=hl)
    return slide


def add_comparison_slide(prs, columns, style=None, font_size=28, alignment="left", reference=None):
    """Comparison layout — heading + body per column."""
    style = style or {}
    theme = _resolve_theme(style)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _apply_slide_background(slide, style, prs)
    lx, rx, col_w, _, height_in, _ = _two_column_layout(prs, style, "comparison")
    top_in = float(layout_in(style, "comparison", "top_in", 0.75))
    heading_h_in = float(layout_in(style, "comparison", "heading_height_in", 0.55))
    body_gap_in = float(layout_in(style, "comparison", "body_top_gap_in", 0.12))
    heading_pt = int(typography_pt(style, "comparison_heading_size_pt", 28))
    align = _resolve_alignment(alignment)
    cols = (columns or [])[:2]
    while len(cols) < 2:
        cols.append({})
    for x_in, col in zip((lx, rx), cols):
        heading = (col.get("heading") or "").strip()
        body = (col.get("text") or "").strip()
        hl = col.get("highlights")
        if heading:
            ht = slide.shapes.add_textbox(
                Inches(x_in), Inches(top_in), Inches(col_w), Inches(heading_h_in),
            )
            hp = ht.text_frame.paragraphs[0]
            hp.text = heading
            hp.alignment = align
            hp.font.size = Pt(heading_pt)
            hp.font.bold = True
            hp.font.color.rgb = theme["body"]
            if theme["font_name"]:
                hp.font.name = theme["font_name"]
        body_top = top_in + (heading_h_in + body_gap_in if heading else 0)
        body_h = height_in - (heading_h_in + body_gap_in if heading else 0)
        bt = slide.shapes.add_textbox(
            Inches(x_in), Inches(body_top), Inches(col_w), Inches(max(body_h, 1.0)),
        )
        tf = bt.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.TOP
        _write_body_paragraph(
            tf.paragraphs[0], body, font_size, theme, style=style,
            alignment=align, highlights=hl,
        )
    if reference and str(reference).strip():
        ref_pt = int(typography_pt(style, "reference_size_bottom_pt", 22))
        ref_h_in = 0.7
        ref_y = prs.slide_height.inches - ref_h_in - 0.35
        left, width, _, _ = content_box(prs, style, "comparison")
        rt = slide.shapes.add_textbox(left, Inches(ref_y), width, Inches(ref_h_in))
        rp = rt.text_frame.paragraphs[0]
        rp.text = str(reference).strip()
        rp.alignment = PP_ALIGN.CENTER
        rp.font.size = Pt(ref_pt)
        rp.font.italic = True
        rp.font.color.rgb = theme["reference"]
        if theme["font_name"]:
            rp.font.name = theme["font_name"]
    return slide


def add_big_number_slide(prs, number, label, style=None, reference=None):
    """Big Number layout (Google Slides BIG_NUMBER)."""
    style = style or {}
    theme = _resolve_theme(style)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _apply_slide_background(slide, style, prs)
    num_pt = int(typography_pt(style, "big_number_size_pt", 120))
    label_pt = int(typography_pt(style, "big_number_label_size_pt", 32))
    left, width, width_in, _ = content_box(prs, style, "big_number")
    num_str = str(number or "").strip()
    label_str = (label or "").strip()
    ref_str = (reference or "").strip()
    num_lines = _estimate_text_lines(num_str, width_in, num_pt)
    num_h_in = min(1.2 + num_lines * 0.85, 3.5)
    label_h_in = 0.0
    if label_str:
        label_h_in = min(0.5 + _estimate_text_lines(label_str, width_in, label_pt) * 0.42, 2.0)
    ref_h_in = 0.65 if ref_str else 0.0
    block_h = num_h_in + (0.2 + label_h_in if label_str else 0) + (0.25 + ref_h_in if ref_str else 0)
    top_in = (prs.slide_height.inches - block_h) / 2.0
    nt = slide.shapes.add_textbox(left, Inches(top_in), width, Inches(num_h_in))
    np = nt.text_frame.paragraphs[0]
    np.text = num_str
    np.alignment = PP_ALIGN.CENTER
    np.font.size = Pt(num_pt)
    np.font.bold = True
    np.font.color.rgb = theme["title"]
    if theme["font_name"]:
        np.font.name = theme["font_name"]
    y = top_in + num_h_in + 0.2
    if label_str:
        lt = slide.shapes.add_textbox(left, Inches(y), width, Inches(label_h_in))
        lp = lt.text_frame.paragraphs[0]
        lp.text = label_str
        lp.alignment = PP_ALIGN.CENTER
        lp.font.size = Pt(label_pt)
        lp.font.color.rgb = theme["body"]
        if theme["font_name"]:
            lp.font.name = theme["font_name"]
        y += label_h_in + 0.25
    if ref_str:
        rt = slide.shapes.add_textbox(left, Inches(y), width, Inches(ref_h_in))
        rp = rt.text_frame.paragraphs[0]
        rp.text = ref_str
        rp.alignment = PP_ALIGN.CENTER
        rp.font.size = Pt(int(typography_pt(style, "reference_size_bottom_pt", 22)))
        rp.font.italic = True
        rp.font.color.rgb = theme["reference"]
        if theme["font_name"]:
            rp.font.name = theme["font_name"]
    return slide


def add_quote_slide(prs, text, style=None, reference=None, font_size=None, alignment="center"):
    """Centred pull-quote slide."""
    style = style or {}
    theme = _resolve_theme(style)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _apply_slide_background(slide, style, prs)
    from .layout_tokens import pip_top_inches

    quote_pt = int(font_size or typography_pt(style, "quote_size_pt", 36))
    left, width, width_in, _ = content_box(prs, style, "quote")
    slide_h_in = prs.slide_height.inches
    slide_w_in = prs.slide_width.inches
    pip_top = pip_top_inches(style, slide_h_in, slide_w_in)
    top_in = float(layout_in(style, "quote", "top_in", 1.6))
    body_h_in = max(pip_top - top_in - 0.35, 2.5)
    align = _resolve_alignment(alignment)
    tb = slide.shapes.add_textbox(left, Inches(top_in), width, Inches(body_h_in * 0.68))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    _write_body_paragraph(p, text, quote_pt, theme, style=style, alignment=align)
    ref_str = (reference or "").strip()
    if ref_str:
        ref_pt = int(typography_pt(style, "reference_size_pt", 24))
        ref_y = top_in + body_h_in * 0.72
        ref_h = max(min(pip_top - ref_y - 0.15, 1.0), 0.55)
        rt = slide.shapes.add_textbox(left, Inches(ref_y), width, Inches(ref_h))
        rp = rt.text_frame.paragraphs[0]
        rp.text = f"— {ref_str}"
        rp.alignment = align
        rp.font.size = Pt(ref_pt)
        rp.font.italic = True
        rp.font.color.rgb = theme["reference"]
        if theme["font_name"]:
            rp.font.name = theme["font_name"]
    return slide


def add_picture_text_slide(prs, image_path, text, style=None, image_side="left",
                           image_fit="contain", font_size=28, alignment="left",
                           source_file=None):
    """Picture + text side by side (Text and Object layout)."""
    import os
    style = dict(style or {})
    if source_file:
        style["_source_file"] = source_file
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _apply_slide_background(slide, style, prs)
    theme = _resolve_theme(style)
    margin_in = float(layout_in(style, "picture_text", "margin_in", 0.35))
    gap_in = float(layout_in(style, "picture_text", "column_gap_in", 0.35))
    ratio = float(layout_in(style, "picture_text", "image_width_ratio", 0.48))
    margin = Inches(margin_in)
    slide_w = prs.slide_width
    slide_h = prs.slide_height
    content_w = slide_w - margin * 2
    content_h = slide_h - margin * 2
    img_w = int(content_w * ratio)
    txt_w = int(content_w) - img_w - Inches(gap_in)
    side = (image_side or "left").lower()
    if side not in ("left", "right"):
        side = "left"
    img_left = margin if side == "left" else margin + txt_w + Inches(gap_in)
    txt_left = margin + img_w + Inches(gap_in) if side == "left" else margin
    resolved = resolve_asset_path(image_path, source_file=source_file)
    path = resolved if resolved else image_path
    fit = (image_fit or "contain").lower()
    if fit not in ("contain", "cover", "fill"):
        fit = "contain"
    if path and os.path.exists(path):
        if fit == "fill":
            slide.shapes.add_picture(path, img_left, margin, width=img_w, height=content_h)
        else:
            pic = slide.shapes.add_picture(path, img_left, margin, width=img_w)
            _fit_picture_in_box(pic, img_left, margin, img_w, content_h, fit)
    else:
        print(f"Warning: Image not found: {image_path}")
    align = _resolve_alignment(alignment)
    tb = slide.shapes.add_textbox(txt_left, margin, txt_w, content_h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    _write_body_paragraph(tf.paragraphs[0], text, font_size, theme, style=style, alignment=align)
    return slide


def _table_palette(style: dict, theme: dict) -> dict:
    """Contrasting header/body fills — avoids PowerPoint default light zebra stripes."""
    dark = theme["dark_mode"]
    if dark:
        return {
            "header_fill": _parse_color(layout_in(style, "table", "header_fill", "#2563EB")),
            "header_text": _parse_color(layout_in(style, "table", "header_text", "#FFFFFF")),
            "row_fill": _parse_color(layout_in(style, "table", "row_fill", "#1F2937")),
            "row_alt_fill": _parse_color(layout_in(style, "table", "row_alt_fill", "#374151")),
            "body_text": theme["body"],
        }
    return {
        "header_fill": _parse_color(layout_in(style, "table", "header_fill", "#1E40AF")),
        "header_text": _parse_color(layout_in(style, "table", "header_text", "#FFFFFF")),
        "row_fill": _parse_color(layout_in(style, "table", "row_fill", "#F3F4F6")),
        "row_alt_fill": _parse_color(layout_in(style, "table", "row_alt_fill", "#E5E7EB")),
        "body_text": _parse_color(layout_in(style, "table", "body_text", "#111827")),
    }


def _apply_table_cell(cell, text, *, font_pt, text_rgb, fill_rgb, bold=False, font_name=None):
    cell.text = str(text)
    tf = cell.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Pt(4)
    tf.margin_top = tf.margin_bottom = Pt(2)
    cell.fill.solid()
    cell.fill.fore_color.rgb = fill_rgb
    for para in tf.paragraphs:
        para.font.size = Pt(font_pt)
        para.font.color.rgb = text_rgb
        para.font.bold = bold
        if font_name:
            para.font.name = font_name


def _fit_table_layout(rows, width_in, usable_h_in, font_pt, header_row, min_pt):
    """Shrink font and row heights until the table fits the vertical budget."""
    row_count = len(rows)
    col_count = max(len(r) for r in rows)
    col_w_in = width_in / max(col_count, 1)
    min_pt = int(min_pt)
    pt = int(font_pt)
    row_heights = []
    for _ in range(12):
        row_heights = []
        for r_idx, row in enumerate(rows):
            max_lines = 1
            for c_idx in range(col_count):
                val = row[c_idx] if c_idx < len(row) else ""
                max_lines = max(max_lines, _estimate_text_lines(str(val), col_w_in, pt))
            line_h = (pt / 72.0) * 1.38
            row_heights.append(max(0.38, max_lines * line_h + 0.1))
        total = sum(row_heights)
        if total <= usable_h_in or pt <= min_pt:
            break
        pt = max(min_pt, int(pt * 0.88))
    total = sum(row_heights)
    if total > usable_h_in and total > 0:
        scale = usable_h_in / total
        row_heights = [h * scale for h in row_heights]
    return pt, row_heights


def add_table_slide(prs, rows, style=None, font_size=24, header_row=True, reference=None):
    """Table layout with explicit fills, word wrap, and PiP-aware vertical fit."""
    style = style or {}
    theme = _resolve_theme(style)
    palette = _table_palette(style, theme)
    if not rows:
        rows = [[" "]]
    row_count = len(rows)
    col_count = max(len(r) for r in rows)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _apply_slide_background(slide, style, prs)
    from .layout_tokens import pip_top_inches

    top_in = float(layout_in(style, "table", "top_in", 0.75))
    left, width, width_in, _margin_in = content_box(prs, style, "table")
    slide_h_in = prs.slide_height.inches
    slide_w_in = prs.slide_width.inches
    pip_top = pip_top_inches(style, slide_h_in, slide_w_in)
    bottom_in = float(layout_in(style, "table", "bottom_in", 0.35))
    ref_str = (reference or "").strip()
    ref_h_in = 0.0
    ref_pt = int(typography_pt(style, "reference_size_bottom_pt", 22))
    if ref_str:
        ref_lines = _estimate_text_lines(ref_str, width_in, ref_pt)
        ref_h_in = min(0.35 + ref_lines * 0.32, 1.05) + float(layout_in(style, "table", "ref_gap_in", 0.15))
    usable_h = max(1.2, pip_top - top_in - ref_h_in - bottom_in)
    min_pt = int(layout_in(style, "table", "min_font_pt", 11))
    body_pt, row_heights = _fit_table_layout(
        rows, width_in, usable_h, font_size, header_row, min_pt,
    )
    table_h_in = sum(row_heights)
    table = slide.shapes.add_table(
        row_count, col_count, left, Inches(top_in), width, Inches(table_h_in),
    ).table
    font_name = theme["font_name"]
    for r_idx, row in enumerate(rows):
        table.rows[r_idx].height = Inches(row_heights[r_idx])
        is_header = header_row and r_idx == 0
        if is_header:
            fill, text_rgb = palette["header_fill"], palette["header_text"]
        else:
            body_idx = r_idx - (1 if header_row else 0)
            fill = palette["row_alt_fill"] if body_idx % 2 else palette["row_fill"]
            text_rgb = palette["body_text"]
        for c_idx in range(col_count):
            val = row[c_idx] if c_idx < len(row) else ""
            _apply_table_cell(
                table.cell(r_idx, c_idx),
                val,
                font_pt=body_pt,
                text_rgb=text_rgb,
                fill_rgb=fill,
                bold=is_header,
                font_name=font_name,
            )
    if ref_str:
        ref_y = top_in + table_h_in + float(layout_in(style, "table", "ref_gap_in", 0.15))
        ref_box_h = min(ref_h_in, max(0.45, pip_top - ref_y - 0.12))
        rt = slide.shapes.add_textbox(left, Inches(ref_y), width, Inches(ref_box_h))
        rp = rt.text_frame.paragraphs[0]
        rp.text = ref_str
        rp.alignment = PP_ALIGN.LEFT
        rp.font.size = Pt(ref_pt)
        rp.font.italic = True
        rp.font.color.rgb = theme["reference"]
        if font_name:
            rp.font.name = font_name
    return slide


def _estimate_text_lines(text, width_in, pt_size):
    """Rough line count for word-wrapped text in a textbox."""
    chars_per_line = max(14, int(width_in * 72 / max(pt_size, 1) * 0.5))
    lines = 0
    for part in (text or '').split('\n'):
        part = part.strip()
        if not part:
            continue
        lines += max(1, (len(part) + chars_per_line - 1) // chars_per_line)
    return max(1, lines)


def _normalize_ref_position(value: str | None) -> str:
    """Verse/list reference placement: bottom (slide foot), below (under body), top."""
    pos = (value or 'bottom').lower().strip()
    if pos in ('bottom', 'below', 'top'):
        return pos
    return 'bottom'


def _estimate_verse_body_height_in(verse_text, content_w_in, font_size):
    """Approximate rendered verse body height in inches."""
    lines = 0
    for _, v_text in _parse_verse_lines(verse_text):
        lines += _estimate_text_lines(v_text, content_w_in, font_size)
    line_h_in = (float(font_size) / 72.0) * 1.45
    return max(0.35, lines * line_h_in + 0.12)


def _slide_content_width(prs, style, slide_type, default_margin_in=0.6):
    """Centred content width from slide size and optional layout tokens."""
    left, width, _, _ = content_box(prs, style, slide_type, default_margin_in)
    return left, width


def _render_title_textboxes(slide, prs, title, subtitle, style, theme):
    """Word-wrapped title and subtitle blocks, centred on the slide."""
    left, width = _slide_content_width(prs, style, 'title')
    title_top_in = layout_in(style, 'title', 'title_top_in', 2.5)
    title_pt = typography_pt(style, 'title_size_pt', 44)
    title_lines = _estimate_text_lines(title, width.inches, title_pt)
    title_h = Inches(min(0.55 + title_lines * 0.52, 3.2))
    tb = slide.shapes.add_textbox(left, Inches(title_top_in), width, title_h)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(title_pt)
    p.font.bold = True
    p.font.color.rgb = theme['title']
    if theme['font_name']:
        p.font.name = theme['font_name']
    if subtitle:
        subtitle_pt = typography_pt(style, 'subtitle_size_pt', 28)
        subtitle_lines = _estimate_text_lines(subtitle, width.inches, subtitle_pt)
        gap = layout_in(style, 'title', 'subtitle_gap_in', 0.25)
        subtitle_top_in = title_top_in + title_h.inches + gap
        subtitle_h = Inches(min(0.45 + subtitle_lines * 0.38, 2.5))
        tb2 = slide.shapes.add_textbox(left, Inches(subtitle_top_in), width, subtitle_h)
        tf2 = tb2.text_frame
        tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        p2.text = subtitle
        p2.alignment = PP_ALIGN.CENTER
        p2.font.size = Pt(subtitle_pt)
        p2.font.color.rgb = theme['subtitle']
        if theme['font_name']:
            p2.font.name = theme['font_name']


def add_title_slide(prs, title, subtitle="", style=None):
    """
    Add a title slide. Custom word-wrapped layout when a background is set or
    the subtitle is long; otherwise uses the default template layout.
    """
    style = style or {}
    has_background = bool(style.get('background_image') or style.get('background_color'))
    theme = _resolve_theme(style)
    use_custom = has_background or bool(subtitle and len(subtitle) > title_custom_threshold(style))

    if use_custom:
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
        _apply_slide_background(slide, style, prs)
        _render_title_textboxes(slide, prs, title, subtitle, style, theme)
    else:
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        _apply_slide_background(slide, style, prs)
        title_shape = slide.shapes.title
        title_shape.text = title
        if theme['font_name']:
            title_shape.text_frame.paragraphs[0].font.name = theme['font_name']
        if subtitle and len(slide.placeholders) > 1:
            sub = slide.placeholders[1]
            sub.text = subtitle
            if theme['font_name']:
                sub.text_frame.paragraphs[0].font.name = theme['font_name']

    return slide


def add_section_slide(prs, section_name, style=None, section_subtitle=None):
    """
    Add a section title slide — centred, same layout for all themes.
    Background image/color applied when present in slide_style.

    Optional ``section_subtitle`` (str): second line below the title, smaller type
    (e.g. ``Section 1``). Set from YAML as ``section_subtitle`` next to ``section``.
    """
    style = style or {}
    theme = _resolve_theme(style)

    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank for full control
    _apply_slide_background(slide, style, prs)

    name = section_name or ''
    sub = (section_subtitle or '').strip()
    line_count = name.count('\n') + 1 if name else 1
    section_title_pt = typography_pt(style, 'section_title_size_pt', 44)
    section_sub_pt = typography_pt(style, 'section_subtitle_size_pt', 24)
    tb_h_in = 1.5 if line_count <= 1 else min(1.2 + line_count * 0.55, 4.5)
    left, tb_w, tb_w_in, margin_in = content_box(prs, style, 'section')
    if sub:
        sub_lines = _estimate_text_lines(sub, tb_w_in, section_sub_pt)
        tb_h_in += min(0.45 + sub_lines * 0.38, 2.5)
    tb_h = Inches(tb_h_in)
    left = (prs.slide_width - tb_w) / 2
    top = (prs.slide_height - tb_h) / 2

    tb = slide.shapes.add_textbox(left, top, tb_w, tb_h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = section_name
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(section_title_pt)
    p.font.bold = True
    p.font.color.rgb = theme['section']
    if theme['font_name']:
        p.font.name = theme['font_name']

    if sub:
        p.space_after = Pt(20)
        sr, sg, sb = list(theme['section'])
        dim = float(layout_in(style, 'section', 'subtitle_dim_factor', 0.76))
        sub_rgb = RGBColor(
            max(0, min(255, int(sr * dim))),
            max(0, min(255, int(sg * dim))),
            max(0, min(255, int(sb * dim))),
        )
        sp = tf.add_paragraph()
        sp.text = sub
        sp.alignment = PP_ALIGN.CENTER
        sp.font.size = Pt(section_sub_pt)
        sp.font.bold = False
        sp.font.color.rgb = sub_rgb
        if theme['font_name']:
            sp.font.name = theme['font_name']

    return slide


def _parse_color(color_value):
    """Parse a color value (named string or hex) into RGBColor."""
    NAMED_COLORS = {
        'orange': RGBColor(255, 140, 0),
        'yellow': RGBColor(255, 215, 0),
        'red':    RGBColor(220, 50,  50),
        'green':  RGBColor(50,  180, 50),
        'blue':   RGBColor(30,  100, 220),
        'white':  RGBColor(255, 255, 255),
        'cyan':   RGBColor(0,   200, 200),
        'purple': RGBColor(150, 50,  200),
    }
    if not color_value:
        return NAMED_COLORS['orange']
    if isinstance(color_value, str):
        if color_value.lower() in NAMED_COLORS:
            return NAMED_COLORS[color_value.lower()]
        # Hex string e.g. "#FF8C00" or "FF8C00"
        lower = color_value.strip('#')
        if len(lower) == 6:
            try:
                r, g, b = int(lower[0:2], 16), int(lower[2:4], 16), int(lower[4:6], 16)
                return RGBColor(r, g, b)
            except ValueError:
                pass
    return NAMED_COLORS['orange']


def _resolve_theme(style: dict) -> dict:
    """
    Resolve all display colors from a slide_style dict with smart defaults.

    When background_image or background_color is set, text auto-defaults to
    white (dark-background mode). Any individual key in slide_style overrides
    the auto-default.

    JSON keys supported (all optional):
        text_color, reference_color, title_color, subtitle_color,
        section_title_color, highlight_color, annotation_color,
        reference_position, alignment
    """
    style = style or {}
    has_dark_bg = bool(style.get('background_image') or style.get('background_color'))
    raw_text = str(style.get('text_color') or '').lower().strip()
    if raw_text:
        dark_mode = raw_text in ('white', '#ffffff', 'ffffff')
    else:
        dark_mode = has_dark_bg  # auto: dark bg → white text

    def _rc(key, light, dark):
        raw = style.get(key, '')
        return _parse_color(raw) if raw else (_parse_color(dark) if dark_mode else _parse_color(light))

    return {
        'dark_mode':        dark_mode,
        'body':             _rc('text_color',          '#1A1A2E', '#FFFFFF'),
        'reference':        _rc('reference_color',     '#404040', '#CCCCCC'),
        'title':            _rc('title_color',         '#1A1A2E', '#FFFFFF'),
        'subtitle':         _rc('subtitle_color',      '#505050', '#AAAAAA'),
        'section':          _rc('section_title_color', '#003366', '#FFFFFF'),
        'highlight':        _rc('highlight_color',     '#FF8C00', '#FFD700'),
        'annotation':       _rc('annotation_color',    '#1E50C8', '#1E50C8'),
        'ref_position':     _normalize_ref_position(style.get('reference_position')),
        'global_alignment': style.get('alignment', 'left'),
        'font_name':        style.get('font_name') or 'Palatino',
    }


def _normalise_highlights(highlights, highlight_rgb=None):
    """
    Normalise highlights list. highlight_rgb overrides the default orange.
    """
    default_hl = highlight_rgb or RGBColor(255, 140, 0)
    BUBBLES = {1: '\u2776', 2: '\u2777', 3: '\u2778', 4: '\u2779', 5: '\u277a',
                6: '\u277b', 7: '\u277c', 8: '\u277d', 9: '\u277e'}
    result = []
    for h in highlights:
        if isinstance(h, str):
            result.append({'text': h, 'color': default_hl,
                           'bold': True, 'italic': False,
                           'underline': False, 'annotation': None})
        elif isinstance(h, dict) and h.get('text'):
            ann = h.get('annotation', None)
            result.append({
                'text': h['text'],
                'color': _parse_color(h.get('color', 'orange')),
                'bold': h.get('bold', True),
                'italic': h.get('italic', False),
                'underline': h.get('underline', True if ann else False),
                'annotation': BUBBLES.get(ann) if ann else None,
            })
    return result


def _apply_highlights(paragraph, text, highlights, large_text=None,
                      body_rgb=None, highlight_rgb=None, annotation_rgb=None,
                      font_name=None, base_font_size=32, annotation_size_pt=46):
    """
    Apply per-phrase rich text formatting.
    body_rgb, highlight_rgb, annotation_rgb, font_name all come from _resolve_theme.
    base_font_size: point size for normal body and highlight runs (``large_text`` overrides per match).
    """
    import re
    _body = body_rgb or RGBColor(26, 26, 46)
    _ann  = annotation_rgb or RGBColor(30, 80, 200)
    _base = int(base_font_size) if base_font_size else 32

    def _sf(run, size_pt):
        """Set font size and optional font name on a run."""
        run.font.size = Pt(size_pt)
        if font_name:
            run.font.name = font_name

    matches = []

    if highlights:
        for fmt in _normalise_highlights(highlights, highlight_rgb=highlight_rgb):
            pattern = re.escape(fmt['text'])
            for match in re.finditer(pattern, text, re.IGNORECASE):
                matches.append((match.start(), match.end(), match.group(), 'highlight', fmt))

    if large_text:
        for word, font_size in large_text.items():
            pattern = re.escape(word)
            for match in re.finditer(pattern, text, re.IGNORECASE):
                matches.append((match.start(), match.end(), match.group(), 'large', font_size))

    matches.sort(key=lambda x: x[0])
    filtered, last_end = [], -1
    for m in matches:
        if m[0] >= last_end:
            filtered.append(m)
            last_end = m[1]

    _body = body_rgb or RGBColor(26, 26, 46)

    if not filtered:
        run = paragraph.add_run()
        run.text = text
        _sf(run, _base)
        run.font.color.rgb = _body
        if font_name:
            run.font.name = font_name
        return

    current_pos = 0

    for start, end, matched_text, fmt_type, fmt in filtered:
        # Plain text before this match
        if start > current_pos:
            run = paragraph.add_run()
            run.text = text[current_pos:start]
            _sf(run, _base)
            run.font.color.rgb = _body
            run.font.bold = False
            run.font.italic = False
            run.font.underline = False

        # Formatted run — always add_run() to preserve any pre-existing runs
        run = paragraph.add_run()
        run.text = matched_text

        if fmt_type == 'highlight':
            _sf(run, _base)
            run.font.color.rgb = fmt['color']
            run.font.bold = fmt['bold']
            run.font.italic = fmt['italic']
            run.font.underline = fmt['underline']
            if fmt.get('annotation'):
                ann_run = paragraph.add_run()
                ann_run.text = fmt['annotation']
                _sf(ann_run, int(annotation_size_pt))
                ann_run.font.bold = False
                ann_run.font.color.rgb = _ann
                rPr = ann_run._r.get_or_add_rPr()
                rPr.set('baseline', '30000')
        elif fmt_type == 'large':
            _sf(run, fmt)
            run.font.color.rgb = _body

        current_pos = end

    # Remaining plain text
    if current_pos < len(text):
        run = paragraph.add_run()
        run.text = text[current_pos:]
        _sf(run, _base)
        run.font.color.rgb = _body
        run.font.bold = False
        run.font.italic = False
        run.font.underline = False

def _resolve_alignment(align_str):
    """Convert alignment string to PP_ALIGN constant."""
    return {
        'left':   PP_ALIGN.LEFT,
        'right':  PP_ALIGN.RIGHT,
        'center': PP_ALIGN.CENTER,
    }.get((align_str or 'center').lower(), PP_ALIGN.CENTER)


def add_list_slide(prs, items, reference, list_type='bullet', font_size=32,
                   alignment='left', style=None):
    """
    Add a bullet/numbered list slide. Default alignment is left.
    All colors and font resolved via slide_style.
    Respects slide_style reference_position (top, bottom, or below).
    """
    style = style or {}
    theme = _resolve_theme(style)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _apply_slide_background(slide, style, prs)

    from .layout_tokens import pip_top_inches

    ref_position = theme['ref_position']
    margin_in = float(layout_in(style, 'list', 'margin_in', 0.75))
    slide_h_in = prs.slide_height.inches
    slide_w_in = prs.slide_width.inches
    content_w_in = content_width_inches(prs, style, 'list', margin_in)
    pip_top = pip_top_inches(style, slide_h_in, slide_w_in)
    ref_pt = int(typography_pt(style, 'section_title_size_pt', 36) * 0.82)
    body_pt = int(font_size or typography_pt(style, 'body_size_pt', 28) * 0.88)

    ref_h_in = 0.0
    list_top_in = 0.75
    if reference and ref_position == 'top':
        ref_lines = _estimate_text_lines(reference, content_w_in, ref_pt)
        ref_h_in = min(0.4 + ref_lines * 0.38, 1.4)
        ref_tb = slide.shapes.add_textbox(
            Inches(margin_in), Inches(list_top_in), Inches(content_w_in), Inches(ref_h_in),
        )
        ref_tf = ref_tb.text_frame
        ref_tf.word_wrap = True
        ref_p = ref_tf.paragraphs[0]
        ref_p.text = reference
        ref_p.alignment = PP_ALIGN.LEFT
        ref_p.font.size = Pt(ref_pt)
        ref_p.font.bold = True
        ref_p.font.color.rgb = theme['title']
        if theme['font_name']:
            ref_p.font.name = theme['font_name']
        list_top_in += ref_h_in + float(layout_in(style, 'list', 'ref_gap_in', 0.18))

    item_line_h = body_pt / 72.0 * 0.62
    ref_gap = float(layout_in(style, 'list', 'ref_gap_in', 0.18))
    bottom_ref_h = 0.0
    if reference and ref_position in ('bottom', 'below'):
        ref_pt_b = int(typography_pt(style, 'list_ref_bottom_pt', 22))
        ref_lines_b = _estimate_text_lines(reference, content_w_in, ref_pt_b)
        bottom_ref_h = min(0.35 + ref_lines_b * 0.32, 1.0) + ref_gap

    list_max_bottom = min(slide_h_in - 0.35, pip_top - 0.25)
    if bottom_ref_h:
        list_max_bottom = min(list_max_bottom, slide_h_in - bottom_ref_h - 0.35)

    est_lines = sum(max(1, _estimate_text_lines(str(it), content_w_in, body_pt)) for it in items)
    list_h_in = max(est_lines * (item_line_h + 0.12), 1.0)
    list_h_in = min(list_h_in, max(1.0, list_max_bottom - list_top_in))

    block_h = list_h_in + (ref_h_in + ref_gap if reference and ref_position == 'top' else 0)
    block_h += bottom_ref_h if bottom_ref_h else 0
    if list_top_in + block_h > pip_top - 0.25:
        list_top_in = max(0.55, pip_top - block_h - 0.25)
        list_h_in = min(list_h_in, max(1.0, list_max_bottom - list_top_in))

    tb = slide.shapes.add_textbox(
        Inches(margin_in), Inches(list_top_in), Inches(content_w_in), Inches(list_h_in),
    )
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    align = _resolve_alignment(alignment)

    for idx, item in enumerate(items):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        prefix = f"{idx + 1}. " if list_type == 'numbered' else "\u2022  "
        p.text = prefix + item
        p.alignment = align
        p.font.size = Pt(body_pt)
        p.font.color.rgb = theme['body']
        if theme['font_name']:
            p.font.name = theme['font_name']
        p.space_after = Pt(12)
        p.line_spacing = 1.15

    if reference and ref_position in ('bottom', 'below'):
        ref_pt = int(typography_pt(style, 'list_ref_bottom_pt', 22))
        ref_lines = _estimate_text_lines(reference, content_w_in, ref_pt)
        ref_h_in = min(0.35 + ref_lines * 0.32, 1.0)
        ref_y_in = min(list_top_in + list_h_in + ref_gap, pip_top - ref_h_in - 0.15)
        ref_y_in = max(ref_y_in, list_top_in + list_h_in + 0.08)
        ref_tb = slide.shapes.add_textbox(
            Inches(margin_in), Inches(ref_y_in), Inches(content_w_in), Inches(ref_h_in),
        )
        ref_tf = ref_tb.text_frame
        ref_tf.word_wrap = True
        ref_p = ref_tf.paragraphs[0]
        ref_p.text = reference
        ref_p.alignment = PP_ALIGN.LEFT
        ref_p.font.size = Pt(ref_pt)
        ref_p.font.color.rgb = theme['reference']
        ref_p.font.italic = True
        if theme['font_name']:
            ref_p.font.name = theme['font_name']

    return slide


def _fit_picture_in_box(pic, box_left, box_top, box_w, box_h, fit):
    """Scale picture for contain/cover within a content box; crop cover overflow."""
    scale_w = box_w / pic.width
    scale_h = box_h / pic.height
    scale = min(scale_w, scale_h) if fit == 'contain' else max(scale_w, scale_h)
    pic.width = int(pic.width * scale)
    pic.height = int(pic.height * scale)
    pic.left = box_left + (box_w - pic.width) // 2
    pic.top = box_top + (box_h - pic.height) // 2
    if fit != 'cover':
        return
    crop_l = max(0.0, float(box_left - pic.left) / float(pic.width))
    crop_r = max(0.0, float(pic.left + pic.width - (box_left + box_w)) / float(pic.width))
    crop_t = max(0.0, float(box_top - pic.top) / float(pic.height))
    crop_b = max(0.0, float(pic.top + pic.height - (box_top + box_h)) / float(pic.height))
    if crop_l or crop_r or crop_t or crop_b:
        pic.crop_left = crop_l
        pic.crop_right = crop_r
        pic.crop_top = crop_t
        pic.crop_bottom = crop_b
        pic.left = box_left
        pic.top = box_top
        pic.width = box_w
        pic.height = box_h


def add_image_slide(prs, image_path, style=None, caption=None, reference=None,
                    image_fit='contain', source_file=None):
    """
    Add a slide with an embedded image and optional caption.

    ``image_fit``: ``contain`` (default, keep aspect ratio), ``cover`` (fill area,
    keep ratio, may crop), or ``fill`` (stretch to content box).
    ``reference`` / ``caption`` appear below the image when provided.
    """
    import os
    style = dict(style or {})
    if source_file:
        style['_source_file'] = source_file

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _apply_slide_background(slide, style, prs)

    resolved = resolve_asset_path(image_path, source_file=source_file)
    path = resolved if resolved else image_path
    if not path or not os.path.exists(path):
        print(f"Warning: Image not found: {image_path}")
        return slide

    theme = _resolve_theme(style)
    slide_w = prs.slide_width
    slide_h = prs.slide_height
    caption_lines = []
    if reference and str(reference).strip():
        caption_lines.append(str(reference).strip())
    if caption and str(caption).strip():
        caption_lines.append(str(caption).strip())
    caption_h_in = float(layout_in(style, 'image', 'caption_height_in', 0.9))
    caption_h = Inches(caption_h_in) if caption_lines else Inches(0)
    fit = (image_fit or 'contain').lower()
    margin_in = float(layout_in(style, 'image', 'margin_in', 0.35))
    margin = Inches(0) if not caption_lines and fit in ('cover', 'fill') else Inches(margin_in)

    box_top = margin
    box_h = slide_h - margin * 2 - caption_h
    box_w = slide_w - margin * 2
    if fit not in ('contain', 'cover', 'fill'):
        fit = 'contain'

    if fit == 'fill':
        slide.shapes.add_picture(path, margin, box_top, width=box_w, height=box_h)
    else:
        pic = slide.shapes.add_picture(path, margin, box_top, width=box_w)
        _fit_picture_in_box(pic, margin, box_top, box_w, box_h, fit)

    if caption_lines:
        cap_ref_pt = int(typography_pt(style, 'caption_ref_size_pt', 22))
        cap_body_pt = int(typography_pt(style, 'caption_body_size_pt', 18))
        cap_top = slide_h - margin - caption_h
        cap_tb = slide.shapes.add_textbox(margin, cap_top, box_w, caption_h)
        cap_tf = cap_tb.text_frame
        cap_tf.word_wrap = True
        for i, line in enumerate(caption_lines):
            p = cap_tf.paragraphs[0] if i == 0 else cap_tf.add_paragraph()
            p.text = line
            p.alignment = PP_ALIGN.CENTER
            p.font.size = Pt(cap_ref_pt if i == 0 and reference else cap_body_pt)
            p.font.bold = bool(i == 0 and reference)
            p.font.color.rgb = theme['reference'] if i == 0 and reference else theme['body']
            if theme['font_name']:
                p.font.name = theme['font_name']

    return slide


def _hebrew_font_name(style: dict) -> str:
    """Prefer a Hebrew-capable font when available."""
    custom = (style or {}).get("hebrew_font_name")
    if custom:
        return custom
    import os
    for path in (
        "/System/Library/Fonts/SFHebrew.ttf",
        "/System/Library/Fonts/Supplemental/Arial Hebrew.ttf",
        "/Library/Fonts/Arial Hebrew.ttf",
    ):
        if os.path.exists(path):
            return "Arial Hebrew" if "Arial" in path else "SF Hebrew"
    return (style or {}).get("font_name") or "Palatino"


def _fill_hebrew_runs(paragraph, text, highlight_substring, body_rgb, highlight_rgb, font_size_pt, font_name):
    """Write Hebrew text; colour ``highlight_substring`` in the right-hand name."""
    hl = highlight_substring or ""
    if not hl or hl not in text:
        run = paragraph.add_run()
        run.text = text
        run.font.size = Pt(font_size_pt)
        run.font.color.rgb = body_rgb
        if font_name:
            run.font.name = font_name
        return

    parts = text.split(hl, 1)
    before, after = parts[0], parts[1] if len(parts) > 1 else ""

    def _add(part, rgb):
        if not part:
            return
        r = paragraph.add_run()
        r.text = part
        r.font.size = Pt(font_size_pt)
        r.font.color.rgb = rgb
        if font_name:
            r.font.name = font_name

    _add(before, body_rgb)
    _add(hl, highlight_rgb)
    _add(after, body_rgb)


def add_hebrew_rename_slide(prs, rows, style=None, font_size=110, reference=None, caption=None,
                            highlight_color=None):
    """
    Slide matching the original Why Delay Hebrew layout: large names left/right,
    purple highlight on the changed letter in the new name (e.g. הָ in אַבְרָהָם).
    """
    from pptx.enum.shapes import MSO_CONNECTOR_TYPE

    style = style or {}
    theme = _resolve_theme(style)
    body_rgb = theme["body"]
    hl_hex = highlight_color or style.get("hebrew_highlight_color") or "#9900FF"
    hl_rgb = _parse_color(hl_hex) or RGBColor(153, 0, 255)
    hebrew_pt = int(font_size or style.get("hebrew_font_size") or 110)
    fn = _hebrew_font_name(style)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _apply_slide_background(slide, style, prs)

    slide_w_in = prs.slide_width.inches
    ref_w_in = float(layout_in(style, 'hebrew_rename', 'reference_width_in', 10.0))
    sx = slide_w_in / ref_w_in
    row_y = layout_in(style, 'hebrew_rename', 'row_y_in', [1.15, 4.05])
    if not isinstance(row_y, list):
        row_y = [1.15, 4.05]
    left_x = float(layout_in(style, 'hebrew_rename', 'left_x_factor', 0.35)) * sx
    right_x = float(layout_in(style, 'hebrew_rename', 'right_x_factor', 5.15)) * sx
    box_w = float(layout_in(style, 'hebrew_rename', 'box_width_factor', 4.2)) * sx
    box_h = float(layout_in(style, 'hebrew_rename', 'box_height_in', 1.35))

    for i, row in enumerate(rows[:2]):
        left = (row.get("left") or "").strip()
        right = (row.get("right") or "").strip()
        hl = (row.get("highlight_in_right") or "ה").strip()
        y = Inches(row_y[i] if i < len(row_y) else row_y[-1] + i * 2.9)

        lt = slide.shapes.add_textbox(Inches(left_x), y, Inches(box_w), Inches(box_h))
        lt_tf = lt.text_frame
        lt_tf.word_wrap = False
        lp = lt_tf.paragraphs[0]
        lp.alignment = PP_ALIGN.CENTER
        lr = lp.add_run()
        lr.text = left
        lr.font.size = Pt(hebrew_pt)
        lr.font.color.rgb = body_rgb
        if fn:
            lr.font.name = fn

        rt = slide.shapes.add_textbox(Inches(right_x), y, Inches(box_w), Inches(box_h))
        rt_tf = rt.text_frame
        rt_tf.word_wrap = False
        rp = rt_tf.paragraphs[0]
        rp.alignment = PP_ALIGN.CENTER
        _fill_hebrew_runs(rp, right, hl, body_rgb, hl_rgb, hebrew_pt, fn)

        y_mid = y.inches + box_h * 0.45
        x1 = left_x + box_w * 0.85
        x2 = right_x + box_w * 0.05
        conn = slide.shapes.add_connector(
            MSO_CONNECTOR_TYPE.STRAIGHT,
            Inches(x1), Inches(y_mid), Inches(x2), Inches(y_mid),
        )
        conn.line.color.rgb = theme["highlight"]
        conn.line.width = Pt(2.5)

    cap_h_in = float(layout_in(style, 'hebrew_rename', 'caption_height_in', 0.85))
    cap_h = Inches(cap_h_in)
    if reference or caption:
        cap_bottom = float(layout_in(style, 'hebrew_rename', 'caption_bottom_in', 0.45))
        cap_margin = float(layout_in(style, 'hebrew_rename', 'caption_margin_in', 0.5))
        cap_top = prs.slide_height - Inches(cap_bottom) - cap_h
        cap_tb = slide.shapes.add_textbox(Inches(cap_margin), cap_top, prs.slide_width - Inches(cap_margin * 2), cap_h)
        cap_tf = cap_tb.text_frame
        cap_tf.word_wrap = True
        cap_ref_pt = int(typography_pt(style, 'caption_ref_size_pt', 22))
        cap_body_pt = int(typography_pt(style, 'caption_body_size_pt', 18))
        lines = []
        if reference:
            lines.append(str(reference).strip())
        if caption:
            lines.append(str(caption).strip())
        for idx, line in enumerate(lines):
            p = cap_tf.paragraphs[0] if idx == 0 else cap_tf.add_paragraph()
            p.text = line
            p.alignment = PP_ALIGN.CENTER
            p.font.size = Pt(cap_ref_pt if idx == 0 and reference else cap_body_pt)
            p.font.bold = bool(idx == 0 and reference)
            p.font.color.rgb = theme["reference"] if idx == 0 and reference else theme["body"]
            if theme["font_name"]:
                p.font.name = theme["font_name"]

    return slide


def _parse_verse_lines(text):
    """
    Parse verse text into [(verse_num_or_None, line_text), ...] pairs.
    Detects lines (or inline segments) starting with 1–3 digit verse numbers.
    e.g. '11 For the grace...\n12 Teaching us...' → [('11','For the grace...'),('12','Teaching us...')]
    e.g. '1 Therefore, holy brethren...' → [('1', 'Therefore, holy brethren...')]
    Returns [(None, full_text)] if no verse numbers detected.
    """
    import re
    VERSE_NUM_RE = re.compile(r"^(\d{1,3})\s+(.*)", re.DOTALL)

    # Books that appear with a numeric prefix (1/2/3 Timothy, etc.)
    NUMBERED_BOOKS = frozenset([
        "timothy", "corinthians", "thessalonians", "peter", "john",
        "chronicles", "samuel", "kings", "esdras", "maccabees",
    ])

    raw_lines = [l.strip() for l in text.split("\n") if l.strip()]
    result = []
    for line in raw_lines:
        m = VERSE_NUM_RE.match(line)
        if m:
            num_str = m.group(1)
            remainder = m.group(2)
            # If 1/2/3 followed by a numbered Bible book name -> plain text
            first_word = remainder.split()[0].lower().rstrip(",:)") if remainder.split() else ""
            if int(num_str) in (1, 2, 3) and first_word in NUMBERED_BOOKS:
                result.append((None, line))
            else:
                result.append((num_str, remainder))
        else:
            result.append((None, line))

    # If no verse numbers found at all, return original
    if not any(num for num, _ in result):
        return [(None, text)]
    return result


def _add_superscript_num_run(paragraph, num_str, font_size, body_rgb, font_name):
    """Add a small superscript verse-number run to a paragraph."""
    from pptx.oxml.ns import qn
    run = paragraph.add_run()
    run.text = num_str + '\u2009'  # narrow space after number
    run.font.size = Pt(int(font_size * 0.52))
    run.font.color.rgb = body_rgb
    run.font.bold = False
    if font_name:
        run.font.name = font_name
    # Set superscript baseline (30 000 = 30% above normal)
    rPr = run._r.get_or_add_rPr()
    rPr.set('baseline', '30000')


def add_verse_slide(prs, verse_text, reference, part_num=None, highlights=None,
                    large_text=None, alignment='left', font_size=32, style=None,
                    reference_font_size=None, leading_title=None,
                    text_below_reference=None, text_below_reference_highlights=None,
                    text_below_reference_large_text=None, reference_position=None):
    """
    Add a verse slide. All colors and font resolved via slide_style.
    Supported slide_style keys: background_image, background_color,
    text_color, reference_color, highlight_color, annotation_color,
    title_color, section_title_color, font_name,
    reference_position ('bottom'/'below'/'top'), alignment.

    ``bottom`` (default): reference anchored at the foot of the slide.
    ``below``: reference placed directly under the verse body text.
    ``top``: reference above the verse body.

    Optional verse YAML key ``leading_title`` (str): large line at the top of the
    slide; the reference is drawn directly under that title, then the verse body
    beneath the reference. Optional ``text_below_reference`` adds a further block
    below the verse body (e.g. a second passage on the same slide).
    Optional ``text_below_reference_highlights`` / ``text_below_reference_large_text``:
    same shape as ``highlights`` / ``large_text`` for that block.
    """
    style = style or {}
    theme = _resolve_theme(style)
    ann_pt = int(typography_pt(style, 'annotation_size_pt', 46))
    left, content_w, content_w_in, margin_in = content_box(prs, style, 'verse')
    verse_slide = prs.slides.add_slide(prs.slide_layouts[6])
    _apply_slide_background(verse_slide, style, prs)
    align = _resolve_alignment(alignment)
    ref_position = _normalize_ref_position(reference_position or theme['ref_position'])
    fn = theme['font_name']
    leading = (leading_title or '').strip()
    extra_ref = (text_below_reference or '').strip()

    slide_h_in = prs.slide_height.inches
    bottom_margin_in = float(layout_in(style, 'verse', 'bottom_margin_in', 0.15))
    leading_top_in = float(layout_in(style, 'verse', 'leading_title_top_in', 0.35))
    leading_pt = int(typography_pt(style, 'leading_title_size_pt', 38))
    ref_top_in = float(layout_in(style, 'verse', 'ref_top_in', 0.3))
    body_gap_in = float(layout_in(style, 'verse', 'body_gap_in', 0.15))
    default_body_h_in = float(layout_in(style, 'verse', 'default_body_height_in', 4.5))
    bottom_ref_top_in = float(layout_in(style, 'verse', 'bottom_ref_top_in', 6.0))
    bottom_ref_h_in = float(layout_in(style, 'verse', 'bottom_ref_height_in', 0.7))
    no_ref_top_in = float(layout_in(style, 'verse', 'no_ref_body_top_in', 1.5))
    no_ref_h_in = float(layout_in(style, 'verse', 'no_ref_body_height_in', 3.8))
    extra_reserve_in = float(layout_in(style, 'verse', 'extra_ref_reserve_in', 1.32))
    title_ref_gap_in = float(layout_in(style, 'verse', 'leading_title_ref_gap_in', 0.2))
    ref_h_default_in = float(layout_in(style, 'verse', 'ref_height_in', 0.7))
    ref_h_large_in = float(layout_in(style, 'verse', 'ref_height_large_in', 0.95))
    ref_pt_default = int(typography_pt(style, 'reference_size_pt', 28))
    ref_pt_small = int(typography_pt(style, 'reference_size_small_pt', 24))
    ref_pt_bottom = int(typography_pt(style, 'reference_size_bottom_pt', 22))

    leading_verse_top_in = None
    leading_verse_h_in = None

    def _set_ref(tb_shape, text, align_const, size, bold=False, italic=False):
        rp = tb_shape.text_frame.paragraphs[0]
        rp.text = text
        rp.alignment = align_const
        rp.font.size = Pt(size)
        rp.font.bold = bold
        rp.font.italic = italic
        rp.font.color.rgb = theme['reference']
        if fn:
            rp.font.name = fn

    if leading:
        leading_lines = max(1, leading.count('\n') + 1)
        leading_lines = max(leading_lines, _estimate_text_lines(leading, content_w_in, leading_pt))
        lt_h_in = min(0.45 + leading_lines * 0.42, 2.8)
        lt_tb = verse_slide.shapes.add_textbox(
            left, Inches(leading_top_in), content_w, Inches(lt_h_in),
        )
        lt_tf = lt_tb.text_frame
        lt_tf.word_wrap = True
        lt_tf.vertical_anchor = MSO_ANCHOR.TOP
        lt_p = lt_tf.paragraphs[0]
        lt_p.text = leading
        lt_p.alignment = align
        lt_p.font.size = Pt(leading_pt)
        lt_p.font.bold = True
        lt_p.font.color.rgb = theme['body']
        if fn:
            lt_p.font.name = fn
        if reference and ref_position == 'top':
            ref_text = reference + (f' (Part {part_num})' if part_num is not None else '')
            ref_pt = int(reference_font_size) if reference_font_size is not None else ref_pt_small
            ref_h_in = 0.72 if extra_ref else 0.85
            ref_y_in = leading_top_in + lt_h_in + title_ref_gap_in
            ref_tb = verse_slide.shapes.add_textbox(
                left, Inches(ref_y_in), content_w, Inches(ref_h_in),
            )
            _set_ref(ref_tb, ref_text, PP_ALIGN.LEFT, ref_pt, bold=True, italic=False)
            ref_tb.text_frame.paragraphs[0].font.color.rgb = theme['body']
            if fn:
                ref_tb.text_frame.paragraphs[0].font.name = fn
            vtop = ref_y_in + ref_h_in + 0.1
            if extra_ref:
                vh = slide_h_in - vtop - extra_reserve_in - bottom_margin_in
            else:
                vh = slide_h_in - vtop - bottom_margin_in
            leading_verse_top_in = vtop
            leading_verse_h_in = max(vh, 2.0)
            verse_top = Inches(vtop)
            verse_height = Inches(leading_verse_h_in)
        elif reference and ref_position in ('bottom', 'below'):
            ref_reserve_in = 0.85
            vtop = leading_top_in + lt_h_in + 0.15
            vh = slide_h_in - vtop - ref_reserve_in - bottom_margin_in
            leading_verse_top_in = vtop
            leading_verse_h_in = max(vh, 1.5)
            verse_top = Inches(vtop)
            verse_height = Inches(leading_verse_h_in)
        else:
            leading_verse_top_in = leading_top_in + lt_h_in + 0.15
            leading_verse_h_in = (
                slide_h_in - leading_verse_top_in - bottom_margin_in
                if not extra_ref else 3.25
            )
            verse_top = Inches(leading_verse_top_in)
            verse_height = Inches(max(leading_verse_h_in, 2.0))
    elif ref_position == 'top' and reference:
        ref_text = reference + (f' (Part {part_num})' if part_num is not None else '')
        ref_pt = int(reference_font_size) if reference_font_size is not None else ref_pt_default
        ref_h_in = ref_h_large_in if ref_pt >= 36 else ref_h_default_in
        ref_h = Inches(ref_h_in)
        ref_tb = verse_slide.shapes.add_textbox(left, Inches(ref_top_in), content_w, ref_h)
        _set_ref(ref_tb, ref_text, PP_ALIGN.LEFT, ref_pt, bold=True)
        ref_tb.text_frame.paragraphs[0].font.color.rgb = theme['body']
        verse_top = Inches(ref_top_in + ref_h_in + body_gap_in)
        verse_height = Inches(default_body_h_in)
    else:
        verse_top, verse_height = Inches(no_ref_top_in), Inches(no_ref_h_in)

    textbox = verse_slide.shapes.add_textbox(left, verse_top, content_w, verse_height)
    tf = textbox.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.alignment = align

    verse_lines = _parse_verse_lines(verse_text)
    has_verse_nums = any(num for num, _ in verse_lines)

    first_para = True
    for v_num, v_text in verse_lines:
        if first_para:
            p = tf.paragraphs[0]
            first_para = False
        else:
            p = tf.add_paragraph()
        p.alignment = align

        if (highlights and len(highlights) > 0) or (large_text and len(large_text) > 0):
            if has_verse_nums and v_num:
                _add_superscript_num_run(p, v_num, font_size, theme['body'], fn)
            _apply_highlights(p, v_text, highlights, large_text,
                              body_rgb=theme['body'],
                              highlight_rgb=theme['highlight'],
                              annotation_rgb=theme['annotation'],
                              font_name=fn,
                              base_font_size=int(font_size),
                              annotation_size_pt=ann_pt)
        else:
            if has_verse_nums and v_num:
                _add_superscript_num_run(p, v_num, font_size, theme['body'], fn)
            run = p.add_run()
            run.text = v_text
            run.font.size = Pt(font_size)
            run.font.color.rgb = theme['body']
            if fn:
                run.font.name = fn

    if leading and extra_ref and leading_verse_top_in is not None and leading_verse_h_in is not None:
        below_top = leading_verse_top_in + leading_verse_h_in + 0.06
        below_fs = max(int(font_size) - 2, 22)
        ex_h_in = max(slide_h_in - below_top - bottom_margin_in, 0.5)
        ex_tb = verse_slide.shapes.add_textbox(left, Inches(below_top), content_w, Inches(ex_h_in))
        ex_tf = ex_tb.text_frame
        ex_tf.word_wrap = True
        ex_tf.vertical_anchor = MSO_ANCHOR.TOP
        ex_p = ex_tf.paragraphs[0]
        ex_p.alignment = align
        hl2 = text_below_reference_highlights
        lt2 = text_below_reference_large_text
        if (hl2 and len(hl2) > 0) or (lt2 and len(lt2) > 0):
            _apply_highlights(ex_p, extra_ref, hl2, lt2,
                              body_rgb=theme['body'],
                              highlight_rgb=theme['highlight'],
                              annotation_rgb=theme['annotation'],
                              font_name=fn,
                              base_font_size=below_fs,
                              annotation_size_pt=ann_pt)
        else:
            ex_run = ex_p.add_run()
            ex_run.text = extra_ref
            ex_run.font.size = Pt(below_fs)
            ex_run.font.color.rgb = theme['body']
            if fn:
                ex_run.font.name = fn
    elif reference and ref_position == 'below':
        ref_text = reference + (f' (Part {part_num})' if part_num is not None else '')
        body_h_in = _estimate_verse_body_height_in(verse_text, content_w_in, font_size)
        if leading and leading_verse_top_in is not None:
            ref_y_in = leading_verse_top_in + body_h_in + body_gap_in
        else:
            ref_y_in = verse_top.inches + body_h_in + body_gap_in
        ref_lines = _estimate_text_lines(ref_text, content_w_in, ref_pt_bottom)
        ref_h_in = min(0.3 + ref_lines * 0.32, bottom_ref_h_in)
        ref_tb = verse_slide.shapes.add_textbox(
            left, Inches(ref_y_in), content_w, Inches(ref_h_in),
        )
        _set_ref(ref_tb, ref_text, align, ref_pt_bottom, italic=True)
    elif reference and ref_position == 'bottom':
        ref_text = reference + (f' (Part {part_num})' if part_num is not None else '')
        ref_tb = verse_slide.shapes.add_textbox(
            left, Inches(bottom_ref_top_in), content_w, Inches(bottom_ref_h_in),
        )
        _set_ref(ref_tb, ref_text, PP_ALIGN.CENTER, ref_pt_bottom, italic=True)

    return verse_slide


def create_presentation(data, output_file=None, custom_title=None, 
                        convert_to_pdf=False, pdf_options=None, pdf_backend='auto'):
    """
    Create a PowerPoint presentation from Bible verses data.
    
    Args:
        data (dict): Verses data dictionary with structure:
                     {
                         "presentation_title": "Title",
                         "presentation_subtitle": "Subtitle",
                         "sections": [
                             {
                                 "section": "Section Name",
                                 "verses": [
                                     {"reference": "Book 1:1", "text": "Verse text"}
                                 ]
                             }
                         ]
                     }
        output_file (str): Output filename (optional, auto-generated if not provided)
        custom_title (str): Custom presentation title (optional, overrides JSON title)
        convert_to_pdf (bool): Whether to also convert to PDF (default: False)
        pdf_options (PDFOptions): PDF conversion options (optional)
        pdf_backend (str): PDF conversion backend ('aspose', 'libreoffice', 'auto')
    
    Returns:
        str or dict: Path to the created presentation file, or dict with both PPTX and PDF paths
                    if convert_to_pdf is True, or None if error
    """
    if not data:
        print("Error: No data provided")
        return None
    
    # Create presentation
    prs = Presentation()

    # Apply slide size if specified in JSON
    # Supported: "widescreen"/"16:9", "standard"/"4:3", or {"width": W, "height": H} in inches
    slide_size = data.get("slide_size")
    if slide_size:
        from pptx.util import Inches as _Inches
        PRESETS = {
            "widescreen": (13.33, 7.5),
            "16:9":       (13.33, 7.5),
            "standard":   (10.0,  7.5),
            "4:3":        (10.0,  7.5),
            "16:10":      (12.8,  8.0),
        }
        if isinstance(slide_size, str) and slide_size.lower() in PRESETS:
            w, h = PRESETS[slide_size.lower()]
            prs.slide_width  = _Inches(w)
            prs.slide_height = _Inches(h)
        elif isinstance(slide_size, dict):
            if "width"  in slide_size: prs.slide_width  = _Inches(slide_size["width"])
            if "height" in slide_size: prs.slide_height = _Inches(slide_size["height"])

    # Get verses data from JSON
    verses_data = data.get("sections", [])
    # Presentation-level slide style (background, text color, alignment, etc.)
    slide_style = dict(data.get("slide_style", {}) or {})
    source_file = data.get("_source_file")
    if source_file:
        slide_style["_source_file"] = source_file

    # Add title slide
    if custom_title:
        title = custom_title
        subtitle = ""
    else:
        title = data.get("presentation_title", "Bible Verses Collection")
        subtitle = data.get("presentation_subtitle", "Selected Scriptures")
    
    add_title_slide(prs, title, subtitle, style=slide_style)
    
    # Add slides for each verse with section slides
    for section_data in verses_data:
        # Add section title slide if section name exists (skip if custom title is provided)
        if section_data.get("section") and not custom_title:
            add_section_slide(
                prs, section_data["section"], style=slide_style,
                section_subtitle=section_data.get("section_subtitle"),
            )

        # Add verse slides if there are any verses
        if section_data.get("verses") and len(section_data["verses"]) > 0:
            from .slide_renderers import resolve_renderer

            for v_idx, verse in enumerate(section_data["verses"]):
                if not isinstance(verse, dict):
                    continue
                resolve_renderer(verse).render(
                    prs, verse, slide_style, source_file=source_file,
                )
    
    # Generate output filename if not provided
    if not output_file:
        if custom_title:
            base_name = sanitize_filename(custom_title)
        else:
            base_name = sanitize_filename(data.get("presentation_title", "presentation"))
        output_file = f"{base_name}.pptx"
    
    # Ensure .pptx extension
    if not output_file.endswith('.pptx'):
        output_file += '.pptx'
    
    # Save presentation
    try:
        prs.save(output_file)
        print(f"✓ Presentation created successfully: {output_file}")
        
        # Convert to PDF if requested
        if convert_to_pdf:
            try:
                # Use default PDF options if none provided
                if pdf_options is None:
                    pdf_options = PDFOptions()
                
                # Generate PDF filename
                from pathlib import Path
                pdf_file = str(Path(output_file).with_suffix('.pdf'))
                
                # Convert to PDF
                pdf_result = convert_pptx_to_pdf(
                    output_file, 
                    pdf_file, 
                    backend=pdf_backend, 
                    options=pdf_options
                )
                
                print(f"✓ PDF created successfully: {pdf_result}")
                
                # Return both files
                return {
                    'pptx': output_file,
                    'pdf': pdf_result
                }
                
            except Exception as e:
                print(f"Warning: PDF conversion failed: {e}")
                print("Presentation was created successfully at:", output_file)
                return output_file
        
        return output_file
        
    except Exception as e:
        print(f"Error saving presentation: {e}")
        return None
