#!/usr/bin/env python3
"""
PowerPoint Generator for Bible Verses Collection
Creates a presentation with each verse on its own slide.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.dml import MSO_THEME_COLOR
import os
import json
import sys
import argparse

def split_long_text(text, max_length=200):
    """Split long text into multiple parts at sentence boundaries."""
    if len(text) <= max_length:
        return [text]
    
    # Split at sentences first
    sentences = text.replace('. ', '.|').replace('! ', '!|').replace('? ', '?|').split('|')
    parts = []
    current_part = ""
    
    for sentence in sentences:
        if len(current_part + sentence) <= max_length:
            current_part += sentence
        else:
            if current_part:
                parts.append(current_part.strip())
            current_part = sentence
    
    if current_part:
        parts.append(current_part.strip())
    
    return parts if parts else [text]

def load_verses_data(filename="verses.json"):
    """Load verses data from JSON file."""
    try:
        with open(filename, 'r', encoding='utf-8') as f:
            return json.load(f)
    except FileNotFoundError:
        print(f"Error: {filename} not found. Please create the file with your verses.")
        return None
    except json.JSONDecodeError:
        print(f"Error: Invalid JSON format in {filename}")
        return None

def create_bible_verses_presentation(verses_file="verses.json", output_file=None, custom_title=None):
    """Create PowerPoint presentation with Bible verses."""
    
    # Load verses data from JSON file
    data = load_verses_data(verses_file)
    if not data:
        return None
    
    # Create presentation
    prs = Presentation()
    
    # Get verses data from JSON
    verses_data = data.get("sections", [])
    
    # Title slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    # Use custom title if provided, otherwise use JSON title
    if custom_title:
        title.text = custom_title
        subtitle.text = ""  # No subtitle when using custom title
    else:
        title.text = data.get("presentation_title", "Bible Verses Collection")
        subtitle.text = data.get("presentation_subtitle", "Selected Scriptures")
    
    # Add slides for each verse (skip section slides when custom title is provided)
    for section_data in verses_data:
        # Only add section title slide if no custom title is provided
        if not custom_title:
            section_slide_layout = prs.slide_layouts[1]
            section_slide = prs.slides.add_slide(section_slide_layout)
            section_title = section_slide.shapes.title
            section_title.text = section_data["section"]
            
            # Style section title
            section_title.text_frame.paragraphs[0].font.size = Pt(36)
            section_title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
        
        # Add verse slides
        for verse in section_data["verses"]:
            # Split long verses into multiple parts
            verse_parts = split_long_text(verse["text"])
            
            for i, part in enumerate(verse_parts):
                verse_slide_layout = prs.slide_layouts[6]  # Blank layout to avoid bullets
                verse_slide = prs.slides.add_slide(verse_slide_layout)
                
                # Add verse reference as title using textbox
                left = Inches(1)
                top = Inches(1)
                width = Inches(8)
                height = Inches(1)
                
                title_box = verse_slide.shapes.add_textbox(left, top, width, height)
                title_frame = title_box.text_frame
                title_p = title_frame.paragraphs[0]
                
                # Add part indicator for multi-part verses
                if len(verse_parts) > 1:
                    title_p.text = f"{verse['reference']} (Part {i+1}/{len(verse_parts)})"
                else:
                    title_p.text = verse["reference"]
                    
                title_p.font.size = Pt(28)
                title_p.font.bold = True
                title_p.font.color.rgb = RGBColor(0, 51, 102)
                title_p.alignment = PP_ALIGN.CENTER
                
                # Add verse text using textbox (no bullets)
                left = Inches(1)
                top = Inches(2.5)
                width = Inches(8)
                height = Inches(4)
                
                text_box = verse_slide.shapes.add_textbox(left, top, width, height)
                text_frame = text_box.text_frame
                text_frame.word_wrap = True
                
                text_p = text_frame.paragraphs[0]
                text_p.text = f'"{part}"'
                text_p.font.size = Pt(40)
                text_p.font.color.rgb = RGBColor(51, 51, 51)
                text_p.alignment = PP_ALIGN.CENTER
    
    # Save presentation
    if output_file is None:
        if custom_title:
            # Generate output filename from custom title
            # Clean the title for use as filename
            safe_title = "".join(c for c in custom_title if c.isalnum() or c in (' ', '-', '_')).rstrip()
            safe_title = safe_title.replace(' ', '_')
            output_file = f"{safe_title}.pptx"
        else:
            # Generate output filename based on input file
            base_name = os.path.splitext(os.path.basename(verses_file))[0]
            output_file = f"{base_name}_presentation.pptx"
    
    prs.save(output_file)
    print(f"PowerPoint presentation saved as: {output_file}")
    return output_file

def main():
    """Main function to handle command line arguments."""
    parser = argparse.ArgumentParser(
        description="Create PowerPoint presentation from Bible verses JSON file",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
Examples:
  python create_bible_verses_presentation.py
  python create_bible_verses_presentation.py -v my_verses.json
  python create_bible_verses_presentation.py -v verses.json -t "Why Delay?"
  python create_bible_verses_presentation.py -v verses.json -t "God's Promises" -o promises.pptx
        """
    )
    
    parser.add_argument(
        '--verses', '-v',
        default='verses.json',
        help='Input JSON file containing verses (default: verses.json)'
    )
    
    parser.add_argument(
        '--title', '-t',
        help='Custom title for the presentation (overrides JSON title)'
    )
    
    parser.add_argument(
        '--output', '-o',
        help='Output PowerPoint file name (default: auto-generated from input filename)'
    )
    
    args = parser.parse_args()
    
    # Check if input file exists
    if not os.path.exists(args.verses):
        print(f"Error: Input file '{args.verses}' not found.")
        print("Please make sure the file exists and try again.")
        sys.exit(1)
    
    # Create presentation
    result = create_bible_verses_presentation(args.verses, args.output, args.title)
    
    if result:
        print(f"✅ Successfully created presentation: {result}")
    else:
        print("❌ Failed to create presentation. Please check your input file.")
        sys.exit(1)

if __name__ == "__main__":
    main()
