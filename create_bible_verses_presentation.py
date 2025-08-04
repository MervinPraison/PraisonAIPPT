#!/usr/bin/env python3
"""
PowerPoint Generator for Bible Verses Collection
Creates a presentation with each verse on its own slide.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

def split_long_text(text, max_length=200):
    """Split long text into multiple parts at sentence boundaries."""
    if len(text) <= max_length:
        return [text]
    
    # Split at sentences first
    sentences = text.replace('. ', '.|').replace('! ', '!|').replace('? ', '?|').split('|')
    parts = []
    current_part = ""
    
    for sentence in sentences:
        if len(current_part + sentence) <= max_length:
            current_part += sentence
        else:
            if current_part:
                parts.append(current_part.strip())
            current_part = sentence
    
    if current_part:
        parts.append(current_part.strip())
    
    return parts if parts else [text]

def create_bible_verses_presentation():
    """Create PowerPoint presentation with Bible verses."""
    
    # Create presentation
    prs = Presentation()
    
    # Define verses data
    verses_data = [
        {
            "section": "Bible Verses",
            "verses": [
                {
                    "reference": "Psalm 71:2",
                    "text": "In your righteousness, rescue me and deliver me; turn your ear to me and save me."
                },
                {
                    "reference": "Psalm 71:15",
                    "text": "My mouth will tell of your righteous deeds, of your saving acts all day long—though I know not how to relate them all."
                },
                {
                    "reference": "Psalm 71:16",
                    "text": "I will come and proclaim your mighty acts, Sovereign Lord; I will proclaim your righteous deeds, yours alone."
                },
                {
                    "reference": "Psalm 71:19",
                    "text": "Your righteousness, God, reaches to the heavens, you who have done great things. Who is like you, God?"
                },
                {
                    "reference": "Psalm 71:24",
                    "text": "My tongue will tell of your righteous acts all day long, for those who wanted to harm me have been put to shame and confusion."
                },
                {
                    "reference": "Proverbs 22:3",
                    "text": "The prudent see danger and take refuge, but the simple keep going and pay the penalty."
                },
                {
                    "reference": "Psalm 119:71",
                    "text": "It was good for me to be afflicted so that I might learn your decrees."
                },
                {
                    "reference": "Deuteronomy 28:1",
                    "text": "If you fully obey the Lord your God and carefully follow all his commands I give you today, the Lord your God will set you high above all the nations on earth."
                },
                {
                    "reference": "Deuteronomy 28:15",
                    "text": "However, if you do not obey the Lord your God and do not carefully follow all his commands and decrees I am giving you today, all these curses will come on you and overtake you:"
                },
                {
                    "reference": "Deuteronomy 28:68",
                    "text": "The Lord will send you back in ships to Egypt on a journey I said you should never make again. There you will offer yourselves for sale to your enemies as male and female slaves, but no one will buy you."
                },
                {
                    "reference": "Psalm 109:17",
                    "text": "He loved to pronounce a curse—may it come back on him. He found no pleasure in blessing—may it be far from him."
                },
                {
                    "reference": "Genesis 13:15",
                    "text": "All the land that you see I will give to you and your offspring forever."
                },
                {
                    "reference": "Joshua 1:3",
                    "text": "I will give you every place where you set your foot, as I promised Moses."
                },
                {
                    "reference": "Genesis 3:7",
                    "text": "Then the eyes of both of them were opened, and they realized they were naked; so they sewed fig leaves together and made coverings for themselves."
                },
                {
                    "reference": "Genesis 3:21",
                    "text": "The Lord God made garments of skin for Adam and his wife and clothed them."
                },
                {
                    "reference": "Genesis 4:3-4",
                    "text": "In the course of time Cain brought some of the fruits of the soil as an offering to the Lord. And Abel also brought an offering—fat portions from some of the firstborn of his flock. The Lord looked with favor on Abel and his offering,"
                },
                {
                    "reference": "Jeremiah 23:4",
                    "text": "I will place shepherds over them who will tend them, and they will no longer be afraid or terrified, nor will any be missing,\" declares the Lord."
                },
                {
                    "reference": "Jeremiah 23:6",
                    "text": "In his days Judah will be saved and Israel will live in safety. This is the name by which he will be called: The Lord Our Righteous Savior."
                },
                {
                    "reference": "Jeremiah 33:16",
                    "text": "In those days Judah will be saved and Jerusalem will live in safety. This is the name by which it will be called: The Lord Our Righteous Savior.'"
                },
                {
                    "reference": "Matthew 6:31",
                    "text": "So do not worry, saying, 'What shall we eat?' or 'What shall we drink?' or 'What shall we wear?'"
                },
                {
                    "reference": "Matthew 6:33",
                    "text": "But seek first his kingdom and his righteousness, and all these things will be given to you as well."
                }
            ]
        }
    ]
    
    # Title slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "Bible Verses Collection"
    subtitle.text = "Selected Scriptures"
    
    # Add slides for each verse
    for section_data in verses_data:
        # Add section title slide
        section_slide_layout = prs.slide_layouts[1]
        section_slide = prs.slides.add_slide(section_slide_layout)
        section_title = section_slide.shapes.title
        section_title.text = section_data["section"]
        
        # Style section title
        section_title.text_frame.paragraphs[0].font.size = Pt(36)
        section_title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
        
        # Add verse slides
        for verse in section_data["verses"]:
            # Split long verses into multiple parts
            verse_parts = split_long_text(verse["text"])
            
            for i, part in enumerate(verse_parts):
                verse_slide_layout = prs.slide_layouts[5]  # Blank layout
                verse_slide = prs.slides.add_slide(verse_slide_layout)
                
                # Add verse reference as title
                left = Inches(1)
                top = Inches(1)
                width = Inches(8)
                height = Inches(1)
                
                title_box = verse_slide.shapes.add_textbox(left, top, width, height)
                title_frame = title_box.text_frame
                title_p = title_frame.paragraphs[0]
                
                # Add part indicator for multi-part verses
                if len(verse_parts) > 1:
                    title_p.text = f"{verse['reference']} (Part {i+1}/{len(verse_parts)})"
                else:
                    title_p.text = verse["reference"]
                    
                title_p.font.size = Pt(28)
                title_p.font.bold = True
                title_p.font.color.rgb = RGBColor(0, 51, 102)
                title_p.alignment = PP_ALIGN.CENTER
                
                # Add verse text
                left = Inches(1)
                top = Inches(2.5)
                width = Inches(8)
                height = Inches(4)
                
                text_box = verse_slide.shapes.add_textbox(left, top, width, height)
                text_frame = text_box.text_frame
                text_frame.word_wrap = True
                
                text_p = text_frame.paragraphs[0]
                text_p.text = f'"{part}"'
                text_p.font.size = Pt(40)
                text_p.font.color.rgb = RGBColor(51, 51, 51)
                text_p.alignment = PP_ALIGN.CENTER
    
    # Save presentation
    output_file = "Bible_Verses_Collection.pptx"
    prs.save(output_file)
    print(f"PowerPoint presentation saved as: {output_file}")
    return output_file

if __name__ == "__main__":
    create_bible_verses_presentation()
