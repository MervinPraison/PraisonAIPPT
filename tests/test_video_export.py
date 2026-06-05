"""Tests for video export (manifest, timing, mocked ffmpeg)."""

from pathlib import Path
from unittest.mock import MagicMock, patch

import pytest
from pptx import Presentation
from pptx.util import Inches

from praisonaippt.avatar_layouts import export_slide_regions, region_box_to_pixels
from praisonaippt.ffmpeg_composer import (
    OverlaySpec,
    ToolCheck,
    check_video_tools,
    pick_video_encoder,
    render_slide_segment,
)
from praisonaippt.video_exporter import (
    VideoOptions,
    build_video_manifest,
    iter_slide_plan,
    resolve_slide_durations,
    resolve_video_backend,
    _resolve_narration_mode,
    _srt_time,
    write_srt,
    SlideVideoEntry,
)

PKG = Path(__file__).resolve().parent.parent


def test_video_options_preset():
    opts = VideoOptions(preset="draft")
    assert opts.width == 1280
    assert opts.height == 720
    assert opts.fps == 24


def test_video_options_from_dict():
    raw = {
        "preset": "standard",
        "narration_mode": "fixed",
        "resolution": {"width": 1280, "height": 720},
        "avatar": {"fit": "contain"},
    }
    opts = VideoOptions.from_dict(raw)
    assert opts.width == 1280
    assert opts.avatar_fit == "contain"


def test_iter_slide_plan_avatar_deck():
    data = {
        "presentation_title": "T",
        "sections": [
            {
                "section": "S",
                "verses": [
                    {"slide_type": "avatar_media_1", "media_path": "x.jpg"},
                    {"slide_type": "avatar_media_2", "media_path": "x.jpg"},
                ],
            }
        ],
    }
    plan = list(iter_slide_plan(data))
    assert plan[0]["slide_role"] == "title"
    assert plan[1]["slide_role"] == "section"
    assert len(plan) == 4


def test_manifest_avatar_region_ratios():
    data = {
        "sections": [{"verses": [
            {"slide_type": "avatar_media_1", "media_path": "a.jpg"},
            {"slide_type": "avatar_media_2", "media_path": "a.jpg"},
        ]}],
    }
    prs = Presentation()
    blank = prs.slide_layouts[6]
    for _ in range(3):
        prs.slides.add_slide(blank)
    opts = VideoOptions()
    entries = build_video_manifest(data, prs, opts)
    avatar_entries = [e for e in entries if e.slide_type and e.slide_type.startswith("avatar_media")]
    assert len(avatar_entries) == 2
    w1 = avatar_entries[0].media_box_px["width"]
    w2 = avatar_entries[1].media_box_px["width"]
    assert w1 != w2


def test_manifest_avatar_pip_on_standard_slide():
    data = {
        "presentation_title": "T",
        "sections": [{"verses": [
            {
                "slide_type": "list",
                "text": "One\nTwo",
                "reference": "Demo",
                "avatar_video_path": "heygen.mp4",
            },
        ]}],
    }
    prs = Presentation()
    prs.slides.add_slide(prs.slide_layouts[6])
    prs.slides.add_slide(prs.slide_layouts[6])
    entries = build_video_manifest(data, prs, VideoOptions())
    content = [e for e in entries if e.slide_role == "content"][0]
    assert content.slide_type == "list"
    assert content.avatar_box_px is not None
    assert content.avatar_box_px["width"] > 0


def test_resolve_narration_mode_auto_fixed():
    assert _resolve_narration_mode(None, "fixed") == "fixed"
    assert _resolve_narration_mode({"audio_path": "a.mp3"}, "auto") == "audio_file"


@patch("praisonaippt.video_exporter.ffprobe_has_audio_safe", return_value=True)
def test_resolve_narration_mode_auto_prefers_heygen_video(mock_has_audio):
    verse = {
        "audio_path": "n.mp3",
        "avatar_video_path": "heygen.mp4",
    }
    assert _resolve_narration_mode(verse, "auto") == "avatar"
    mock_has_audio.assert_called()


def test_video_options_audio_source_maps_to_narration_mode():
    opts = VideoOptions.from_dict({"audio_source": "heygen_video"})
    assert opts.narration_mode == "avatar"
    opts2 = VideoOptions.from_dict({"audio_source": "external"})
    assert opts2.narration_mode == "audio_file"
    opts3 = VideoOptions.from_dict(
        {"audio_source": "external", "narration_mode": "avatar"},
    )
    assert opts3.narration_mode == "avatar"


def test_resolve_slide_durations_fixed(tmp_path):
    entries = [
        SlideVideoEntry(index=0, slide_role="content", slide_type="avatar_only", verse={}),
    ]
    opts = VideoOptions(slide_duration_sec=7.0, narration_mode="fixed")
    resolve_slide_durations(entries, opts, temp_dir=tmp_path)
    assert entries[0].duration_sec == 7.0
    assert entries[0].audio_primary == "none"


@patch("praisonaippt.video_exporter.ffprobe_duration", return_value=10.5)
def test_resolve_slide_durations_avatar(mock_dur, tmp_path):
    entries = [
        SlideVideoEntry(
            index=0,
            slide_role="content",
            slide_type="avatar_only",
            verse={"avatar_video_path": "speaker.mp4"},
            narration_mode="avatar",
            avatar_video_path="speaker.mp4",
        ),
    ]
    opts = VideoOptions()
    with patch("praisonaippt.video_exporter.resolve_asset_path", return_value="speaker.mp4"):
        with patch("pathlib.Path.is_file", return_value=True):
            resolve_slide_durations(entries, opts, temp_dir=tmp_path)
    assert entries[0].duration_sec == 10.5
    assert entries[0].audio_primary == "avatar"


def test_srt_time_format():
    assert _srt_time(0) == "00:00:00,000"
    assert _srt_time(65.5) == "00:01:05,500"


def test_write_srt(tmp_path):
    entries = [
        SlideVideoEntry(
            index=0, slide_role="content", slide_type="verse", verse=None,
            duration_sec=2.0, caption_text="Hello",
        ),
    ]
    out = tmp_path / "test.srt"
    write_srt(entries, str(out))
    text = out.read_text(encoding="utf-8")
    assert "Hello" in text
    assert "-->" in text


@patch("praisonaippt.ffmpeg_composer._run")
def test_render_slide_segment_silent_audio(mock_run):
    mock_run.return_value = MagicMock(returncode=0)
    render_slide_segment("slide.png", 3.0, "out.mp4")
    cmd = mock_run.call_args[0][0]
    assert "anullsrc" in " ".join(cmd)
    assert "-an" not in cmd


@patch("praisonaippt.ffmpeg_composer._run")
def test_render_slide_segment_calls_ffmpeg(mock_run):
    mock_run.return_value = MagicMock(returncode=0)
    render_slide_segment(
        "slide.png", 3.0, "out.mp4",
        overlays=[OverlaySpec("v.mp4", 10, 10, 100, 100, is_video=True)],
    )
    assert mock_run.called
    cmd = mock_run.call_args[0][0]
    assert cmd[0] == "ffmpeg"
    assert "-filter_complex" in cmd


@patch("praisonaippt.ffmpeg_composer.probe_tool")
def test_check_video_tools(mock_probe):
    from praisonaippt.pdf_converter import PDFConverter

    mock_probe.side_effect = lambda name, **kw: ToolCheck(name=name, found=True, path=f"/usr/{name}")

    def _fake_detect(self):
        self._libreoffice_path = "/Applications/LibreOffice.app/Contents/MacOS/soffice"
        return ["libreoffice"]

    with patch.object(PDFConverter, "_detect_backends", _fake_detect):
        tools = check_video_tools()
    assert tools["ffmpeg"].found
    assert tools["libreoffice"].found


def test_region_box_to_pixels():
    prs = Presentation()
    box = export_slide_regions(prs, "avatar_media_1", {})["media"]
    px = region_box_to_pixels(box, prs.slide_width.inches, prs.slide_height.inches, 1920, 1080)
    assert px["width"] > 0
    assert px["x"] >= 0


def test_parse_video_options_respects_yaml_narration_mode():
    from argparse import Namespace
    from praisonaippt.cli import parse_video_options

    data = {"video_export": {"narration_mode": "auto", "preset": "draft"}}
    args = Namespace(
        video_options=None,
        video_backend=None,
        video_preset=None,
        narration_mode=None,
        video_output=None,
        keep_temp=False,
        slide_range=None,
    )
    opts = parse_video_options(args, data)
    assert opts.narration_mode == "auto"
    assert opts.preset == "draft"
    assert opts.width == 1280


def test_windows_worker_not_implemented():
    from praisonaippt.workers.ppt_com import create_video_via_powerpoint

    with pytest.raises(NotImplementedError):
        create_video_via_powerpoint("a.pptx", "b.mp4")


def test_video_options_invalid_narration_mode():
    with pytest.raises(ValueError, match="narration_mode"):
        VideoOptions(narration_mode="invalid")


def test_resolve_video_backend_auto():
    assert resolve_video_backend(VideoOptions(backend="auto")) == "compositor"


def test_resolve_video_backend_aspose_raises():
    with pytest.raises(NotImplementedError):
        resolve_video_backend(VideoOptions(backend="aspose_frames"))


def test_manifest_plan_mismatch_raises():
    data = {"sections": [{"verses": [{"slide_type": "avatar_only"}]}]}
    prs = Presentation()
    with pytest.raises(RuntimeError, match="slide plan"):
        build_video_manifest(data, prs, VideoOptions())


@patch("praisonaippt.video_exporter.render_slide_segment")
@patch("praisonaippt.video_exporter.concat_segments")
def test_compose_video_slide_range(mock_concat, mock_render, tmp_path):
    from praisonaippt.video_exporter import compose_video

    entries = [
        SlideVideoEntry(index=0, slide_role="title", slide_type="title", verse=None, duration_sec=1.0),
        SlideVideoEntry(index=1, slide_role="content", slide_type="avatar_only", verse=None, duration_sec=2.0),
    ]
    opts = VideoOptions(slide_range=(2, 2))
    compose_video(entries, ["a.png", "b.png"], "out.mp4", opts, tmp_path)
    assert mock_render.call_count == 1
    mock_concat.assert_called_once()


@patch("praisonaippt.video_exporter.ffprobe_duration", return_value=12.0)
def test_apply_sync_mode_longest(mock_dur, tmp_path):
    from praisonaippt.video_exporter import _apply_sync_mode

    entry = SlideVideoEntry(
        index=0, slide_role="content", slide_type="avatar_only", verse={},
        duration_sec=5.0, avatar_video_path="a.mp4",
    )
    with patch("pathlib.Path.is_file", return_value=True):
        _apply_sync_mode(entry, {"sync_mode": "longest"}, VideoOptions(), tmp_path, None)
    assert entry.duration_sec == 12.0


def test_write_srt_contiguous_indices(tmp_path):
    entries = [
        SlideVideoEntry(index=0, slide_role="c", slide_type="v", verse=None, duration_sec=1.0, caption_text=""),
        SlideVideoEntry(index=1, slide_role="c", slide_type="v", verse=None, duration_sec=2.0, caption_text="Hi"),
    ]
    out = tmp_path / "t.srt"
    write_srt(entries, str(out))
    assert out.read_text(encoding="utf-8").startswith("1\n")


def test_load_deck_sidecar_yaml(tmp_path):
    from praisonaippt.video_sidecar import load_deck_sidecar

    yaml_path = tmp_path / "deck.yaml"
    yaml_path.write_text(
        "presentation_title: T\nsections:\n  - verses:\n      - slide_type: avatar_only\n",
        encoding="utf-8",
    )
    pptx = tmp_path / "deck.pptx"
    pptx.write_bytes(b"pk")
    data = load_deck_sidecar(str(pptx))
    assert data is not None
    assert data["presentation_title"] == "T"


def test_letterbox_region_pixels():
    from praisonaippt.avatar_layouts import letterbox_content_rect, region_box_to_pixels

    prs = Presentation()
    box = export_slide_regions(prs, "avatar_media_1", {})["media"]
    plain = region_box_to_pixels(
        box, prs.slide_width.inches, prs.slide_height.inches, 1920, 1080,
    )
    pad_x, pad_y, _, _ = letterbox_content_rect(
        prs.slide_width.inches, prs.slide_height.inches, 1920, 1080,
    )
    assert plain["x"] >= pad_x or pad_x == 0


@patch("praisonaippt.ffmpeg_composer._run")
def test_concat_segments_single(mock_run, tmp_path):
    from praisonaippt.ffmpeg_composer import concat_segments

    src = tmp_path / "one.mp4"
    src.write_bytes(b"fake")
    out = tmp_path / "out.mp4"
    concat_segments([str(src)], str(out))
    assert out.is_file()
    mock_run.assert_not_called()


@patch("praisonaippt.video_exporter.ffprobe_duration", return_value=57.0)
def test_title_slide_silent_in_audio_file(mock_dur, tmp_path):
    entries = [
        SlideVideoEntry(index=0, slide_role="title", slide_type="title", verse=None, narration_mode="audio_file"),
        SlideVideoEntry(
            index=1, slide_role="content", slide_type="avatar_only",
            verse={"audio_path": "n.mp3", "duration_sec": 4.0, "audio_start_sec": 0.0},
            narration_mode="audio_file", duration_sec=4.0, audio_path="n.mp3",
        ),
    ]
    opts = VideoOptions(slide_duration_sec=3.0)
    opts._slide_timestamps = [0.0, 3.0, 7.0]  # type: ignore[attr-defined]
    with patch("praisonaippt.video_exporter.resolve_asset_path", return_value="n.mp3"):
        with patch("pathlib.Path.is_file", return_value=True):
            resolve_slide_durations(entries, opts, temp_dir=tmp_path)
    assert entries[0].duration_sec == 3.0
    assert entries[0].audio_primary == "none"
    assert entries[1].duration_sec == 4.0
    assert entries[1].audio_start_sec == 0.0


@patch("praisonaippt.video_exporter.ffprobe_duration", return_value=57.0)
def test_duration_preserved_avatar_mode(mock_dur, tmp_path):
    entries = [
        SlideVideoEntry(
            index=0,
            slide_role="content",
            slide_type="avatar_only",
            verse={"avatar_video_path": "heygen.mp4", "duration_sec": 4.0, "audio_start_sec": 0.0},
            narration_mode="avatar",
            duration_sec=4.0,
            avatar_video_path="heygen.mp4",
        ),
    ]
    opts = VideoOptions()
    with patch("praisonaippt.video_exporter.resolve_asset_path", return_value="heygen.mp4"):
        with patch("pathlib.Path.is_file", return_value=True):
            resolve_slide_durations(entries, opts, temp_dir=tmp_path)
    assert entries[0].duration_sec == 4.0
    mock_dur.assert_not_called()


@patch("praisonaippt.video_exporter.ffprobe_duration", return_value=57.0)
def test_duration_preserved_audio_file(mock_dur, tmp_path):
    entries = [
        SlideVideoEntry(
            index=0,
            slide_role="content",
            slide_type="avatar_only",
            verse={"audio_path": "n.mp3", "duration_sec": 8.1, "audio_start_sec": 18.2},
            narration_mode="audio_file",
            duration_sec=8.1,
            audio_path="n.mp3",
            audio_start_sec=18.2,
        ),
    ]
    opts = VideoOptions()
    with patch("praisonaippt.video_exporter.resolve_asset_path", return_value="n.mp3"):
        with patch("pathlib.Path.is_file", return_value=True):
            resolve_slide_durations(entries, opts, temp_dir=tmp_path)
    assert entries[0].duration_sec == 8.1
    assert entries[0].audio_start_sec == 18.2


@patch("praisonaippt.video_exporter.ffprobe_duration", return_value=57.0)
def test_slide_timestamps_with_audio_path(mock_dur, tmp_path):
    entries = [
        SlideVideoEntry(
            index=0,
            slide_role="content",
            slide_type="avatar_only",
            verse={"audio_path": "n.mp3"},
            narration_mode="audio_file",
            audio_path="n.mp3",
        ),
    ]
    opts = VideoOptions()
    opts._slide_timestamps = [0.0, 4.0, 12.0]  # type: ignore[attr-defined]
    with patch("praisonaippt.video_exporter.resolve_asset_path", return_value="n.mp3"):
        with patch("pathlib.Path.is_file", return_value=True):
            resolve_slide_durations(entries, opts, temp_dir=tmp_path)
    assert entries[0].duration_sec == 4.0
    assert entries[0].audio_start_sec == 0.0


@patch("praisonaippt.video_exporter.ffprobe_duration", return_value=57.0)
def test_sync_mode_skipped_when_explicit_duration(mock_dur, tmp_path):
    from praisonaippt.video_exporter import _apply_sync_mode

    entry = SlideVideoEntry(
        index=0, slide_role="content", slide_type="avatar_only", verse={},
        duration_sec=5.0, avatar_video_path="a.mp4",
    )
    with patch("pathlib.Path.is_file", return_value=True):
        _apply_sync_mode(
            entry,
            {"sync_mode": "longest", "duration_sec": 5.0},
            VideoOptions(),
            tmp_path,
            None,
        )
    assert entry.duration_sec == 5.0


@patch("praisonaippt.ffmpeg_composer._run")
def test_render_slide_segment_audio_start(mock_run):
    mock_run.return_value = MagicMock(returncode=0)
    render_slide_segment("slide.png", 3.0, "out.mp4", audio_path="a.mp3", audio_start_sec=12.5)
    cmd = mock_run.call_args[0][0]
    joined = " ".join(cmd)
    assert "-ss" in joined
    assert "12.500" in joined


@patch("praisonaippt.video_exporter.render_slide_segment")
@patch("praisonaippt.video_exporter.concat_segments")
def test_compose_passes_audio_start(mock_concat, mock_render, tmp_path):
    from praisonaippt.video_exporter import compose_video

    entries = [
        SlideVideoEntry(
            index=0, slide_role="content", slide_type="avatar_only", verse=None,
            duration_sec=4.0, narration_mode="audio_file", audio_primary="file",
            audio_path="n.mp3", audio_start_sec=18.2,
        ),
    ]
    compose_video(entries, ["a.png"], "out.mp4", VideoOptions(), tmp_path)
    kwargs = mock_render.call_args[1]
    assert kwargs.get("audio_start_sec") == 18.2


@pytest.mark.parametrize("slide_type,expect_avatar,expect_rect,skip_media", [
    ("deck_title_split", True, True, True),
    ("deck_exec_summary", True, False, True),
    ("deck_agenda", False, False, True),
    ("deck_intro_split", False, False, True),
    ("deck_thank_you", True, True, True),
])
def test_manifest_deck_layout_regions(slide_type, expect_avatar, expect_rect, skip_media):
    from praisonaippt.deck_slides import DECK_SLIDE_TYPES

    assert slide_type in DECK_SLIDE_TYPES
    verse = {"slide_type": slide_type, "text": "Title"}
    if slide_type == "deck_exec_summary":
        verse["items"] = [{"text": "one"}]
    elif slide_type == "deck_agenda":
        verse["items"] = ["One"]
    elif slide_type == "deck_intro_split":
        verse["reference"] = "Body"
        verse["media_path"] = "assets/background_alt.jpg"
    elif slide_type == "deck_thank_you":
        verse["reference"] = "YOU"
    if expect_avatar:
        verse["avatar_video_path"] = "examples/heygen-article-50590.mp4"

    data = {"sections": [{"verses": [verse]}]}
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    prs.slides.add_slide(prs.slide_layouts[6])
    prs.slides.add_slide(prs.slide_layouts[6])
    data["_source_file"] = str(PKG)
    entries = build_video_manifest(data, prs, VideoOptions(), source_file=str(PKG))
    content = [e for e in entries if e.slide_role == "content"][0]
    assert content.slide_type == slide_type
    if expect_avatar:
        assert content.avatar_box_px is not None
        assert content.avatar_box_px["width"] > 0
        if expect_rect:
            assert content.avatar_shape == "rect"
        else:
            assert content.avatar_shape == "circle"
    else:
        assert content.avatar_box_px is None
    if skip_media:
        assert content.skip_media_overlay is True
        assert content.media_box_px is None


def test_avatar_overlay_uses_audio_start_sec():
    from praisonaippt.ffmpeg_composer import OverlaySpec
    from praisonaippt.video_exporter import SlideVideoEntry, VideoOptions, _apply_avatar_overlay_timing

    entry = SlideVideoEntry(
        index=2,
        slide_role="content",
        slide_type="deck_exec_summary",
        verse={"audio_start_sec": 21.2, "avatar_video_path": "heygen.mp4"},
        narration_mode="avatar",
        duration_sec=8.0,
        avatar_video_path="heygen.mp4",
    )
    ov = OverlaySpec(path="heygen.mp4", x=0, y=0, width=100, height=100, is_video=True)
    _apply_avatar_overlay_timing(
        [ov], entry, VideoOptions(), avatar_offset=99.0, timeline="continuous",
    )
    assert ov.video_start_sec == 21.2


@patch("praisonaippt.video_exporter.ffprobe_duration", return_value=57.0)
def test_timestamps_do_not_set_audio_seek(mock_dur, tmp_path):
    entries = [
        SlideVideoEntry(index=0, slide_role="title", slide_type="title", verse=None, narration_mode="avatar"),
        SlideVideoEntry(
            index=1,
            slide_role="content",
            slide_type="avatar_only",
            verse={"avatar_video_path": "heygen.mp4"},
            narration_mode="avatar",
        ),
    ]
    opts = VideoOptions(slide_duration_sec=3.0)
    opts._slide_timestamps = [0.0, 3.0, 10.0]  # type: ignore[attr-defined]
    with patch("praisonaippt.video_exporter.resolve_asset_path", return_value="heygen.mp4"):
        with patch("pathlib.Path.is_file", return_value=True):
            resolve_slide_durations(entries, opts, temp_dir=tmp_path)
    assert entries[1].audio_start_sec == 0.0


def test_overlays_skip_baked_deck_avatar():
    from praisonaippt.video_exporter import _overlays_for_entry

    entry = SlideVideoEntry(
        index=0,
        slide_role="content",
        slide_type="deck_title_split",
        verse={"avatar_video_path": "examples/heygen-article-50590.mp4"},
        avatar_video_path="examples/heygen-article-50590.mp4",
        avatar_box_px={"x": 960, "y": 0, "width": 960, "height": 1080},
        skip_avatar_overlay=True,
    )
    overlays = _overlays_for_entry(entry, VideoOptions(), source_file=str(PKG))
    assert overlays == []


def test_overlays_play_on_deck_thank_you():
    from praisonaippt.video_exporter import _overlays_for_entry

    entry = SlideVideoEntry(
        index=7,
        slide_role="content",
        slide_type="deck_thank_you",
        verse={
            "avatar_video_path": "examples/heygen-article-50590.mp4",
            "audio_start_sec": 51.88,
        },
        avatar_video_path="examples/heygen-article-50590.mp4",
        avatar_box_px={"x": 960, "y": 0, "width": 960, "height": 1080},
        skip_avatar_overlay=False,
        avatar_shape="h_rect",
    )
    with patch("pathlib.Path.is_file", return_value=True):
        overlays = _overlays_for_entry(entry, VideoOptions(), source_file=str(PKG))
    assert len(overlays) == 1
    assert overlays[0].is_video
    assert overlays[0].shape == "h_rect"


def test_overlays_skip_baked_deck_media():
    from praisonaippt.video_exporter import _overlays_for_entry

    entry = SlideVideoEntry(
        index=0,
        slide_role="content",
        slide_type="deck_intro_split",
        verse={"media_path": "assets/background_alt.jpg"},
        media_path="assets/background_alt.jpg",
        media_box_px={"x": 0, "y": 100, "width": 640, "height": 360},
        skip_media_overlay=True,
    )
    overlays = _overlays_for_entry(entry, VideoOptions(), source_file=str(PKG))
    assert overlays == []


def test_video_options_transitions_default_none():
    opts = VideoOptions.from_dict({"preset": "standard"}, {})
    assert opts.transition_defaults.default == "none"


def test_build_compose_plan_default_none_edges():
    from praisonaippt.video_exporter import build_compose_plan

    entries = [
        SlideVideoEntry(0, "title", "title", None, duration_sec=3.0),
        SlideVideoEntry(1, "content", "verse", {}, duration_sec=5.0),
    ]
    _, edges = build_compose_plan(entries, {}, VideoOptions())
    assert all(e.type == "none" for e in edges)


def test_write_srt_with_xfade_timeline(tmp_path):
    entries = [
        SlideVideoEntry(0, "content", "verse", {}, duration_sec=10.0, caption_text="A"),
        SlideVideoEntry(1, "content", "verse", {}, duration_sec=10.0, caption_text="B"),
    ]
    from praisonaippt.video_protocol import ResolvedEdgeTransition

    edges = [ResolvedEdgeTransition(1, "crossfade", 2.0, "test")]
    out = tmp_path / "cap.srt"
    write_srt(entries, str(out), edges=edges)
    text = out.read_text(encoding="utf-8")
    assert "00:00:00,000 --> 00:00:10,000" in text
    assert "00:00:08,000 --> 00:00:18,000" in text
