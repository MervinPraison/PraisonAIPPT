"""Tests for video overlay protocol parsing and manifest integration."""

from pathlib import Path

import pytest
from pptx import Presentation

from praisonaippt.exceptions import SchemaError
from praisonaippt.video_protocol import (
    OverlayPlacement,
    apply_pixel_offset,
    merge_placement,
    parse_placement,
    region_from_placement,
    resolve_slide_overlays,
    validate_overlay_placement,
)
from praisonaippt.video_exporter import VideoOptions, build_video_manifest

PKG = Path(__file__).resolve().parent.parent / "examples"


def test_parse_placement_explicit_box():
    p = parse_placement({
        "box": {"left_in": 1.0, "top_in": 2.0, "width_in": 3.0, "height_in": 4.0},
        "zoom_ratio": 1.2,
        "offset_px": {"x": 10, "y": -5},
    })
    assert p.has_explicit_box()
    assert p.zoom_ratio == 1.2
    assert p.offset_px == (10, -5)


def test_merge_placement_precedence():
    base = parse_placement({"zoom_ratio": 1.1, "anchor": "bottom_right"})
    override = parse_placement({"zoom_ratio": 1.35, "anchor": "top_right"})
    merged = merge_placement(base, override)
    assert merged.zoom_ratio == 1.35
    assert merged.anchor == "top_right"


def test_deck_layout_pip_width_beats_global_pip_ratio():
    """``layouts.deck_*``.pip_width_ratio must not be overridden by ``layouts.pip``.width_ratio."""
    resolved = resolve_slide_overlays(
        verse={},
        slide_type="deck_exec_summary",
        style={
            "layouts": {
                "pip": {"width_ratio": 0.24, "margin_in": 0.38},
                "deck_exec_summary": {"pip_position": "top_right", "pip_width_ratio": 0.20},
            },
        },
        video_export={},
        framing_kind="deck_exec_summary",
    )
    assert resolved.avatar.width_ratio == 0.20


def test_resolve_slide_overlays_verse_wins():
    verse = {
        "video_overlay": {
            "avatar": {"zoom_ratio": 1.5, "anchor": "top_right", "width_ratio": 0.22},
        },
        "avatar_zoom_ratio": 1.4,
    }
    resolved = resolve_slide_overlays(
        verse=verse,
        slide_type="avatar_quote",
        style={"layouts": {"pip": {"zoom_ratio": 1.1}}},
        video_export={"avatar": {"zoom_ratio": 1.2}},
        framing_kind="avatar_quote",
    )
    assert resolved.avatar.zoom_ratio == 1.5
    assert resolved.avatar.anchor == "top_right"


def test_region_from_anchor_overrides_layout():
    from pptx.util import Inches

    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    placement = OverlayPlacement(anchor="top_right", width_ratio=0.25, margin_in=0.4)
    box = region_from_placement(None, placement, 13.33, 7.5, {}, "pip")
    assert box is not None
    assert box.left_in > 8.0
    assert box.top_in < 1.5


def test_apply_pixel_offset():
    px = {"x": 100, "y": 200, "width": 50, "height": 50}
    out = apply_pixel_offset(px, (10, 20))
    assert out["x"] == 110
    assert out["y"] == 220


def test_validate_overlay_placement_rejects_bad_anchor():
    with pytest.raises(SchemaError):
        validate_overlay_placement({"anchor": "middle"}, "video_export.avatar")


def test_manifest_applies_verse_video_overlay():
    data = {
        "sections": [{
            "verses": [{
                "slide_type": "avatar_quote",
                "text": "Hello",
                "avatar_video_path": "examples/heygen-article-50590.mp4",
                "video_overlay": {
                    "avatar": {
                        "anchor": "bottom_right",
                        "width_ratio": 0.22,
                        "margin_in": 0.35,
                        "zoom_ratio": 1.25,
                    },
                },
            }],
        }],
        "video_export": {"narration_mode": "fixed", "preset": "draft"},
    }
    prs = Presentation()
    prs.slides.add_slide(prs.slide_layouts[6])
    prs.slides.add_slide(prs.slide_layouts[6])
    entries = build_video_manifest(
        data, prs, VideoOptions(narration_mode="fixed", preset="draft"),
        source_file=str(PKG / "heygen-50590-content.yaml"),
    )
    content = [e for e in entries if e.slide_role == "content"][0]
    assert content.avatar_zoom_ratio == 1.25
    assert content.avatar_box_px is not None
    assert content.avatar_box_px["width"] > 0


def test_heygen_content_passes_overlay_validation():
    from praisonaippt.loader import load_verses_from_file

    data = load_verses_from_file(str(PKG / "heygen-50590-content.yaml"))
    assert data.get("video_export")
