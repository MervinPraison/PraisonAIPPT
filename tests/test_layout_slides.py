"""Tests for standard layout slide types (title_only, two_column, etc.)."""

from pathlib import Path

import pytest
from pptx import Presentation

from praisonaippt import create_presentation, load_verses_from_dict, list_renderers
from praisonaippt.exceptions import SchemaError
from praisonaippt.slide_renderers import resolve_renderer, validate_verse

PKG = Path(__file__).resolve().parent.parent
IMG = PKG / "assets" / "background_alt.jpg"


def test_list_renderers_includes_new_layout_kinds():
    kinds = set(list_renderers())
    assert kinds >= {
        "title_only", "two_column", "comparison", "big_number",
        "quote", "picture_text", "table", "verse", "image",
    }


@pytest.mark.parametrize("slide_type", [
    "title_only", "two_column", "comparison", "big_number",
    "quote", "table",
])
def test_resolve_renderer_new_types(slide_type):
    verse = {"slide_type": slide_type}
    if slide_type == "title_only":
        verse["text"] = "Headline"
    elif slide_type in ("two_column",):
        verse["left"] = "A"
        verse["right"] = "B"
    elif slide_type == "comparison":
        verse["columns"] = [{"heading": "H1", "text": "A"}, {"heading": "H2", "text": "B"}]
    elif slide_type == "big_number":
        verse["number"] = "100"
    elif slide_type == "quote":
        verse["text"] = "Quote text"
    elif slide_type == "table":
        verse["table_rows"] = [["Col1", "Col2"], ["a", "b"]]
    assert resolve_renderer(verse).kind == slide_type


def test_title_only_builds(tmp_path):
    data = {
        "presentation_title": "T",
        "sections": [{"section": "", "verses": [
            {"slide_type": "title_only", "text": "Main headline", "reference": "Subtitle line"},
        ]}],
    }
    out = tmp_path / "title_only.pptx"
    create_presentation(load_verses_from_dict(data), str(out))
    assert out.is_file()
    assert len(Presentation(out).slides) == 2


def test_comparison_and_big_number_build(tmp_path):
    data = {
        "presentation_title": "T",
        "sections": [{"section": "", "verses": [
            {"slide_type": "comparison", "columns": [
                {"heading": "Before", "text": "Old way"},
                {"heading": "After", "text": "New way"},
            ]},
            {"slide_type": "big_number", "number": "100", "label": "Fold blessing"},
            {"slide_type": "quote", "text": "Faith comes by hearing.", "reference": "Romans 10:17"},
        ]}],
    }
    out = tmp_path / "layouts.pptx"
    create_presentation(load_verses_from_dict(data), str(out))
    assert len(Presentation(out).slides) == 4


@pytest.mark.skipif(not IMG.is_file(), reason="sample image missing")
def test_picture_text_build(tmp_path):
    data = {
        "presentation_title": "T",
        "sections": [{"section": "", "verses": [{
            "slide_type": "picture_text",
            "image_path": str(IMG),
            "image_side": "left",
            "text": "Explanation beside the image.",
        }]}],
    }
    out = tmp_path / "pic_text.pptx"
    create_presentation(load_verses_from_dict(data), str(out))
    assert len(Presentation(out).slides) == 2


def test_speaker_notes_on_verse(tmp_path):
    data = {
        "presentation_title": "T",
        "sections": [{"section": "", "verses": [{
            "reference": "John 3:16",
            "text": "For God so loved the world.",
            "notes": "Mention Nicodemus context.",
        }]}],
    }
    out = tmp_path / "notes.pptx"
    create_presentation(load_verses_from_dict(data), str(out))
    slide = Presentation(out).slides[1]
    assert "Nicodemus" in slide.notes_slide.notes_text_frame.text


def test_speaker_notes_on_title_only(tmp_path):
    data = {
        "presentation_title": "T",
        "sections": [{"section": "", "verses": [{
            "slide_type": "title_only",
            "text": "Headline",
            "notes": "Presenter cue for this slide.",
        }]}],
    }
    out = tmp_path / "notes_title.pptx"
    create_presentation(load_verses_from_dict(data), str(out))
    slide = Presentation(out).slides[1]
    assert "Presenter cue" in slide.notes_slide.notes_text_frame.text


def test_validate_comparison_requires_columns():
    with pytest.raises(SchemaError):
        validate_verse({"slide_type": "comparison"}, "v[0]")
