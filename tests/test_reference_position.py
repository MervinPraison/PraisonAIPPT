"""Tests for verse reference placement (bottom / below / top)."""

from pptx import Presentation

from praisonaippt import create_presentation, load_verses_from_dict
from praisonaippt.core import _normalize_ref_position, _resolve_theme


def _ref_top_inches(pptx_path, slide_index=1):
    prs = Presentation(pptx_path)
    slide = prs.slides[slide_index]
    refs = [
        s.top.inches
        for s in slide.shapes
        if s.has_text_frame and s.text_frame.text and "Romans" in s.text_frame.text
    ]
    return min(refs) if refs else None


def test_default_reference_position_is_bottom():
    theme = _resolve_theme({})
    assert theme["ref_position"] == "bottom"
    assert _normalize_ref_position(None) == "bottom"


def test_below_places_reference_higher_than_bottom(tmp_path):
    text = "For I am not ashamed of the gospel of Christ."
    base = {
        "presentation_title": "Ref",
        "sections": [
            {
                "section": "",
                "verses": [
                    {
                        "reference": "Romans 1:16 (NKJV)",
                        "text": text,
                    }
                ],
            }
        ],
    }
    bottom_out = tmp_path / "bottom.pptx"
    below_out = tmp_path / "below.pptx"
    create_presentation(
        load_verses_from_dict({**base, "slide_style": {"reference_position": "bottom"}}),
        str(bottom_out),
    )
    create_presentation(
        load_verses_from_dict({**base, "slide_style": {"reference_position": "below"}}),
        str(below_out),
    )
    bottom_y = _ref_top_inches(bottom_out)
    below_y = _ref_top_inches(below_out)
    assert bottom_y is not None and below_y is not None
    assert below_y < bottom_y - 0.5


def test_top_places_reference_above_verse_body(tmp_path):
    out = tmp_path / "top.pptx"
    data = {
        "presentation_title": "Ref",
        "slide_style": {"reference_position": "top"},
        "sections": [
            {
                "section": "",
                "verses": [
                    {
                        "reference": "Romans 1:16 (NKJV)",
                        "text": "For I am not ashamed of the gospel.",
                    }
                ],
            }
        ],
    }
    create_presentation(load_verses_from_dict(data), str(out))
    prs = Presentation(out)
    slide = prs.slides[1]
    tops = sorted(s.top.inches for s in slide.shapes if s.has_text_frame)
    assert tops[0] < 1.0


def test_per_verse_reference_position_override(tmp_path):
    out = tmp_path / "mixed.pptx"
    data = {
        "presentation_title": "Ref",
        "slide_style": {"reference_position": "bottom"},
        "sections": [
            {
                "section": "",
                "verses": [
                    {
                        "reference": "Romans 1:16 (NKJV)",
                        "text": "Short verse.",
                        "reference_position": "below",
                    }
                ],
            }
        ],
    }
    create_presentation(load_verses_from_dict(data), str(out))
    below_y = _ref_top_inches(out)
    assert below_y is not None
    assert below_y < 5.5
