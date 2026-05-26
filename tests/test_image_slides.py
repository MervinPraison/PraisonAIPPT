"""Tests for image slide support."""

import pytest
from pathlib import Path

from praisonaippt import create_presentation, load_verses_from_dict
from praisonaippt.exceptions import SchemaError
from praisonaippt.schema import validate_verses
from praisonaippt.utils import resolve_asset_path

PKG_ROOT = Path(__file__).resolve().parent.parent
SAMPLE_IMAGE = PKG_ROOT / "assets" / "background_alt.jpg"


@pytest.mark.skipif(not SAMPLE_IMAGE.is_file(), reason="sample image missing")
def test_resolve_asset_path_from_repo_root():
    got = resolve_asset_path("assets/background_alt.jpg")
    assert got == str(SAMPLE_IMAGE.resolve())


def test_resolve_asset_path_relative_to_source_file(tmp_path):
    img = tmp_path / "pic.jpg"
    img.write_bytes(b"\xff\xd8\xff")  # minimal placeholder; pptx may still load
    deck = tmp_path / "deck.yaml"
    deck.write_text("x: 1\n", encoding="utf-8")
    got = resolve_asset_path("pic.jpg", source_file=deck)
    assert got == str(img.resolve())


def test_validate_image_slide_requires_path():
    with pytest.raises(SchemaError):
        validate_verses(
            {
                "presentation_title": "T",
                "sections": [{"section": "", "verses": [{"slide_type": "image"}]}],
            }
        )


def test_validate_image_slide_accepts_path_only():
    out = validate_verses(
        {
            "presentation_title": "T",
            "sections": [
                {
                    "section": "",
                    "verses": [
                        {
                            "slide_type": "image",
                            "image_path": "assets/background_alt.jpg",
                        }
                    ],
                }
            ],
        }
    )
    assert out["sections"][0]["verses"][0]["image_path"]


@pytest.mark.skipif(not SAMPLE_IMAGE.is_file(), reason="sample image missing")
def test_create_presentation_with_image_slide(tmp_path):
    data = {
        "presentation_title": "Image test",
        "slide_style": {"background_color": "#1A1A2E"},
        "sections": [
            {
                "section": "",
                "verses": [
                    {
                        "slide_type": "image",
                        "image_path": str(SAMPLE_IMAGE),
                        "reference": "Diagram",
                        "text": "Optional caption",
                        "image_fit": "contain",
                    }
                ],
            }
        ],
    }
    out = tmp_path / "image_test.pptx"
    path = create_presentation(load_verses_from_dict(data), output_file=str(out))
    assert path == str(out)
    assert out.is_file()

    from pptx import Presentation

    prs = Presentation(str(out))
    assert len(prs.slides) == 2  # title + image
    has_picture = any(
        s.shape_type == 13 for s in prs.slides[1].shapes
    )
    assert has_picture
