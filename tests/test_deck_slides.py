"""Tests for HeyGen-style deck slide layouts."""

from pathlib import Path

import pytest
from pptx import Presentation
from pptx.util import Inches

from praisonaippt import create_presentation, load_verses_from_dict, list_renderers
from praisonaippt.deck_slides import (
    DECK_COLOR_PRESETS,
    DECK_RECT_AVATAR_TYPES,
    DECK_SLIDE_TYPES,
    _title_block_height,
    deck_avatar_shape,
    deck_skips_avatar_overlay,
    deck_skips_media_overlay,
    export_deck_slide_regions,
    render_deck_slide,
    resolve_deck_style,
)
from praisonaippt.exceptions import SchemaError
from praisonaippt.slide_renderers import resolve_renderer

PKG = Path(__file__).resolve().parent.parent


def _minimal_verse(slide_type: str) -> dict:
    verse = {"slide_type": slide_type, "text": "Title"}
    if slide_type == "deck_exec_summary":
        verse["items"] = [{"badge": "01", "heading": "One", "text": "x"}]
    elif slide_type in ("deck_split_performance", "deck_channel_analysis"):
        verse["rows"] = [{"badge": "Q1", "number": "1", "text": "x"}]
    elif slide_type == "deck_region_grid":
        verse["cells"] = [{"number": "1%", "label": "A", "text": "x"}]
    elif slide_type in ("deck_product_columns", "deck_customer_segments"):
        verse["columns"] = [{"number": "1%", "label": "A", "text": "x"}]
    elif slide_type == "deck_agenda":
        verse["items"] = ["Introduction", "Summary"]
    elif slide_type == "deck_intro_split":
        verse["reference"] = "Overview body text."
        verse["media_path"] = "assets/background_alt.jpg"
    elif slide_type == "deck_opportunity_cards":
        verse["columns"] = [
            {"badge": "01", "heading": "A", "text": "x", "image_path": "assets/background_alt.jpg"}
        ]
    elif slide_type == "deck_forecast_split":
        verse["items"] = [{"badge": "01", "text": "Forecast point."}]
        verse["media_path"] = "assets/background_alt.jpg"
    return verse


@pytest.mark.parametrize("slide_type", DECK_SLIDE_TYPES)
def test_deck_renderer_registered(slide_type):
    verse = _minimal_verse(slide_type)
    assert resolve_renderer(verse).kind == slide_type
    assert slide_type in list_renderers()


@pytest.mark.parametrize("slide_type", DECK_SLIDE_TYPES)
def test_deck_validate_minimal_verse(slide_type):
    verse = _minimal_verse(slide_type)
    resolve_renderer(verse).validate(verse, f"test.{slide_type}")


@pytest.mark.parametrize("slide_type", DECK_SLIDE_TYPES)
def test_deck_render_smoke(slide_type, tmp_path):
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    verse = _minimal_verse(slide_type)
    render_deck_slide(prs, slide_type, verse, source_file=str(PKG))
    assert len(prs.slides) == 1
    assert len(prs.slides[0].shapes) >= 1


@pytest.mark.parametrize("slide_type", DECK_SLIDE_TYPES)
def test_export_deck_regions_known_kind(slide_type):
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    regions = export_deck_slide_regions(prs, slide_type, {})
    assert "avatar" in regions
    assert "media" in regions
    assert "content" in regions


def test_all_color_presets_in_gallery():
    import yaml

    gallery = yaml.safe_load((PKG / "examples" / "deck_template_gallery.yaml").read_text())
    used = {
        v.get("color_scheme")
        for s in gallery["sections"]
        for v in s.get("verses", [])
        if v.get("color_scheme")
    }
    assert used <= set(DECK_COLOR_PRESETS.keys())


def test_rect_avatar_shape_for_split_layouts():
    for kind in DECK_RECT_AVATAR_TYPES:
        assert deck_avatar_shape(kind, {}, "circle") == "rect"


def test_circle_avatar_shape_for_pip_layouts():
    assert deck_avatar_shape("deck_exec_summary", {}, "circle") == "circle"
    assert deck_avatar_shape("deck_region_grid", {}, "circle") == "circle"


def test_deck_skips_avatar_overlay():
    assert deck_skips_avatar_overlay("deck_thank_you")
    assert deck_skips_avatar_overlay("deck_title_split")
    assert not deck_skips_avatar_overlay("deck_exec_summary")


def test_deck_skips_media_overlay():
    for kind in DECK_SLIDE_TYPES:
        assert deck_skips_media_overlay(kind)


def test_resolve_deck_style_preset():
    style = resolve_deck_style({}, {"color_scheme": "split_blue"}, "deck_split_performance")
    assert style.get("left_panel_color") == "#4338CA"


def test_title_split_avatar_right_half():
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    regions = export_deck_slide_regions(prs, "deck_title_split", {})
    avatar = regions["avatar"]
    assert avatar is not None
    assert avatar.left_in > prs.slide_width.inches * 0.45


def test_deck_rect_avatar_forces_rect_shape_despite_circle_pip():
    from praisonaippt.avatar_layouts import RegionBox, resolve_avatar_shape

    style = {"layouts": {"pip": {"shape": "circle"}}}
    wide = RegionBox(6.0, 0.0, 6.5, 7.5)
    assert deck_avatar_shape("deck_title_split", style, "circle", box=wide) == "rect"
    strip = RegionBox(0.3, 4.0, 5.0, 3.0)
    assert resolve_avatar_shape(style, layout_kind="deck_split_performance", box=strip) == "h_rect"
    assert deck_avatar_shape("deck_exec_summary", style, "circle") == "circle"


def test_public_slide_image_exports():
    from praisonaippt import SlideImageOptions, default_slide_images_dir, export_pptx_slide_jpegs

    assert SlideImageOptions is not None
    assert callable(export_pptx_slide_jpegs)
    assert callable(default_slide_images_dir)


def test_exec_summary_pip_top_right():
    prs = Presentation()
    style = {"layouts": {"deck_exec_summary": {"pip_position": "top_right"}}}
    regions = export_deck_slide_regions(prs, "deck_exec_summary", style)
    avatar = regions["avatar"]
    assert avatar.top_in < 1.0
    assert avatar.left_in > prs.slide_width.inches * 0.7
    assert regions["content"] is not None
    assert regions["media"] is None


def test_region_grid_pip_bottom_left():
    prs = Presentation()
    style = {"layouts": {"deck_region_grid": {"pip_position": "bottom_left"}}}
    regions = export_deck_slide_regions(prs, "deck_region_grid", style)
    avatar = regions["avatar"]
    assert avatar.left_in < 1.5
    assert avatar.top_in > prs.slide_height.inches * 0.6


def test_region_grid_content_clears_bottom_left_pip():
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    style = {"layouts": {"deck_region_grid": {"pip_position": "bottom_left"}}}
    regions = export_deck_slide_regions(prs, "deck_region_grid", style)
    pip, content = regions["avatar"], regions["content"]
    assert content.left_in >= pip.left_in + pip.width_in + 0.1
    assert content.top_in + content.height_in <= pip.top_in - 0.1


def test_thank_you_avatar_framing_zoomed_out():
    from praisonaippt.avatar_layouts import avatar_framing

    _, _, zoom = avatar_framing({}, "deck_thank_you")
    _, _, pip_zoom = avatar_framing({}, "pip")
    assert zoom < pip_zoom
    assert zoom == 1.02


def test_agenda_list_starts_below_title():
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    verse = {
        "slide_type": "deck_agenda",
        "text": "Agenda",
        "items": ["Introduction", "Executive Summary"],
    }
    render_deck_slide(prs, "deck_agenda", verse, source_file=str(PKG))
    title_h = _title_block_height("Agenda", "", 13.33 * 0.4, {})
    min_list_top = 0.65 + title_h + 0.28
    badge_tops = []
    for shape in prs.slides[0].shapes:
        if shape.has_text_frame and shape.text_frame.text.strip() == "01":
            badge_tops.append(shape.top.inches)
    assert badge_tops
    assert min(badge_tops) >= min_list_top - 0.05


def test_product_columns_start_below_title_block():
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    verse = {
        "slide_type": "deck_product_columns",
        "text": "Top Product Performance",
        "reference": "Product Revenue Distribution",
        "columns": [{"number": "40%", "label": "Product A", "text": "Enterprise sector"}],
    }
    render_deck_slide(prs, "deck_product_columns", verse, source_file=str(PKG))
    margin = 0.65
    title_w = 6.5
    title_h = _title_block_height(verse["text"], verse["reference"], title_w, {})
    min_col_top = margin + title_h + 0.28
    metric_tops = []
    for shape in prs.slides[0].shapes:
        if shape.has_text_frame and "40%" in shape.text_frame.text:
            metric_tops.append(shape.top.inches)
    assert metric_tops
    assert min(metric_tops) >= min_col_top - 0.05


def test_agenda_content_region():
    prs = Presentation()
    regions = export_deck_slide_regions(prs, "deck_agenda", {})
    assert regions["content"] is not None
    assert regions["content"].height_in > 3.0
    assert regions["media"] is None


def test_intro_split_media_bottom():
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    regions = export_deck_slide_regions(prs, "deck_intro_split", {})
    assert regions["media"] is not None
    assert regions["media"].top_in > prs.slide_height.inches * 0.4


def test_unknown_kind_returns_empty_regions():
    prs = Presentation()
    regions = export_deck_slide_regions(prs, "not_a_deck", {})
    assert regions["avatar"] is None
    assert regions["media"] is None


def test_rows_string_coercion_renders(tmp_path):
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    verse = {"slide_type": "deck_split_performance", "text": "T", "rows": ["First row text"]}
    render_deck_slide(prs, "deck_split_performance", verse, source_file=str(PKG))
    texts = []
    for shape in prs.slides[0].shapes:
        if shape.has_text_frame:
            texts.append(shape.text_frame.text)
    assert any("First row text" in t for t in texts)


def test_validate_rejects_bad_rows():
    verse = {"slide_type": "deck_split_performance", "text": "T", "rows": [123]}
    with pytest.raises(SchemaError):
        resolve_renderer(verse).validate(verse, "test")


def test_validate_region_grid_accepts_columns():
    verse = {"slide_type": "deck_region_grid", "text": "T", "columns": [{"label": "A", "text": "x"}]}
    resolve_renderer(verse).validate(verse, "test")


def test_validate_agenda_accepts_agenda_key():
    verse = {"slide_type": "deck_agenda", "text": "Agenda", "agenda": ["One"]}
    resolve_renderer(verse).validate(verse, "test")


def test_validate_opportunity_accepts_items_key():
    verse = {"slide_type": "deck_opportunity_cards", "text": "O", "items": [{"heading": "H", "text": "x"}]}
    resolve_renderer(verse).validate(verse, "test")


def test_deck_gallery_builds(tmp_path):
    yaml_path = PKG / "examples" / "deck_template_gallery.yaml"
    if not yaml_path.is_file():
        pytest.skip("gallery yaml missing")
    import yaml

    data = yaml.safe_load(yaml_path.read_text(encoding="utf-8"))
    data["_source_file"] = str(PKG)
    out = tmp_path / "deck_gallery.pptx"
    create_presentation(load_verses_from_dict(data), str(out))
    prs = Presentation(out)
    assert len(prs.slides) == len(DECK_SLIDE_TYPES) + 2  # title + section + 12 layouts
