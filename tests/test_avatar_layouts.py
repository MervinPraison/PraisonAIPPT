"""Tests for avatar layout slide types."""

from pathlib import Path

import pytest
from pptx import Presentation

from praisonaippt import create_presentation, load_verses_from_dict, list_renderers
from praisonaippt.avatar_layouts import AVATAR_SLIDE_TYPES
from praisonaippt.exceptions import SchemaError
from praisonaippt.slide_renderers import resolve_renderer, validate_verse

PKG = Path(__file__).resolve().parent.parent
IMG = PKG / "assets" / "background_alt.jpg"


def test_list_renderers_includes_avatar_kinds():
    kinds = set(list_renderers())
    assert kinds >= set(AVATAR_SLIDE_TYPES)


@pytest.mark.parametrize("slide_type", AVATAR_SLIDE_TYPES)
def test_resolve_avatar_renderer(slide_type):
    verse = {"slide_type": slide_type}
    if slide_type in ("media_only", "media_border"):
        verse["media_path"] = "assets/x.jpg"
    elif slide_type in ("avatar_name_card", "avatar_headline", "avatar_headline_full"):
        verse["headline"] = "Title"
    elif slide_type == "avatar_quote":
        verse["text"] = "Quote"
    assert resolve_renderer(verse).kind == slide_type


def test_media_only_requires_media_path():
    with pytest.raises(SchemaError, match="media_path"):
        validate_verse({"slide_type": "media_only"}, "sections[0].verses[0]")


def test_avatar_name_card_requires_headline():
    with pytest.raises(SchemaError, match="headline"):
        validate_verse({"slide_type": "avatar_name_card"}, "sections[0].verses[0]")


def test_invalid_media_fit():
    with pytest.raises(SchemaError, match="media_fit"):
        validate_verse(
            {"slide_type": "avatar_media_1", "media_fit": "stretch"},
            "sections[0].verses[0]",
        )


def test_split_geometry_media_left():
    from pptx import Presentation
    from praisonaippt.avatar_layouts import _slide_regions

    prs = Presentation()
    m1, a1 = _slide_regions(prs, "avatar_media_1", {})["media"], _slide_regions(
        prs, "avatar_media_1", {}
    )["avatar"]
    m2, a2 = _slide_regions(prs, "avatar_media_2", {})["media"], _slide_regions(
        prs, "avatar_media_2", {}
    )["avatar"]
    assert m1.left_in < a1.left_in
    assert abs(m1.width_in / (m1.width_in + a1.width_in) - 0.5) < 0.02
    assert abs(m2.width_in / (m2.width_in + a2.width_in) - 0.4) < 0.02
    assert m1.width_in != m2.width_in


def test_border_split_ratios_differ():
    from pptx import Presentation
    from praisonaippt.avatar_layouts import _slide_regions

    prs = Presentation()
    m1 = _slide_regions(prs, "avatar_media_border_1", {})["media"]
    m2 = _slide_regions(prs, "avatar_media_border_2", {})["media"]
    assert m1.width_in > m2.width_in


def test_avatar_media_3_media_below_text_panel():
    from pptx import Presentation
    from praisonaippt.avatar_layouts import _slide_regions

    prs = Presentation()
    regions = _slide_regions(prs, "avatar_media_3", {})
    panel = regions["text_panel"]
    media = regions["media"]
    assert panel is not None
    assert media is not None
    assert media.top_in >= panel.top_in + panel.height_in - 0.02
    assert media.height_in > 3.0


def test_avatar_media_3_full_bleed_uses_full_media():
    from pptx import Presentation
    from praisonaippt.avatar_layouts import _content_area, _slide_regions

    prs = Presentation()
    style = {"layouts": {"avatar_media_3": {"hero_layout": "full_bleed"}}}
    verse = {"text_panel": {"anchor": "top_right"}}
    regions = _slide_regions(prs, "avatar_media_3", style, verse=verse)
    cx, cy, cw, ch = _content_area(prs, style, "avatar_media_3")
    media = regions["media"]
    assert media is not None
    assert abs(media.top_in - cy) < 0.02
    assert abs(media.height_in - ch) < 0.02
    panel = regions["text_panel"]
    assert panel.left_in > cx + cw * 0.5


def test_avatar_media_3_full_bleed_bottom_left_anchor():
    from pptx import Presentation
    from praisonaippt.avatar_layouts import _content_area, _slide_regions

    prs = Presentation()
    style = {"layouts": {"avatar_media_3": {"hero_layout": "full_bleed"}}}
    verse = {"text_panel": {"anchor": "bottom_left"}}
    regions = _slide_regions(prs, "avatar_media_3", style, verse=verse)
    cx, cy, cw, ch = _content_area(prs, style, "avatar_media_3")
    panel = regions["text_panel"]
    pip = regions["avatar"]
    assert panel.left_in <= cx + cw * 0.45
    overlaps = (
        panel.left_in + panel.width_in > pip.left_in - 0.1
        and panel.top_in + panel.height_in > pip.top_in - 0.1
        and panel.left_in < pip.left_in + pip.width_in
        and panel.top_in < pip.top_in + pip.height_in
    )
    assert not overlaps


def test_avatar_quote_has_no_media_region():
    from pptx import Presentation
    from praisonaippt.avatar_layouts import _slide_regions

    prs = Presentation()
    regions = _slide_regions(prs, "avatar_quote", {})
    assert regions["media"] is None
    assert regions["avatar"] is not None


def test_avatar_quote_does_not_bake_movie_shape():
    """Quote slides use FFmpeg PiP only — no embedded video shape (avoids double avatar)."""
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    from praisonaippt.avatar_layouts import render_avatar_slide

    prs = Presentation()
    verse = {
        "text": "How do you run agents?",
        "reference": "Subtitle",
        "avatar_video_path": "nonexistent.mp4",
    }
    slide = render_avatar_slide(prs, "avatar_quote", verse, {})
    movies = [s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.MEDIA]
    pictures = [s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE]
    assert len(movies) == 0
    assert len(pictures) == 0


def test_avatar_media_3_does_not_bake_pip_still():
    """Media-3 PiP is FFmpeg-only — baked still + overlay caused a double white ring."""
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    from praisonaippt.avatar_layouts import render_avatar_slide

    prs = Presentation()
    verse = {
        "headline": "Dreaming",
        "subheader": "Agents",
        "media_path": str(IMG),
        "avatar_video_path": "nonexistent.mp4",
    }
    slide = render_avatar_slide(prs, "avatar_media_3", verse, {})
    movies = [s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.MEDIA]
    pictures = [s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE]
    assert len(movies) == 0
    assert len(pictures) == 1  # hero media only — no baked PiP still


def test_avatar_headline_uses_pip_not_full_frame():
    from pptx import Presentation
    from praisonaippt.avatar_layouts import _slide_regions

    prs = Presentation()
    regions = _slide_regions(prs, "avatar_headline", {})
    avatar = regions["avatar"]
    panel = regions["text_panel"]
    assert avatar is not None
    assert panel is not None
    sw, sh = prs.slide_width.inches, prs.slide_height.inches
    assert avatar.width_in < sw * 0.25
    assert avatar.left_in > sw * 0.7
    assert panel.left_in + panel.width_in <= avatar.left_in + 0.05
    assert panel.top_in + panel.height_in <= sh - 0.5


def test_avatar_headline_full_uses_full_frame_and_top_panel():
    from pptx import Presentation
    from praisonaippt.avatar_layouts import _slide_regions

    prs = Presentation()
    regions = _slide_regions(prs, "avatar_headline_full", {})
    avatar = regions["avatar"]
    panel = regions["text_panel"]
    sw, sh = prs.slide_width.inches, prs.slide_height.inches
    assert avatar.width_in >= sw * 0.95
    assert avatar.height_in >= sh * 0.95
    assert panel is not None
    assert panel.top_in < sh * 0.2
    assert panel.left_in < sw * 0.2


def test_avatar_name_card_has_separate_pill_regions():
    from pptx import Presentation
    from praisonaippt.avatar_layouts import _slide_regions

    prs = Presentation()
    regions = _slide_regions(prs, "avatar_name_card", {})
    name = regions["name_pill"]
    title = regions["title_pill"]
    assert name is not None
    assert title is not None
    assert name.top_in < title.top_in
    assert name.left_in == title.left_in
    assert name.width_in >= title.width_in


def test_avatar_intro_has_no_avatar_region():
    from pptx import Presentation
    from praisonaippt.avatar_layouts import _slide_regions

    prs = Presentation()
    regions = _slide_regions(prs, "avatar_intro", {})
    assert regions["avatar"] is None
    assert regions["media"] is None


def test_avatar_outro_full_frame_avatar():
    from pptx import Presentation
    from praisonaippt.avatar_layouts import _slide_regions

    prs = Presentation()
    regions = _slide_regions(prs, "avatar_outro", {})
    avatar = regions["avatar"]
    sw, sh = prs.slide_width.inches, prs.slide_height.inches
    assert avatar is not None
    assert avatar.width_in >= sw * 0.95


@pytest.mark.skipif(not IMG.is_file(), reason="sample image missing")
def test_avatar_layouts_build(tmp_path):
    verses = [{"slide_type": "avatar_only"}]
    for kind in AVATAR_SLIDE_TYPES:
        if kind == "avatar_only":
            continue
        entry = {"slide_type": kind}
        if kind in ("media_only", "media_border") or kind.startswith("avatar_media"):
            entry["media_path"] = "assets/background_alt.jpg"
        if kind in ("avatar_name_card", "avatar_headline", "avatar_headline_full"):
            entry["headline"] = "Demo"
            entry["subheader"] = "Role"
        if kind == "avatar_quote":
            entry["text"] = "Sample quote."
            entry["reference"] = "Author"
        verses.append(entry)

    data = {
        "presentation_title": "Avatar layouts",
        "_source_file": str(PKG),
        "sections": [{"section": "", "verses": verses}],
    }
    out = tmp_path / "avatar.pptx"
    create_presentation(load_verses_from_dict(data), str(out))
    prs = Presentation(out)
    assert len(prs.slides) == len(AVATAR_SLIDE_TYPES) + 1  # title slide
