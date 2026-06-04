"""Tests for layout SDK helpers."""

from pathlib import Path

import pytest
from pptx import Presentation
from pptx.util import Inches

from praisonaippt import create_presentation, load_verses_from_dict
from praisonaippt.layout_tokens import (
    content_box,
    content_width_inches,
    layout_in,
    pip_reserve_inches,
    split_max_length_default,
    typography_pt,
)


def test_content_width_widescreen_caps_at_nine():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    assert content_width_inches(prs, {}, "verse") == 9.0


def test_content_width_standard_slide():
    prs = Presentation()
    # Default python-pptx slide width is 10"
    assert content_width_inches(prs, {}, "verse") == pytest.approx(8.8, abs=0.01)


def test_content_width_override():
    prs = Presentation()
    style = {"layouts": {"verse": {"content_width_in": 7.5}}}
    assert content_width_inches(prs, style, "verse") == 7.5


def test_title_slide_ignores_pip_reserve():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    style = {"layouts": {"pip": {"width_ratio": 0.24, "margin_in": 0.38}}}
    assert pip_reserve_inches(style, prs.slide_width.inches, kind="title") == 0.0
    left, width, width_in, _ = content_box(prs, style, "title")
    assert left.inches == pytest.approx((13.333 - width_in) / 2, abs=0.02)


def test_content_box_centres_on_slide():
    prs = Presentation()
    prs.slide_width = Inches(10)
    left, width, width_in, margin_in = content_box(prs, {}, "verse")
    assert margin_in == 0.6
    assert width_in == pytest.approx(8.8, abs=0.01)
    assert left.inches == pytest.approx((10 - 8.8) / 2, abs=0.01)


def test_layout_in_and_typography_pt_overrides():
    style = {
        "layouts": {"verse": {"margin_in": 1.0}},
        "typography": {"body_size_pt": 28},
    }
    assert layout_in(style, "verse", "margin_in", 0.6) == 1.0
    assert typography_pt(style, "body_size_pt") == 28


def test_split_max_length_default():
    assert split_max_length_default({}) == 200
    assert split_max_length_default({"split_max_length": 300}) == 300
    assert split_max_length_default({"split_max_length": "bad"}) == 200


def test_create_presentation_with_layout_override(tmp_path):
    out = tmp_path / "layout_override.pptx"
    data = {
        "presentation_title": "Layout test",
        "slide_style": {
            "background_color": "#1A1A2E",
            "layouts": {"verse": {"margin_in": 0.9}},
            "typography": {"body_size_pt": 28},
        },
        "sections": [
            {
                "section": "",
                "verses": [{"reference": "Ps 23:1", "text": "The Lord is my shepherd."}],
            }
        ],
    }
    create_presentation(load_verses_from_dict(data), str(out))
    assert out.is_file() and out.stat().st_size > 0
